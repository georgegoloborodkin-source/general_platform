"""
Company Second Brain V2 — Anthropic/Claude 2026 Edition
Enterprise-grade VC intelligence backend with Claude-native document ingestion,
structured outputs, GraphRAG retrieval, SSE streaming, and JWT-based ACLs.
"""

from __future__ import annotations

import asyncio
import base64
import csv
import json
import os
import random
import re
import sys
import time
from functools import lru_cache
from io import BytesIO, StringIO
from typing import Any, AsyncGenerator, Dict, List, Optional, Tuple

import httpx
from fastapi import Depends, FastAPI, HTTPException, Header, Request, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, RedirectResponse
from pydantic import BaseModel, Field

# ---------- High-Performance Foundation ----------
# ORJSONResponse: 20-50 % faster JSON serialisation than stdlib json
# IMPORTANT: FastAPI lets you import ORJSONResponse even without orjson,
# but it crashes at runtime. So we must check for the orjson module itself.
try:
    import orjson  # noqa: F401
    from fastapi.responses import ORJSONResponse
except ImportError:
    from fastapi.responses import JSONResponse as ORJSONResponse  # type: ignore[assignment]

# Anthropic SDK — async client for all Claude interactions
try:
    import anthropic
    _anthropic_sdk_available = True
except ImportError:
    _anthropic_sdk_available = False


def _is_anthropic_overloaded_or_rate_limited(e: Exception) -> bool:
    """True if the error is 529 (overloaded) or 429 (rate limit) — safe to retry with backoff."""
    msg = str(e).lower()
    if "529" in str(e) or "overloaded" in msg or "429" in str(e) or "rate" in msg:
        return True
    if _anthropic_sdk_available and anthropic is not None:
        if getattr(e, "status_code", None) in (429, 529, 503):
            return True
        if type(e).__name__ == "APIStatusError" and getattr(e, "status_code", None) in (429, 529, 503):
            return True
    return False


def _friendly_anthropic_error_message(e: Exception) -> str:
    """User-facing message for 529/429/overloaded so chat doesn't show raw API errors."""
    if _is_anthropic_overloaded_or_rate_limited(e):
        return "The AI service is temporarily busy. Please try again in a moment."
    return str(e)[:300]

# Optional: ollama for legacy/local model support
try:
    import ollama
except ImportError:
    ollama = None  # type: ignore[assignment]

# Supabase — backend-side retrieval for Agentic RAG
try:
    from supabase import create_client as _supabase_create_client
    _supabase_available = True
except ImportError:
    _supabase_available = False

app = FastAPI(
    title="Company Second Brain V2 API",
    default_response_class=ORJSONResponse,
)


# ---------------------------------------------------------------------------
#  Enterprise Security — JWT-based Document ACLs
# ---------------------------------------------------------------------------

JWT_SECRET = os.getenv("JWT_SECRET", "")
JWT_ALGORITHM = os.getenv("JWT_ALGORITHM", "HS256")
# Set to "true" to enforce JWT auth on all protected endpoints
ENFORCE_AUTH = os.getenv("ENFORCE_AUTH", "false").lower() == "true"


class AuthContext(BaseModel):
    """Decoded JWT claims — used to scope document retrieval."""
    user_id: str = "anonymous"
    group_ids: List[str] = Field(default_factory=list)
    org_id: str = ""
    role: str = "viewer"  # viewer | editor | admin


async def get_auth_context(
    authorization: Optional[str] = Header(default=None),
) -> AuthContext:
    """
    FastAPI dependency that extracts user identity from the Authorization header.
    When ENFORCE_AUTH is true, a valid JWT is required.
    When false (default), returns an anonymous context — useful during dev.
    """
    if not authorization or not JWT_SECRET:
        if ENFORCE_AUTH:
            raise HTTPException(status_code=401, detail="Authorization header required.")
        return AuthContext()

    token = authorization.replace("Bearer ", "").strip()
    if not token:
        if ENFORCE_AUTH:
            raise HTTPException(status_code=401, detail="Bearer token required.")
        return AuthContext()

    try:
        import jwt  # PyJWT
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        return AuthContext(
            user_id=payload.get("sub", payload.get("user_id", "anonymous")),
            group_ids=payload.get("group_ids", payload.get("groups", [])),
            org_id=payload.get("org_id", payload.get("org", "")),
            role=payload.get("role", "viewer"),
        )
    except ImportError:
        if ENFORCE_AUTH:
            raise HTTPException(
                status_code=503,
                detail="JWT auth enabled but PyJWT not installed. Run: pip install PyJWT",
            )
        return AuthContext()
    except Exception as e:
        raise HTTPException(status_code=401, detail=f"Invalid token: {e}")


def acl_metadata_filter(auth: AuthContext) -> Dict[str, Any]:
    """
    Build a metadata filter dict that can be passed to a vector DB query.
    Ensures users can only retrieve documents they have access to.

    Usage in your vector search:
        filters = acl_metadata_filter(auth)
        results = vector_db.query(embedding, filter=filters, top_k=10)
    """
    if auth.role == "admin":
        # Admins can see everything in their org
        if auth.org_id:
            return {"org_id": auth.org_id}
        return {}

    # Non-admins: must match user_id OR one of their group_ids
    conditions: List[Dict[str, Any]] = [{"user_id": auth.user_id}]
    for gid in auth.group_ids:
        conditions.append({"group_id": gid})
    if auth.org_id:
        conditions.append({"org_id": auth.org_id, "visibility": "org"})

    return {"$or": conditions} if len(conditions) > 1 else conditions[0]

# Limit how much extracted text we send to the model (large PDFs often cause truncated JSON output).
MAX_MODEL_INPUT_CHARS = int(os.environ.get("MAX_MODEL_INPUT_CHARS", "24000"))

# PDF page limits
MAX_PDF_PAGES = int(os.environ.get("MAX_PDF_PAGES", "12"))
MAX_PARALLEL_PAGES = int(os.environ.get("MAX_PARALLEL_PAGES", "10"))


# ---------------------------------------------------------------------------
#  Claude-Native Document Ingestion (V2 — replaces Tesseract/OCR)
# ---------------------------------------------------------------------------

def _get_anthropic_async_client() -> "anthropic.AsyncAnthropic":
    """Return a cached Anthropic async client (SDK-based)."""
    if not _anthropic_sdk_available:
        raise HTTPException(
            status_code=503,
            detail="anthropic SDK not installed. Run: pip install anthropic",
        )
    if not ANTHROPIC_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="ANTHROPIC_API_KEY not set.",
        )
    return anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY)


async def _call_claude_with_fallback(
    client,
    messages: list,
    max_tokens: int = 4096,
    temperature: float = 0.0,
    system: str | None = None,
    tools: list | None = None,
    tool_choice: dict | None = None,
    preferred_model: str | None = None,
    # Cost tracking params (optional)
    endpoint: str = "unknown",
    organization_id: str | None = None,
    event_id: str | None = None,
    user_id: str | None = None,
    **extra_kwargs,
):
    """
    Call Claude with automatic model fallback on 404 (retired model).
    Tries preferred_model → ANTHROPIC_MODEL → each fallback in ANTHROPIC_MODEL_FALLBACKS.
    Automatically logs API usage for cost tracking.
    """
    models_to_try = list(dict.fromkeys(filter(None, [
        preferred_model,
        ANTHROPIC_MODEL,
        *ANTHROPIC_MODEL_FALLBACKS,
    ])))

    last_error = None
    for model_name in models_to_try:
        try:
            kwargs = {
                "model": model_name,
                "max_tokens": max_tokens,
                "temperature": temperature,
                "messages": messages,
                **extra_kwargs,
            }
            if system:
                kwargs["system"] = system
            if tools:
                kwargs["tools"] = tools
            if tool_choice:
                kwargs["tool_choice"] = tool_choice
            result = await client.messages.create(**kwargs)
            
            # Log usage for cost tracking (non-blocking)
            if hasattr(result, 'usage') and result.usage:
                usage = result.usage
                input_tokens = getattr(usage, 'input_tokens', 0)
                output_tokens = getattr(usage, 'output_tokens', 0)
                # Fire-and-forget logging (don't await to avoid blocking)
                import asyncio
                asyncio.create_task(log_api_usage(
                    provider="anthropic",
                    model=model_name,
                    endpoint=endpoint,
                    input_tokens=input_tokens,
                    output_tokens=output_tokens,
                    organization_id=organization_id,
                    event_id=event_id,
                    user_id=user_id,
                    request_id=getattr(result, 'id', None),
                ))
            
            return result
        except Exception as e:
            err_str = str(e)
            if "404" in err_str or "not_found" in err_str:
                last_error = e
                continue
            # Log error for failed requests too
            import asyncio
            asyncio.create_task(log_api_usage(
                provider="anthropic",
                model=model_name,
                endpoint=endpoint,
                input_tokens=0,
                output_tokens=0,
                organization_id=organization_id,
                event_id=event_id,
                user_id=user_id,
                error_message=str(e)[:500],
            ))
            raise  # Non-404 error → don't retry

    raise last_error or RuntimeError(f"All models failed: {models_to_try}")


async def extract_pdf_with_claude_native(pdf_bytes: bytes, max_pages: int = 10) -> str:
    """
    Send a PDF directly to Claude 3.5/3.7 Sonnet as a *document content block*
    (base64-encoded).  Claude "sees" layouts, charts, tables, and scanned text
    natively — no Tesseract, no pdf2image, no brittle OCR pipeline.

    Falls back to image-based extraction if the PDF is too large for a single
    document block (>25 MB after base64).
    """
    if not ANTHROPIC_API_KEY:
        return ""

    pdf_b64 = base64.b64encode(pdf_bytes).decode("utf-8")

    # Claude document blocks accept PDFs up to ~32 MB base64
    if len(pdf_b64) > 32_000_000:
        # Too large — fall back to page-image extraction
        return await _extract_pdf_as_page_images(pdf_bytes, max_pages)

    prompt = (
        "Extract ALL text from this document. Preserve:\n"
        "- Table structure (use markdown tables)\n"
        "- Headers, footers, page numbers\n"
        "- Multi-column layouts\n"
        "- Bullet points and numbered lists\n"
        "- All numbers, dates, currency values, and proper nouns\n\n"
        "Return the extracted text in a clear, structured format. "
        "Separate pages with '--- Page N ---' markers."
    )

    try:
        client = _get_anthropic_async_client()
        message = await _call_claude_with_fallback(
            client,
            max_tokens=8192,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "document",
                            "source": {
                                "type": "base64",
                                "media_type": "application/pdf",
                                "data": pdf_b64,
                            },
                        },
                        {"type": "text", "text": prompt},
                    ],
                }
            ],
        )
        # Extract text from response
        text_parts = [
            block.text for block in message.content if hasattr(block, "text")
        ]
        return "\n".join(text_parts).strip()
    except Exception as e:
        # Fallback to page-image approach
        return await _extract_pdf_as_page_images(pdf_bytes, max_pages)


async def _extract_pdf_as_page_images(pdf_bytes: bytes, max_pages: int = 10) -> str:
    """
    Fallback: render PDF pages as PNG images, send each to Claude Vision.
    Used when the PDF is too large for a single document content block.
    """
    if not ANTHROPIC_API_KEY:
        return ""

    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""

    doc = fitz.open(stream=BytesIO(pdf_bytes).getvalue(), filetype="pdf")
    page_limit = min(doc.page_count, max_pages)

    parts: list[str] = []
    client = _get_anthropic_async_client()

    for i in range(page_limit):
        try:
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(200 / 72, 200 / 72))
            img_b64 = base64.b64encode(pix.tobytes("png")).decode("utf-8")

            message = await _call_claude_with_fallback(
                client,
                max_tokens=4096,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": "image/png",
                                    "data": img_b64,
                                },
                            },
                            {
                                "type": "text",
                                "text": "Extract all text from this page. Preserve tables (markdown), lists, and formatting.",
                            },
                        ],
                    }
                ],
            )
            page_text = "".join(
                b.text for b in message.content if hasattr(b, "text")
            )
            if page_text:
                parts.append(f"\n--- Page {i + 1} (Claude Vision) ---\n{page_text}")
        except Exception as e:
            parts.append(f"\n--- Page {i + 1} (extraction failed: {e}) ---\n")

    return "\n".join(parts).strip() if parts else ""


async def _describe_image_with_vision(image_bytes: bytes, media_type: str, file_name: str) -> str:
    """Describe an image using Claude Vision so the model can 'read' pictures."""
    if not ANTHROPIC_API_KEY:
        return f"[Image file: {file_name}. Set ANTHROPIC_API_KEY to enable vision description.]"
    media_type = media_type or "image/png"
    if media_type not in ("image/png", "image/jpeg", "image/gif", "image/webp"):
        media_type = "image/png"
    try:
        img_b64 = base64.b64encode(image_bytes).decode("utf-8")
        client = _get_anthropic_async_client()
        message = await _call_claude_with_fallback(
            client,
            max_tokens=2048,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": img_b64,
                            },
                        },
                        {
                            "type": "text",
                            "text": "Describe this image in detail: any text (OCR), labels, charts, diagrams, people, objects, and context. Output plain text suitable for search and retrieval.",
                        },
                    ],
                }
            ],
        )
        text = "".join(b.text for b in message.content if hasattr(b, "text"))
        return text.strip() if text else f"[Image: {file_name} — no description generated]"
    except Exception as e:
        return f"[Image: {file_name} — vision description failed: {str(e)[:150]}]"


# Converter provider settings
_provider_env = os.getenv("CONVERTER_PROVIDER")
CONVERTER_PROVIDER = (_provider_env or "ollama").lower().strip()

# Ollama connection settings
OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://127.0.0.1:11434").rstrip("/")
PREFERRED_OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "vc-converter:latest")

# Anthropic (Claude) settings
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
_raw_anthropic_model = os.getenv("ANTHROPIC_MODEL", "claude-sonnet-4-20250514")

# ── Validate the configured model: retired models get auto-replaced ──
_RETIRED_MODELS = {
    "claude-3-5-sonnet-20241022",
    "claude-3-5-sonnet-20240620",
    "claude-3-5-haiku-20241022",
    "claude-3-haiku-20240307",
    "claude-3-opus-20240229",
}
if _raw_anthropic_model in _RETIRED_MODELS:
    ANTHROPIC_MODEL = "claude-sonnet-4-20250514"
else:
    ANTHROPIC_MODEL = _raw_anthropic_model

ANTHROPIC_API_URL = os.getenv("ANTHROPIC_API_URL", "https://api.anthropic.com/v1/messages")

# Model fallback chain — only non-retired models
ANTHROPIC_MODEL_FALLBACKS = [
    m for m in dict.fromkeys([          # dedupe while preserving order
        ANTHROPIC_MODEL,
        "claude-sonnet-4-20250514",
        "claude-3-7-sonnet-latest",
    ]) if m
]

# ---------------------------------------------------------------------------
#  Supabase — backend-side database access for Agentic RAG
# ---------------------------------------------------------------------------
SUPABASE_URL = os.getenv("SUPABASE_URL", "")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")

_sb_client = None

def get_supabase():
    """Lazy-init Supabase client (service-role key for backend retrieval)."""
    global _sb_client
    if _sb_client is not None:
        return _sb_client
    if not _supabase_available:
        raise RuntimeError("supabase-py not installed — run: pip install supabase")
    if not SUPABASE_URL or not SUPABASE_SERVICE_KEY:
        raise RuntimeError("SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY env vars required for Agentic RAG")
    _sb_client = _supabase_create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
    return _sb_client

# ---------------------------------------------------------------------------
#  API Cost Tracking — log Anthropic usage for unit economics
# ---------------------------------------------------------------------------

# Anthropic pricing (per 1M tokens, as of 2025)
ANTHROPIC_PRICING = {
    "claude-sonnet-4-20250514": {"input": 3.0, "output": 15.0},
    "claude-opus-4-20250514": {"input": 15.0, "output": 75.0},
    "claude-haiku-4-20250514": {"input": 0.25, "output": 1.25},
    "claude-3-7-sonnet-latest": {"input": 3.0, "output": 15.0},
    # Fallback for unknown models
    "default": {"input": 3.0, "output": 15.0},
}

def estimate_anthropic_cost(model: str, input_tokens: int, output_tokens: int) -> float:
    """Calculate estimated cost in USD for an Anthropic API call."""
    pricing = ANTHROPIC_PRICING.get(model, ANTHROPIC_PRICING["default"])
    input_cost = (input_tokens / 1_000_000) * pricing["input"]
    output_cost = (output_tokens / 1_000_000) * pricing["output"]
    return round(input_cost + output_cost, 6)

async def log_api_usage(
    provider: str,
    model: str,
    endpoint: str,
    input_tokens: int,
    output_tokens: int,
    organization_id: str | None = None,
    event_id: str | None = None,
    user_id: str | None = None,
    request_id: str | None = None,
    error_message: str | None = None,
):
    """Log API usage to database for cost tracking. Non-blocking (fire-and-forget)."""
    if not _supabase_available or not SUPABASE_URL or not SUPABASE_SERVICE_KEY:
        return  # Skip logging if Supabase not configured
    
    try:
        total_tokens = input_tokens + output_tokens
        estimated_cost = 0.0
        
        if provider == "anthropic":
            estimated_cost = estimate_anthropic_cost(model, input_tokens, output_tokens)
        
        sb = get_supabase()
        # Insert usage log (Supabase client is sync, but we're in async context)
        import asyncio
        await asyncio.to_thread(
            lambda: sb.table("api_usage_logs").insert({
                "provider": provider,
                "model": model,
                "endpoint": endpoint,
                "input_tokens": input_tokens,
                "output_tokens": output_tokens,
                "total_tokens": total_tokens,
                "estimated_cost_usd": estimated_cost,
                "organization_id": organization_id,
                "event_id": event_id,
                "user_id": user_id,
                "request_id": request_id,
                "error_message": error_message,
            }).execute()
        )
    except Exception as e:
        # Don't break API calls if logging fails
        pass

# Ask-the-fund settings (generous tokens for comprehensive answers)
ASK_MAX_TOKENS = int(os.getenv("ASK_MAX_TOKENS", "4000"))  # Increased from 1000 for more detailed responses
ASK_MAX_SOURCES = int(os.getenv("ASK_MAX_SOURCES", "30"))  # Cap sources to avoid token waste
ASK_MAX_SNIPPET_CHARS = int(os.getenv("ASK_MAX_SNIPPET_CHARS", "500"))  # Larger snippets for better answers
# Use Haiku for simple questions (3-5x faster, 75% cheaper)
USE_HAIKU_FOR_SIMPLE = os.getenv("USE_HAIKU_FOR_SIMPLE", "true").lower() == "true"

# Embeddings settings (semantic search)
OLLAMA_EMBEDDING_MODEL = os.getenv("OLLAMA_EMBEDDING_MODEL", "nomic-embed-text")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_EMBEDDING_MODEL = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")  # 1536 dimensions
VOYAGE_API_KEY = os.getenv("VOYAGE_API_KEY") or os.getenv("VOYAGER_API_KEY")
VOYAGE_EMBEDDING_MODEL = os.getenv("VOYAGE_EMBEDDING_MODEL", "voyage-large-2")

# Auto-detect embedding dimension based on model
# voyage-4-large supports 256/512/1024/2048 (NOT 1536) — keeping voyage-large-2 for 1536-dim compatibility
_voyage_model_dims = {
    "voyage-4-large": 1024,
    "voyage-3-lite": 512,
    "voyage-3": 1024,
    "voyage-large-2": 1536,
    "voyage-finance-2": 1024,
    "voyage-code-3": 1024,
    "voyage-law-2": 1024,
}
_auto_dim = _voyage_model_dims.get(VOYAGE_EMBEDDING_MODEL, 1024)
EMBEDDING_DIM = int(os.getenv("EMBEDDING_DIM", str(_auto_dim)))

# Auto-detect embedding provider: use explicit env var, or pick the first available key
_explicit_provider = os.getenv("EMBEDDINGS_PROVIDER", "").lower().strip()
if _explicit_provider:
    EMBEDDINGS_PROVIDER = _explicit_provider
elif VOYAGE_API_KEY:
    EMBEDDINGS_PROVIDER = "voyage"
elif OPENAI_API_KEY:
    EMBEDDINGS_PROVIDER = "openai"
elif ANTHROPIC_API_KEY:
    # Anthropic doesn't have a native embedding API — use Voyage via Anthropic partnership
    # Fall back to OpenAI model if available, else warn
    EMBEDDINGS_PROVIDER = "openai"  # placeholder — will fail without key
else:
    EMBEDDINGS_PROVIDER = "voyage"  # default, but will warn on first call


# Reranking settings (cross-encoder) — Voyage rerank-2.5 is primary; Cohere kept as fallback
COHERE_API_KEY = os.getenv("COHERE_API_KEY")
RERANK_MODEL = os.getenv("RERANK_MODEL", "rerank-2.5")

# Ingestion settings
CLICKUP_API_TOKEN = os.getenv("CLICKUP_API_TOKEN")

def get_anthropic_api_url() -> str:
    """
    Normalize Anthropic API URL.
    Accepts base URLs like:
      - https://api.anthropic.com
      - https://api.anthropic.com/v1
      - https://api.anthropic.com/v1/messages
      - https://api.anthropic.com/v1/messages/
    Returns the full /v1/messages endpoint.
    """
    base = (ANTHROPIC_API_URL or "").strip()
    if not base:
        return "https://api.anthropic.com/v1/messages"
    base = base.rstrip("/")
    if base.endswith("/v1/messages"):
        return base
    if base.endswith("/v1"):
        return f"{base}/messages"
    return f"{base}/v1/messages"


async def fetch_ollama_model_names() -> List[str]:
    """
    More reliable than python ollama.list() on some setups.
    Uses Ollama's HTTP API to list installed models.
    """
    names: List[str] = []

    # First, try the HTTP /api/tags endpoint
    try:
        async with httpx.AsyncClient(timeout=2.0) as client:
            res = await client.get(f"{OLLAMA_HOST}/api/tags")
            res.raise_for_status()
            data = res.json() or {} 
            models = data.get("models", []) or []
            for m in models:
                if isinstance(m, dict) and m.get("name"):
                    names.append(m["name"])
    except Exception:
        # swallow and try python client fallback below
        names = []

    # Fallback: python client list() if HTTP returned nothing
    if not names:
        try:
            client = get_ollama_client()
            models = client.list()
            if isinstance(models, dict):
                for m in models.get("models", []) or []:
                    if isinstance(m, dict) and m.get("name"):
                        names.append(m["name"])
                    elif isinstance(m, str):
                        names.append(m)
        except Exception:
            pass

    return names


def pick_model(available_models: List[str]) -> str:
    """
    Pick a model name to use for conversion.
    Prefer env OLLAMA_MODEL, then vc-converter*, then llama3.1*, then llama3.2*, else first.
    """
    if not available_models:
        return PREFERRED_OLLAMA_MODEL

    if PREFERRED_OLLAMA_MODEL in available_models:
        return PREFERRED_OLLAMA_MODEL

    for prefix in ["vc-converter", "llama3.1", "llama3.2", "llama3"]:
        for name in available_models:
            if name.startswith(prefix):
                return name

    return available_models[0]


def get_ollama_client():
    """Force the host so the python client matches what `ollama list` uses."""
    if ollama is None:
        raise HTTPException(status_code=503, detail="ollama package not installed.")
    return ollama.Client(host=OLLAMA_HOST)

# CORS — allow Vercel frontend (and localhost for dev)
_CORS_ALLOWED = [
    "https://general-platform.vercel.app",
    "http://localhost:5174",
    "http://localhost:5173",
    "http://localhost:3000",
    "http://127.0.0.1:5174",
]


def _cors_origin_allowed(origin: str) -> bool:
    if not origin:
        return False
    if origin in _CORS_ALLOWED:
        return True
    # Allow Vercel preview deployments
    if origin.endswith(".vercel.app"):
        return True
    return False


@app.middleware("http")
async def add_cors_headers(request, call_next):
    """Ensure CORS headers on every response so Vercel can call the API."""
    origin = request.headers.get("origin", "")
    response = await call_next(request)
    if _cors_origin_allowed(origin):
        response.headers["Access-Control-Allow-Origin"] = origin
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, DELETE, OPTIONS"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization, Accept"
    response.headers["Access-Control-Max-Age"] = "86400"
    return response

app.add_middleware(
    CORSMiddleware,
    allow_origins=_CORS_ALLOWED,
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
    max_age=86400,
)

# Data models
class StartupData(BaseModel):
    companyName: str
    geoMarkets: List[str]
    industry: str
    fundingTarget: int
    fundingStage: str
    availabilityStatus: str = "present"

class InvestorData(BaseModel):
    firmName: str
    memberName: str
    geoFocus: List[str]
    industryPreferences: List[str]
    stagePreferences: List[str]
    minTicketSize: int
    maxTicketSize: int
    totalSlots: int = 3
    tableNumber: Optional[str] = None
    availabilityStatus: str = "present"

class MentorData(BaseModel):
    fullName: str
    email: str
    linkedinUrl: Optional[str] = None
    geoFocus: List[str]
    industryPreferences: List[str]
    expertiseAreas: List[str]
    totalSlots: int = 3
    availabilityStatus: str = "present"

class CorporateData(BaseModel):
    firmName: str
    contactName: str
    email: Optional[str] = None
    geoFocus: List[str]
    industryPreferences: List[str]
    partnershipTypes: List[str]
    stages: List[str]
    totalSlots: int = 3
    availabilityStatus: str = "present"

class ConversionRequest(BaseModel):
    data: str  # Unstructured data (text, CSV, JSON, etc.)
    dataType: Optional[str] = None  # 'startup', 'investor', or None for auto-detect
    format: Optional[str] = None  # 'csv', 'text', 'json', etc.

class ConversionResponse(BaseModel):
    startups: List[StartupData] = []
    investors: List[InvestorData] = []
    mentors: List[MentorData] = []
    corporates: List[CorporateData] = []
    detectedType: str
    confidence: float
    warnings: List[str] = []
    errors: List[str] = []
    raw_content: Optional[str] = None

class FileValidationResponse(BaseModel):
    isValid: bool
    errors: List[str] = []
    warnings: List[str] = []
    detectedType: Optional[str] = None
    startupCsvTemplate: Optional[str] = None
    investorCsvTemplate: Optional[str] = None

class ClickUpIngestRequest(BaseModel):
    list_id: str
    include_closed: bool = True

class ClickUpIngestResponse(BaseModel):
    tasks: List[Dict[str, Any]] = []

class ClickUpListsRequest(BaseModel):
    team_id: str

class AskSource(BaseModel):
    title: Optional[str] = None
    snippet: Optional[str] = None
    file_name: Optional[str] = None

class AskDecision(BaseModel):
    startup_name: Optional[str] = None
    action_type: Optional[str] = None
    outcome: Optional[str] = None
    notes: Optional[str] = None

class ChatMessage(BaseModel):
    role: str  # "user" or "assistant"
    content: str

class AskConnection(BaseModel):
    """Company connection from the connections graph — passed into chat for context."""
    source_company_name: str
    target_company_name: str
    connection_type: str = ""          # BD, INV, Knowledge, Partnership, Portfolio
    connection_status: str = ""        # To Connect, Connected, Rejected, In Progress, Completed
    ai_reasoning: Optional[str] = None
    notes: Optional[str] = None

class AskRequest(BaseModel):
    question: str
    sources: List[AskSource] = []
    decisions: List[AskDecision] = []
    connections: List[AskConnection] = Field(default_factory=list)
    previous_messages: List[ChatMessage] = Field(default_factory=list, alias="previousMessages")
    web_search_enabled: bool = Field(default=False, alias="webSearchEnabled")

    model_config = {"populate_by_name": True}

class AskResponse(BaseModel):
    answer: str

class RerankDocument(BaseModel):
    id: str
    text: str

class RerankRequest(BaseModel):
    query: str
    documents: List[RerankDocument]
    top_n: int | None = None

class RerankResult(BaseModel):
    id: str
    score: float

class RerankResponse(BaseModel):
    results: List[RerankResult]

class EmbedRequest(BaseModel):
    text: str
    input_type: Optional[str] = None

class EmbedResponse(BaseModel):
    embedding: List[float]

class RewriteQueryRequest(BaseModel):
    question: str
    previous_messages: Optional[List[ChatMessage]] = []

class RewriteQueryResponse(BaseModel):
    rewritten_question: str

class ClickUpListsResponse(BaseModel):
    lists: List[Dict[str, Any]] = []

class GoogleDriveIngestRequest(BaseModel):
    url: str
    access_token: Optional[str] = None

class GoogleDriveIngestResponse(BaseModel):
    title: str
    content: str
    raw_content: str  # Alias for content, for clarity
    sourceType: str

# ---------- Google Drive Folder-Sync Models ----------

class GDriveListFoldersRequest(BaseModel):
    access_token: str
    folder_id: str  # root folder ID

class GDriveFolderEntry(BaseModel):
    id: str
    name: str
    modifiedTime: Optional[str] = None

class GDriveListFoldersResponse(BaseModel):
    folders: List[GDriveFolderEntry] = []

class GDriveListFilesRequest(BaseModel):
    access_token: str
    folder_id: str

class GDriveFileEntry(BaseModel):
    id: str
    name: str
    mimeType: str
    modifiedTime: Optional[str] = None
    size: Optional[str] = None

class GDriveListFilesResponse(BaseModel):
    files: List[GDriveFileEntry] = []


class GDriveRefreshTokenRequest(BaseModel):
    refresh_token: str


class GDriveRefreshTokenResponse(BaseModel):
    access_token: str


class GDriveDownloadFileRequest(BaseModel):
    access_token: str
    file_id: str
    mime_type: Optional[str] = None
    file_name: Optional[str] = None

class GDriveDownloadFileResponse(BaseModel):
    title: str
    content: str
    raw_content: str
    sourceType: str
    mimeType: str

# ---------- Gmail Sync Models ----------

class GmailListMessagesRequest(BaseModel):
    access_token: str
    query: Optional[str] = None
    label_ids: Optional[List[str]] = None
    max_results: int = 50
    page_token: Optional[str] = None

class GmailMessageSnippet(BaseModel):
    id: str
    threadId: str
    snippet: Optional[str] = None

class GmailListMessagesResponse(BaseModel):
    messages: List[GmailMessageSnippet] = []
    next_page_token: Optional[str] = None
    result_size_estimate: int = 0

class GmailGetMessageRequest(BaseModel):
    access_token: str
    message_id: str

class GmailAttachmentMeta(BaseModel):
    id: str
    filename: str
    mimeType: str
    size: int = 0

class GmailGetMessageResponse(BaseModel):
    id: str
    threadId: str
    subject: str = ""
    sender: str = ""
    to: List[str] = []
    cc: List[str] = []
    date: Optional[str] = None
    body_text: str = ""
    body_html: str = ""
    labels: List[str] = []
    attachments: List[GmailAttachmentMeta] = []

class GmailDownloadAttachmentRequest(BaseModel):
    access_token: str
    message_id: str
    attachment_id: str

class GmailDownloadAttachmentResponse(BaseModel):
    data: str  # base64
    filename: str
    mimeType: str
    size: int = 0

class GmailIngestRequest(BaseModel):
    access_token: str
    message_id: str
    extract_attachments: bool = False

class GmailIngestResponse(BaseModel):
    title: str
    content: str
    raw_content: str
    sourceType: str
    email_from: str = ""
    email_to: List[str] = []
    email_cc: List[str] = []
    email_subject: str = ""
    email_date: Optional[str] = None
    gmail_thread_id: str = ""
    gmail_labels: List[str] = []
    has_attachments: bool = False
    attachments: List[GmailAttachmentMeta] = []

# System prompt for Ollama
SYSTEM_PROMPT = """You are a data extraction and conversion expert. Your task is to extract structured information from unstructured text and convert it into JSON format.

You will receive unstructured data about startups, investors, mentors, or corporates, and you must extract the following information:

FOR STARTUPS:
- companyName: The name of the company/startup
- geoMarkets: List of geographic markets (e.g., ["North America", "Europe"])
- industry: Industry sector (e.g., "AI/ML", "Fintech", "Healthtech")
- fundingTarget: Funding amount as integer (remove currency symbols, commas)
- fundingStage: Stage (e.g., "Pre-seed", "Seed", "Series A", "Series B+")

FOR INVESTORS:
- firmName: Name of the VC firm/investor
- memberName: The specific investor team member/person name (REQUIRED)
- geoFocus: List of geographic focus areas
- industryPreferences: List of preferred industries
- stagePreferences: List of preferred funding stages
- minTicketSize: Minimum investment amount as integer
- maxTicketSize: Maximum investment amount as integer
- totalSlots: Number of meeting slots (default: 3)
- tableNumber: Optional table/booth number

FOR MENTORS:
- fullName: Full name of the mentor (REQUIRED)
- email: Email address (REQUIRED)
- linkedinUrl: LinkedIn profile URL
- geoFocus: List of geographic focus areas
- industryPreferences: List of preferred industries
- expertiseAreas: List of expertise areas (e.g., ["Product Development", "Fundraising"])
- totalSlots: Number of meeting slots (default: 3)

FOR CORPORATES:
- firmName: Name of the corporate/company (REQUIRED)
- contactName: Name of the corporate contact person (REQUIRED)
- email: Email address
- geoFocus: List of geographic focus areas
- industryPreferences: List of preferred industries
- partnershipTypes: List of partnership types (e.g., ["Pilot Program", "Distribution"])
- stages: List of startup stages of interest (e.g., ["Seed", "Series A"])
- totalSlots: Number of meeting slots (default: 3)

IMPORTANT RULES:
1. Always return valid JSON only, no markdown or explanations
2. If multiple entities are found, return an array
3. Extract numbers from text (e.g., "$2M" -> 2000000, "€500K" -> 500000)
4. Parse lists from text (e.g., "North America, Europe" -> ["North America", "Europe"])
5. If information is missing, use reasonable defaults or empty arrays
6. For funding stages, normalize to: "Pre-seed", "Seed", "Series A", "Series B+"
7. For industries, use standard names: "Fintech", "Healthtech", "EdTech", "E-commerce", "Construction", "Transportation/Mobility", "AI/ML", "Logistics", "Consumer Goods", "SaaS", "CleanTech"

Return ONLY the JSON object or array, nothing else."""

def create_conversion_prompt(data: str, data_type: Optional[str] = None) -> str:
    """Create a prompt for Ollama to convert unstructured data"""
    # Keep prompt/model input bounded to reduce truncation.
    trimmed = data if len(data) <= MAX_MODEL_INPUT_CHARS else data[:MAX_MODEL_INPUT_CHARS]
    if len(trimmed) != len(data):
        trimmed = trimmed + "\n\n[TRUNCATED INPUT: content was longer than MAX_MODEL_INPUT_CHARS]"
    if data_type:
        prompt = (
            f"Extract {data_type} information from the following data and convert to JSON format.\n"
            f"Return ONLY valid JSON (no commentary). Keep it minimal: only the required fields.\n\n{trimmed}\n"
        )
    else:
        prompt = (
            "Extract startup or investor information from the following data. Auto-detect the type and convert to JSON.\n"
            "Return ONLY valid JSON (no commentary). Keep it minimal: only the required fields.\n\n"
            f"{trimmed}\n"
        )
    return prompt

def parse_ollama_response(response: str) -> Dict[str, Any]:
    """Parse model response and extract JSON"""
    # Remove markdown code blocks if present
    response = re.sub(r'```json\n?', '', response)
    response = re.sub(r'```\n?', '', response)
    response = response.strip()

    def extract_first_json_block(text: str) -> Optional[str]:
        """
        Extract the first complete JSON object/array from text using bracket balancing.
        This is robust against extra prose before/after JSON and avoids greedy regex traps.
        """
        start_idx = None
        stack: List[str] = []
        in_string = False
        escape = False

        for i, ch in enumerate(text):
            if start_idx is None:
                if ch == '{':
                    start_idx = i
                    stack = ['}']
                elif ch == '[':
                    start_idx = i
                    stack = [']']
                continue

            # We are inside a candidate JSON block
            if in_string:
                if escape:
                    escape = False
                elif ch == '\\':
                    escape = True
                elif ch == '"':
                    in_string = False
                continue

            if ch == '"':
                in_string = True
                continue

            if ch == '{':
                stack.append('}')
            elif ch == '[':
                stack.append(']')
            elif ch in ('}', ']'):
                if stack and ch == stack[-1]:
                    stack.pop()
                    if not stack and start_idx is not None:
                        return text[start_idx:i + 1]
                else:
                    # Mismatched closing bracket; keep scanning but this block is likely invalid.
                    pass

        return None

    # First try: extract a complete JSON block from within the response
    block = extract_first_json_block(response)
    if block:
        try:
            return json.loads(block)
        except json.JSONDecodeError:
            pass

    # Fallback: try parsing the whole response as JSON
    try:
        return json.loads(response)
    except json.JSONDecodeError:
        raise ValueError(
            "Could not parse JSON from model response (likely truncated / non-JSON). "
            f"Response starts with: {response[:200]}"
        )

def is_model_not_found(response: httpx.Response) -> bool:
    try:
        data = response.json() or {}
        err = data.get("error") or {}
        return err.get("type") == "not_found_error" and "model" in str(err.get("message", "")).lower()
    except Exception:
        return False

async def call_anthropic(prompt: str) -> str:
    """
    Call Anthropic (Claude) API and return the raw text response.
    Uses server-side API key from environment variables.
    """
    if not ANTHROPIC_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="ANTHROPIC_API_KEY not set. Set it in the server environment to use Claude."
        )

    headers = {
        "Content-Type": "application/json",
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": "2023-06-01",
    }
    url = get_anthropic_api_url()
    default_url = "https://api.anthropic.com/v1/messages"

    last_error: Optional[str] = None
    async with httpx.AsyncClient(timeout=60.0) as client:
        for model_name in [m for m in ANTHROPIC_MODEL_FALLBACKS if m]:
            payload = {
                "model": model_name,
                "max_tokens": 4096,
                "temperature": 0.1,
                "system": SYSTEM_PROMPT,
                "messages": [
                    {"role": "user", "content": prompt}
                ],
            }

            res = await client.post(url, headers=headers, json=payload)
            # If a misconfigured URL causes 404, retry with the canonical endpoint.
            if res.status_code == 404 and url != default_url:
                res = await client.post(default_url, headers=headers, json=payload)
            if res.status_code == 404 and is_model_not_found(res):
                last_error = f"Model not found: {model_name}"
                continue
            if res.status_code == 404:
                body = res.text[:400].strip()
                raise HTTPException(
                    status_code=502,
                    detail=(
                        f"Claude API 404 at {url}. Check ANTHROPIC_API_URL and account access. "
                        f"Response: {body or 'empty'}"
                    ),
                )
            try:
                res.raise_for_status()
            except httpx.HTTPError as e:
                last_error = str(e)
                continue

            data = res.json()
            content = data.get("content", [])
            if not content or not isinstance(content, list) or "text" not in content[0]:
                last_error = "Claude returned empty content."
                continue
            return content[0]["text"]

    raise HTTPException(
        status_code=502,
        detail=f"Claude API error: {last_error or 'Unknown error. All models failed.'}"
    )


# ---------------------------------------------------------------------------
#  Strict Structured Outputs — tool_choice forces valid JSON from Pydantic schemas
# ---------------------------------------------------------------------------

# Pydantic → JSON Schema conversion for Anthropic tool definitions
def _pydantic_to_tool_schema(name: str, description: str, model: type) -> dict:
    """Convert a Pydantic model class to an Anthropic tool definition."""
    schema = model.model_json_schema() if hasattr(model, "model_json_schema") else model.schema()
    return {
        "name": name,
        "description": description,
        "input_schema": schema,
    }


class StructuredConversionResult(BaseModel):
    """Schema forced onto Claude via tool_choice for conversion results."""
    startups: List[Dict[str, Any]] = Field(default_factory=list, description="Extracted startup records")
    investors: List[Dict[str, Any]] = Field(default_factory=list, description="Extracted investor records")
    mentors: List[Dict[str, Any]] = Field(default_factory=list, description="Extracted mentor records")
    corporates: List[Dict[str, Any]] = Field(default_factory=list, description="Extracted corporate records")
    detected_type: str = Field(default="unknown", description="Detected data type: startup, investor, mentor, corporate, or mixed")
    confidence: float = Field(default=0.0, description="Confidence score 0-1")
    warnings: List[str] = Field(default_factory=list, description="Any parsing warnings")


async def call_anthropic_structured(prompt: str, result_schema: type = StructuredConversionResult) -> dict:
    """
    Call Claude with strict structured output using tool_choice.
    Returns a validated dict matching the Pydantic schema — no regex/bracket-balancing needed.
    
    Uses the Anthropic SDK's tool_choice with strict: true to guarantee valid JSON.
    """
    if not _anthropic_sdk_available or not ANTHROPIC_API_KEY:
        # Fall back to unstructured call + manual parsing
        raw = await call_anthropic(prompt)
        return parse_ollama_response(raw)

    client = _get_anthropic_async_client()
    tool_def = _pydantic_to_tool_schema(
        name="extract_structured_data",
        description="Extract structured startup/investor/mentor/corporate data from the input text.",
        model=result_schema,
    )

    last_error: Optional[str] = None
    for model_name in ANTHROPIC_MODEL_FALLBACKS:
        try:
            message = await client.messages.create(
                model=model_name,
                max_tokens=8192,
                temperature=0.1,
                system=SYSTEM_PROMPT,
                tools=[tool_def],
                tool_choice={"type": "tool", "name": "extract_structured_data"},
                messages=[{"role": "user", "content": prompt}],
            )
            # Extract tool_use block — guaranteed to be valid JSON matching our schema
            for block in message.content:
                if block.type == "tool_use" and block.name == "extract_structured_data":
                    return block.input  # Already a dict matching our schema
            last_error = "No tool_use block in response"
        except Exception as e:
            last_error = str(e)
            continue

    raise HTTPException(
        status_code=502,
        detail=f"Claude structured output failed: {last_error}",
    )


def is_comprehensive_question(question: str) -> bool:
    """Detect if user wants a comprehensive answer (all you know, everything, detailed, etc.)"""
    q_lower = question.lower()
    comprehensive_patterns = [
        "all you know",
        "everything",
        "comprehensive",
        "detailed",
        "full",
        "complete",
        "tell me all",
        "tell me more all",
        "what do you know",
        "what can you tell me",
        "summarize",
        "overview",
        "what is inside",
        "what's inside",
        "what is in",
        "what's in",
        "what does it contain",
        "what does it say",
        "what is the content",
        "what are the contents",
        "just tell",
        "tell what",
        "all about",
        "everything about",
        "from these",
        "from the",
        "from source",
    ]
    return any(pattern in q_lower for pattern in comprehensive_patterns)

def is_raw_text_request(question: str) -> bool:
    """Detect if user is asking for raw/exact text from sources."""
    q_lower = question.lower()
    raw_patterns = [
        "raw text",
        "exact text",
        "verbatim",
        "just text",
        "full text",
        "show the text",
        "give me the text",
        "original text",
        "word for word",
    ]
    return any(pattern in q_lower for pattern in raw_patterns)


def levenshtein_distance(s1: str, s2: str) -> int:
    """
    Calculate Levenshtein distance between two strings.
    Used for fuzzy name matching.
    """
    if len(s1) < len(s2):
        return levenshtein_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    
    previous_row = range(len(s2) + 1)
    for i, c1 in enumerate(s1):
        current_row = [i + 1]
        for j, c2 in enumerate(s2):
            insertions = previous_row[j + 1] + 1
            deletions = current_row[j] + 1
            substitutions = previous_row[j] + (c1 != c2)
            current_row.append(min(insertions, deletions, substitutions))
        previous_row = current_row
    
    return previous_row[-1]


def fuzzy_match_name(query_name: str, document_name: str, max_distance: int = 2) -> bool:
    """
    Check if two names match with fuzzy tolerance.
    Returns True if Levenshtein distance <= max_distance.
    """
    query_lower = query_name.lower().strip()
    doc_lower = document_name.lower().strip()
    
    # Exact match
    if query_lower == doc_lower:
        return True
    
    # Check if one contains the other (for partial matches)
    if query_lower in doc_lower or doc_lower in query_lower:
        return True
    
    # Fuzzy match with Levenshtein distance
    distance = levenshtein_distance(query_lower, doc_lower)
    max_allowed = min(max_distance, len(query_lower) // 3)  # Allow up to 1/3 of length
    return distance <= max_allowed


def extract_proper_nouns(text: str) -> List[str]:
    """
    Extract potential proper nouns (capitalized words, likely names).
    Simple heuristic: words that start with capital letters and are not at sentence start.
    """
    import re
    # Find capitalized words (potential names)
    # Pattern: word boundary, capital letter, followed by lowercase letters
    pattern = r'\b[A-Z][a-z]+\b'
    matches = re.findall(pattern, text)
    # Filter out common words that are always capitalized
    common_caps = {'The', 'A', 'An', 'And', 'Or', 'But', 'In', 'On', 'At', 'To', 'For', 'Of', 'With', 'By'}
    proper_nouns = [m for m in matches if m not in common_caps and len(m) > 2]
    return proper_nouns


def detect_name_in_query(query: str) -> Tuple[bool, List[str]]:
    """
    Detect if query contains person/company names.
    Returns (has_name, list_of_names).
    """
    proper_nouns = extract_proper_nouns(query)
    # Heuristic: if we have 2+ capitalized words together, likely a name
    # Or if we have a capitalized word followed by another capitalized word
    import re
    # Pattern for "FirstName LastName" or "Company Name"
    name_pattern = r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b'
    name_matches = re.findall(name_pattern, query)
    
    all_names = list(set(proper_nouns + name_matches))
    has_name = len(all_names) > 0
    
    return (has_name, all_names)


def classify_query_intent(query: str) -> str:
    """
    Classify query intent: FIND, SUMMARIZE, EXPLAIN, COMPARE, etc.
    """
    q_lower = query.lower()
    
    # FIND intent
    find_patterns = ["find", "search", "locate", "get", "fetch", "retrieve", "show me"]
    if any(pattern in q_lower for pattern in find_patterns):
        return "FIND"
    
    # SUMMARIZE intent
    summarize_patterns = ["summarize", "summarise", "summary", "overview", "brief", "sum up"]
    if any(pattern in q_lower for pattern in summarize_patterns):
        return "SUMMARIZE"
    
    # EXPLAIN intent
    explain_patterns = ["explain", "why", "how does", "how do", "what is", "what are"]
    if any(pattern in q_lower for pattern in explain_patterns):
        return "EXPLAIN"
    
    # COMPARE intent
    compare_patterns = ["compare", "difference", "versus", "vs", "contrast"]
    if any(pattern in q_lower for pattern in compare_patterns):
        return "COMPARE"
    
    # DEFAULT: FIND (most common intent)
    return "FIND"


def extract_source_reference(question: str) -> int | None:
    """
    Extract source number from question (e.g., "source 1", "source [1]", "document 1").
    Returns the source number (1-indexed) or None if not found.
    """
    import re
    q_lower = question.lower()
    # Patterns: "source 1", "source [1]", "document 1", "doc 1", "[1]", etc.
    patterns = [
        r"source\s*\[?(\d+)\]?",
        r"document\s*\[?(\d+)\]?",
        r"doc\s*\[?(\d+)\]?",
        r"\[(\d+)\]",
    ]
    for pattern in patterns:
        match = re.search(pattern, q_lower)
        if match:
            try:
                return int(match.group(1))
            except (ValueError, IndexError):
                continue
    return None


async def extract_search_keywords(user_query: str) -> str:
    """
    Extract core search terms from user query, removing instruction words.
    PRESERVES proper nouns (names) to ensure they're not removed.
    This prevents the vector database from matching on "summarize" instead of "Lily".
    
    Example:
    Input: "Summarize the personal statement for Lily regarding cross-border business"
    Output: "Lily personal statement cross-border business"
    """
    if not user_query:
        return user_query
    
    # Extract proper nouns BEFORE cleaning (to preserve them)
    proper_nouns = extract_proper_nouns(user_query)
    proper_nouns_lower = [pn.lower() for pn in proper_nouns]
    
    if not ANTHROPIC_API_KEY:
        # Fallback: simple keyword extraction using regex
        import re
        # Remove common instruction words
        instruction_patterns = [
            r"\bsummarize\b",
            r"\bsummarise\b",
            r"\btell me about\b",   
            r"\btell me\b",
            r"\bfind\b",
            r"\bsearch for\b",
            r"\bwhat is\b",
            r"\bwhat are\b",
            r"\bwhat does\b",
            r"\bexplain\b",
            r"\bdescribe\b",
            r"\bshow me\b",
            r"\bget\b",
            r"\bfetch\b",
            r"\bretrieve\b",
        ]
        cleaned = user_query
        for pattern in instruction_patterns:
            cleaned = re.sub(pattern, "", cleaned, flags=re.IGNORECASE)
        # Clean up extra spaces
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        
        # Ensure proper nouns are preserved
        if proper_nouns:
            # Add back any proper nouns that might have been removed
            cleaned_lower = cleaned.lower()
            for pn in proper_nouns:
                if pn.lower() not in cleaned_lower:
                    cleaned = f"{pn} {cleaned}".strip()
        
        return cleaned if cleaned else user_query
    
    try:
        headers = {
            "Content-Type": "application/json",
            "x-api-key": ANTHROPIC_API_KEY,
            "anthropic-version": "2023-06-01",
        }
        url = get_anthropic_api_url()
        
        prompt = f"""Extract the core search terms from this user request. 
Remove instructions like "summarize", "tell me about", "find", "explain", "describe".
Focus on Names, Specific Documents, Topics, and Entities.

User Request: "{user_query}"

Output ONLY the search terms. Do not include explanations or additional text."""

        payload = {
            "model": HAIKU_MODEL,  # Use Haiku for cheap, fast extraction
            "max_tokens": 1000,
            "messages": [
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        }
        
        async with httpx.AsyncClient(timeout=5.0) as client:
            res = await client.post(url, headers=headers, json=payload)
            res.raise_for_status()
            data = res.json()
            content = data.get("content", [])
            if isinstance(content, list) and content:
                extracted = content[0].get("text", "").strip()
                if extracted:
                    return extracted
    except Exception as e:
        # Fallback to simple regex extraction on error
        import re
        instruction_patterns = [
            r"\bsummarize\b",
            r"\bsummarise\b",
            r"\btell me about\b",
            r"\btell me\b",
            r"\bfind\b",
            r"\bsearch for\b",
            r"\bwhat is\b",
            r"\bwhat are\b",
            r"\bwhat does\b",
            r"\bexplain\b",
            r"\bdescribe\b",
            r"\bshow me\b",
        ]
        cleaned = user_query
        for pattern in instruction_patterns:
            cleaned = re.sub(pattern, "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return cleaned if cleaned else user_query
    
    return user_query

def is_meta_question(question: str) -> bool:
    """
    Detect if question is about capabilities/system (meta) vs document content.
    Meta questions should be answered with general knowledge.
    Also detects greetings and conversational messages that should always get a response.
    """
    q_lower = question.lower().strip()
    meta_patterns = [
        "what can you do",
        "what could you do",
        "what are you",
        "what do you do",
        "how do you work",
        "what is your purpose",
        "what are your capabilities",
        "what can you help",
        "how can you help",
        "what features",
        "what functionality",
        "what is company assistant",
        "who are you",
        "introduce yourself",
        "what is this",
        "what is this system",
        "what is this platform",
    ]
    # Greetings and conversational messages should ALWAYS get a response
    greeting_patterns = [
        "hello", "hi", "hey", "good morning", "good afternoon", "good evening",
        "thanks", "thank you", "ok", "okay", "sure", "yes", "no", "got it",
        "nice", "great", "awesome", "cool", "bye", "goodbye", "see you",
        "help", "help me", "i need help",
    ]
    # Check if the message is very short (likely conversational)
    words = q_lower.split()
    is_short_conversational = len(words) <= 3
    
    # Check for greetings with punctuation (e.g., "hello?", "hi!")
    q_clean = q_lower.rstrip("?!.,;:")
    is_greeting = any(
        pattern == q_clean or 
        q_clean.startswith(pattern + " ") or 
        q_clean == pattern
        for pattern in greeting_patterns
    )
    
    return (
        any(pattern in q_lower for pattern in meta_patterns) or
        is_greeting or
        (is_short_conversational and any(q_lower.startswith(p) for p in greeting_patterns))
    )


def has_question_overlap(
    question: str,
    sources: List[AskSource],
    previous_messages: List[ChatMessage] | None = None,
    decisions: List[AskDecision] | None = None,
    connections: List[AskConnection] | None = None,
) -> bool:
    """
    Return True if any meaningful keyword from the question appears in the sources, decisions, or connections.
    This is a lightweight guard to avoid hallucinations when sources are unrelated.
    
    IMPORTANT: For follow-up questions, we're more lenient because the user is continuing
    a conversation about a topic that was already validated.
    """
    q_lower = (question or "").lower()
    q_tokens = [t for t in re.split(r"\W+", q_lower) if len(t) > 3]
    
    # CONNECTION-INTENT: If the question is about connections/partnerships, ALWAYS allow
    connection_keywords = ["connect", "connected", "connection", "connections", "partner",
                          "partnership", "partnerships", "introduce", "introduction", "network",
                          "relationship", "relationships", "link", "linked", "help", "collaborate",
                          "synergy", "recommend", "suggest"]
    if any(kw in q_lower for kw in connection_keywords):
        return True
    
    # CRITICAL: If this is a follow-up question (has previous messages), be MORE lenient
    # The user is continuing a conversation, so we should allow it even if the specific
    # words don't match the sources
    is_followup = previous_messages and len(previous_messages) > 0
    
    # Check for follow-up patterns - these should ALWAYS be allowed if there's history
    followup_patterns = ["yes", "more", "tell", "give", "what", "how", "why", "explain", 
                        "elaborate", "detail", "about", "background", "education", "experience"]
    if is_followup and any(pattern in q_lower for pattern in followup_patterns):
        return True
    
    # Extract names from the question - if the question contains a proper name that's in sources, allow it
    name_pattern = r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b'
    question_names = re.findall(name_pattern, question)
    if question_names and sources:
        source_text = " ".join([f"{s.title or ''} {s.file_name or ''} {s.snippet or ''}".strip() for s in sources]).lower()
        for name in question_names:
            if name.lower() in source_text:
                return True
    
    if previous_messages:
        for msg in previous_messages[-5:]:
            q_tokens.extend(
                [t for t in re.split(r"\W+", (msg.content or "").lower()) if len(t) > 3]
            )
    q_tokens = list(dict.fromkeys(q_tokens))
    
    # Check if question is about decisions
    decision_keywords = ["decision", "decisions", "outcome", "outcomes", "action", "actions", 
                         "invest", "investment", "invested", "passed", "declined", "approved",
                         "rejected", "recent decisions", "decision history", "what decisions"]
    if any(keyword in q_lower for keyword in decision_keywords):
        # If decisions are provided, allow the question
        if decisions and len(decisions) > 0:
            return True
    
    # Check sources
    if sources:
        source_text = " ".join(
            [
                f"{s.title or ''} {s.file_name or ''} {s.snippet or ''}".strip()
                for s in sources
            ]
        ).lower()
        if q_tokens and any(token in source_text for token in q_tokens):
            return True
    
    # Check decisions content if provided
    if decisions and q_tokens:
        decision_text = " ".join(
            [
                f"{d.startup_name or ''} {d.action_type or ''} {d.outcome or ''} {d.notes or ''}".strip()
                for d in decisions
            ]
        ).lower()
        if any(token in decision_text for token in q_tokens):
            return True
    
    # Check connections graph if provided
    if connections and q_tokens:
        connection_text = " ".join(
            [
                f"{c.source_company_name or ''} {c.target_company_name or ''} {c.connection_type or ''} {c.ai_reasoning or ''}".strip()
                for c in connections
            ]
        ).lower()
        if any(token in connection_text for token in q_tokens):
            return True
    
    return False


# System prompt for query contextualization (ChatGPT-style)
CONTEXTUALIZE_SYSTEM_PROMPT = """You are an expert search query generator. 
Your task is to rewrite the "Follow-up Question" into a standalone, specific search query based on the "Chat History".

CRITICAL INSTRUCTIONS:
1. **RESOLVE PRONOUNS**: If the Follow-up Question contains "him", "her", "it", "they", replace it with the specific NAME or ENTITY from the *immediately preceding* User/Assistant exchange.
2. **IGNORE OLD TOPICS**: If the chat started about "Giga Energy" but the last message was about "George", and the user asks "tell me about him", you MUST ask about "George". Ignore Giga Energy.
3. **EXTRACT COMPANY/ENTITY NAMES**: If the user's original question mentions a specific company, person, or entity (e.g., "Weego", "Ridelink", "George Goloborodkin"), you MUST include that name in the rewritten query, even if the follow-up question doesn't mention it explicitly.
4. **PRESERVE CONTEXT**: If the original question was "how to make Weego go on IPO" and the follow-up is "you have it, just give the answer", rewrite to "how to make Weego go on IPO" - preserve the original entity and intent.
5. **BE EXPLICIT**: The output must be a full sentence that can be searched in a database, with all entity names explicitly stated.
6. **DO NOT ANSWER**: Do not answer the question. Only output the rewritten query.

---
EXAMPLES (Follow these patterns):

History:
User: Who is Elon Musk?
Assistant: He is the CEO of Tesla.
Follow-up Question: How old is he?
Rewritten Query: How old is Elon Musk?

History:
User: Tell me about the Q3 Report.
Assistant: Here is the Q3 summary...
User: actually, forget that. Who is Sarah Jones?
Assistant: Sarah Jones is the new VP of Sales.
Follow-up Question: Tell me more about her background.
Rewritten Query: Tell me more about Sarah Jones background and resume.

History:
User: Search for Giga Energy.
Assistant: Found 3 documents about Giga Energy.
User: Okay, now look for George Goloborodkin.
Assistant: I found George's resume.
Follow-up Question: What is his email?
Rewritten Query: What is George Goloborodkin's email?
---

Now, rewrite the following:"""


async def rewrite_query_with_llm(question: str, previous_messages: List[ChatMessage] | None = None) -> str:
    """
    ChatGPT-style query rewriting: Use Claude Haiku to contextualize vague follow-up questions.
    
    This is the "Invisible Step" that ChatGPT performs before searching:
    - User: "What do you know about George Goloborodkin?"
    - AI: "George is an intern..."
    - User: "Tell me more about him."
    - System rewrites to: "Tell me more about George Goloborodkin."
    - THEN searches for "George Goloborodkin" (not "him")
    
    This prevents the database from searching for vague words like "him" and returning irrelevant results.
    """
    if not question:
        return question
    
    # If no previous messages, still check if question needs rewriting (might have pronouns from context)
    if not previous_messages:
        previous_messages = []
    
    # Check if question contains pronouns using regex word boundaries (more robust)
    import re
    q_lower = question.lower()
    pronoun_pattern = r'\b(it|its|him|his|her|she|they|them|their|this|that|these|those)\b'
    has_pronouns = bool(re.search(pronoun_pattern, question, re.IGNORECASE))
    
    # Check for affirmative-only responses
    affirmative_only = q_lower.strip() in {
        "yes",
        "yes please",
        "please",
        "ok",
        "okay",
        "sure",
        "go ahead",
        "yep",
        "yeah",
    }
    
    # Check for vague follow-up patterns (case-insensitive, more comprehensive)
    vague_patterns = [
        "tell me more",
        "tell me all",
        "what about",
        "and what",
        "how about",
        "what else",
        "tell more",
        "more about",
        "more details",
        "more info",
        "more information",
        "expand on",
        "elaborate on",
        "go on",
        "all you know",
        "everything about",
        "what's inside",
        "what is inside",
    ]
    has_vague_pattern = any(pattern in q_lower for pattern in vague_patterns)
    
    # ALWAYS rewrite if there's chat history and the question is short/vague
    # This ensures follow-up questions get properly contextualized
    is_short_question = len(question.split()) <= 15  # Increased threshold
    has_chat_history = previous_messages and len(previous_messages) > 0
    
    # Rewrite if: has pronouns, is vague, is affirmative, OR (is short AND has history)
    # ALSO rewrite if question contains "all you know", "everything", "comprehensive" - these need context
    has_comprehensive_intent = any(phrase in q_lower for phrase in ["all you know", "everything", "comprehensive", "all about", "what's inside", "what is inside"])
    
    should_rewrite = has_pronouns or affirmative_only or has_vague_pattern or has_comprehensive_intent or (is_short_question and has_chat_history)
    
    # CRITICAL: Don't rewrite simple greetings or meta questions - they should pass through as-is
    # This prevents "hello?" from being rewritten into an error message
    if is_meta_question(question):
        return question
    
    # If no chat history but question has pronouns or is vague, still try to improve it
    if not should_rewrite and not has_chat_history:
        return question
    
    # If we should rewrite but no history, return as-is (can't contextualize without history)
    if should_rewrite and not has_chat_history:
        return question
    
    # CRITICAL: Extract names from ALL messages (user AND assistant) for robust pronoun resolution
    # The full name might appear in the assistant's response, not the user's question
    all_text = " ".join([m.content for m in previous_messages[-6:]])  # Last 6 messages
    all_names = re.findall(r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b', all_text)  # "FirstName LastName"
    single_names = re.findall(r'\b[A-Z][a-z]{2,}\b', all_text)  # Single capitalized words (potential names)
    # Filter out common words
    common_words = {'The', 'This', 'That', 'Here', 'There', 'What', 'When', 'Where', 'Which', 'Could', 'Would', 'Should', 'Based', 'Found', 'Sorry', 'Please'}
    single_names = [n for n in single_names if n not in common_words]
    
    # CRITICAL: Extract company/entity names from the LAST USER QUESTION (most important context)
    # This ensures we get the right company even if the follow-up is vague
    last_user_question = next((m.content for m in reversed(previous_messages) if m.role == "user"), "")
    company_names_in_last_q = []
    if last_user_question:
        # Look for company names (capitalized words that might be companies)
        # Common company name patterns: "Weego", "Ridelink", "Giga Energy", etc.
        company_patterns = [
            r'\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+)?)\b',  # "Weego" or "Giga Energy"
        ]
        for pattern in company_patterns:
            matches = re.findall(pattern, last_user_question)
            # Filter out common words and keep potential company names
            for match in matches:
                if isinstance(match, tuple):
                    match = match[0] if match[0] else match[1] if len(match) > 1 else ""
                if match and match not in common_words and len(match) > 2:
                    company_names_in_last_q.append(match)
    
    
    # Format the history into a clean dialogue string (last 4 messages for better context)
    history_text = ""
    recent_messages = previous_messages[-4:] if len(previous_messages) >= 4 else previous_messages
    for msg in recent_messages:
        role = "User" if msg.role == "user" else "Assistant"
        # Truncate very long messages to avoid token limits
        content = msg.content[:500] if len(msg.content) > 500 else msg.content
        history_text += f"{role}: {content}\n"
    
    # Construct the final input in the format the system prompt expects
    final_user_content = f"History:\n{history_text}Follow-up Question: {question}\nRewritten Query:"
    
    try:
        if not ANTHROPIC_API_KEY:
            # Fallback to simple replacement if no API key
            # Use names extracted from ALL messages (above)
            if (has_pronouns or affirmative_only) and (all_names or single_names):
                # Prefer full names, fall back to single names
                main_subject = all_names[-1] if all_names else single_names[-1]
                if has_pronouns:
                    rewritten = question
                    for pronoun in ["him", "her", "it", "they", "them", "his", "her", "their", "this", "that"]:
                        rewritten = re.sub(rf'\b{pronoun}\b', main_subject, rewritten, flags=re.IGNORECASE)
                    return rewritten
                if affirmative_only:
                    return f"Tell me more about {main_subject}"
            return question
        
        headers = {
            "Content-Type": "application/json",
            "x-api-key": ANTHROPIC_API_KEY,
            "anthropic-version": "2023-06-01",
        }
        url = get_anthropic_api_url()
        
        # Use system message + user message (ChatGPT-style)
        payload = {
            "model": HAIKU_MODEL,  # Use Haiku for cheap, fast rewriting
            "max_tokens": 100,  # Reduced since we only want the query
            "system": CONTEXTUALIZE_SYSTEM_PROMPT,
            "messages": [
                {
                    "role": "user",
                    "content": final_user_content
                }
            ]
        }
        
        async with httpx.AsyncClient(timeout=10.0) as client:
            res = await client.post(url, headers=headers, json=payload)
            res.raise_for_status()
            data = res.json()
            content = data.get("content", [])
            if isinstance(content, list) and content:
                rewritten = content[0].get("text", "").strip()
                # Clean up any prefixes or extra text the model might add
                rewritten = rewritten.lstrip("Rewritten Query:").lstrip("Query:").lstrip("-").strip()
                # Remove any explanatory text after the question (look for newlines)
                if "\n" in rewritten:
                    rewritten = rewritten.split("\n")[0].strip()
                # Remove trailing periods if it's just a period
                rewritten = rewritten.rstrip(".")
                
                if rewritten and rewritten != question:
                    
                    # CRITICAL VALIDATION: Check if rewritten query contains company/entity names from conversation
                    # This prevents confusion between different companies (e.g., Weego vs Giga Energy)
                    rewritten_lower = rewritten.lower()
                    
                    # Check if rewritten contains any company name from the last user question
                    has_company_from_last_q = False
                    if company_names_in_last_q:
                        has_company_from_last_q = any(name.lower() in rewritten_lower for name in company_names_in_last_q)
                    
                    # Check if rewritten contains any name from history
                    has_name_from_history = any(name.lower() in rewritten_lower for name in all_names) or \
                                           any(name.lower() in rewritten_lower for name in single_names[:5])
                    
                    # If the last user question mentioned a company but the rewritten query doesn't include it, fix it
                    if company_names_in_last_q and not has_company_from_last_q:
                        # For vague follow-ups like "you have it, just give the answer", use the last user question's company
                        if has_vague_pattern or affirmative_only:
                            main_company = company_names_in_last_q[0]
                            # If the rewritten query is too generic, replace it with the original question's intent
                            if len(rewritten.split()) < 5:  # Very short/generic rewrite
                                # Try to extract the intent from the last user question
                                if "ipo" in last_user_question.lower() or "go public" in last_user_question.lower():
                                    rewritten = f"how to make {main_company} go on IPO"
                                elif "about" in last_user_question.lower():
                                    # Extract the topic after "about"
                                    about_match = re.search(r'about\s+([^?]+)', last_user_question.lower())
                                    if about_match:
                                        topic = about_match.group(1).strip()
                                        rewritten = f"tell me about {main_company} {topic}"
                                    else:
                                        rewritten = f"tell me about {main_company}"
                                else:
                                    rewritten = f"{last_user_question} {rewritten}"
                            else:
                                # Insert company name into the rewritten query
                                rewritten = f"{main_company} {rewritten}"
                    
                    # Validate: if original had pronouns and rewritten doesn't contain any name from history, force fix it
                    elif has_pronouns and (all_names or single_names) and not has_name_from_history:
                        # Use the most recent full name, or fall back to single name, or company from last question
                        main_subject = (company_names_in_last_q[0] if company_names_in_last_q else 
                                       (all_names[-1] if all_names else single_names[-1] if single_names else None))
                        if main_subject:
                            rewritten = question
                            for pronoun in ["him", "her", "it", "they", "them", "his", "her", "their", "this", "that"]:
                                rewritten = re.sub(rf'\b{pronoun}\b', main_subject, rewritten, flags=re.IGNORECASE)
                            return rewritten
                    
                    return rewritten
    except Exception as e:
        # Fallback to simple replacement on error
        # Use names extracted from ALL messages (extracted above)
        if has_pronouns and (all_names or single_names):
            main_subject = all_names[-1] if all_names else single_names[-1]
            rewritten = question
            for pronoun in ["him", "her", "it", "they", "them", "his", "her", "their", "this", "that"]:
                rewritten = re.sub(rf'\b{pronoun}\b', main_subject, rewritten, flags=re.IGNORECASE)
            return rewritten
        elif affirmative_only and (all_names or single_names):
            main_subject = all_names[-1] if all_names else single_names[-1]
            return f"Tell me more about {main_subject}"
    
    return question


def resolve_followup_context(question: str, previous_messages: List[ChatMessage] | None = None) -> str:
    """
    Legacy synchronous wrapper. For async contexts, use rewrite_query_with_llm instead.
    """
    if not question or not previous_messages:
        return question
    q_lower = question.lower()
    affirmative_only = q_lower.strip() in {
        "yes",
        "yes please",
        "please",
        "ok",
        "okay",
        "sure",
        "go ahead",
    }
    if affirmative_only:
        last_user = next((m.content for m in reversed(previous_messages) if m.role == "user"), "").strip()
        if last_user:
            return f"{last_user}\n\nFollow-up request: Provide a more complete answer from the available sources."
        return question
    if " it " not in f" {q_lower} " and " it?" not in q_lower and " it." not in q_lower:
        return question
    last_user = next((m.content for m in reversed(previous_messages) if m.role == "user"), "").strip()
    if not last_user:
        return question
    return f"{last_user}\n\nFollow-up question: {question}"


def build_answer_prompt(
    question: str,
    sources: List[AskSource],
    decisions: List[AskDecision],
    previous_messages: List[ChatMessage] = None,
    connections: List[AskConnection] = None,
) -> str:
    is_meta = is_meta_question(question)
    is_comprehensive = is_comprehensive_question(question)
    source_ref = extract_source_reference(question)
    safe_sources = (sources or [])[:ASK_MAX_SOURCES]
    max_snippet_chars = ASK_MAX_SNIPPET_CHARS
    if is_comprehensive:
        max_snippet_chars = max(ASK_MAX_SNIPPET_CHARS, 2000)  # Much larger for comprehensive questions
    source_lines: List[str] = []
    for idx, src in enumerate(safe_sources, start=1):
        title = src.title or src.file_name or f"Source {idx}"
        snippet = (src.snippet or "").strip()
        # If user references a specific source, give it much more space
        if source_ref == idx:
            max_snippet_chars_for_this = max(max_snippet_chars, 3000)  # Even larger for referenced source
        else:
            max_snippet_chars_for_this = max_snippet_chars
        if len(snippet) > max_snippet_chars_for_this:
            snippet = snippet[:max_snippet_chars_for_this] + "…"
        source_lines.append(f"[{idx}] {title}\n{snippet}")

    decision_lines: List[str] = []
    for d in decisions or []:
        summary = " | ".join(
            [part for part in [d.startup_name, d.action_type, d.outcome, d.notes] if part]
        )
        if summary:
            decision_lines.append(f"- {summary}")

    # ── Connections Graph context ──
    connection_lines: List[str] = []
    if connections:
        for conn in connections:
            parts = [
                f"{conn.source_company_name} → {conn.target_company_name}",
                f"type={conn.connection_type}" if conn.connection_type else None,
                f"status={conn.connection_status}" if conn.connection_status else None,
                f"reason: {conn.ai_reasoning[:120]}" if conn.ai_reasoning else None,
                f"notes: {conn.notes[:120]}" if conn.notes else None,
            ]
            connection_lines.append("- " + " | ".join(p for p in parts if p))

    sources_block = "\n\n".join(source_lines) if source_lines else "No sources available."
    decisions_block = "\n".join(decision_lines) if decision_lines else "No decision history available."
    connections_block = "\n".join(connection_lines) if connection_lines else "No company connections in graph yet."
    
    # Build conversation history context - INCLUDE ALL MESSAGES for full context
    # BUT emphasize MOST RECENT messages for pronoun resolution
    conversation_context = ""
    if previous_messages and len(previous_messages) > 0:
        # Include ALL messages (up to 20 to avoid token limits, but prioritize recent)
        recent_messages = previous_messages[-20:] if len(previous_messages) > 20 else previous_messages
        conversation_lines = []
        for i, msg in enumerate(recent_messages):
            role_label = "User" if msg.role == "user" else "Assistant"
            # Truncate very long messages to avoid token limits
            content = msg.content[:1000] if len(msg.content) > 1000 else msg.content
            # Mark most recent messages (last 3) as MOST RECENT
            if i >= len(recent_messages) - 3:
                conversation_lines.append(f"{role_label} (MOST RECENT): {content}")
            else:
                conversation_lines.append(f"{role_label}: {content}")
        
        # Get the last user question explicitly
        last_user_q = next((m.content for m in reversed(recent_messages) if m.role == "user"), "")
        
        conversation_context = f"\n\n=== PREVIOUS CONVERSATION HISTORY (USE THIS TO UNDERSTAND PRONOUNS AND CONTEXT) ===\n" + \
                             f"⚠️ IMPORTANT: Messages marked 'MOST RECENT' are the most recent conversation. If user says 'him', 'her', 'it', check the MOST RECENT messages first!\n" + \
                             f"Last User Question: {last_user_q[:200] if last_user_q else 'N/A'}\n" + \
                             f"\n" + "\n".join(conversation_lines) + "\n=== END OF CONVERSATION HISTORY ===\n"
        # Debug logging - DETAILED
        # Check if history contains names that might be referenced
        all_content = " ".join([m.content for m in recent_messages])
        import re
        names = re.findall(r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b', all_content)
    
    
    is_raw_text = is_raw_text_request(question)
    
    if is_meta:
        # Meta questions: answer with general knowledge about Company Assistant capabilities
        connections_section = ""
        if connection_lines:
            connections_section = f"\n\nYou also have access to a Company Connections Graph with {len(connection_lines)} recorded connections:\n{connections_block}\n\nYou can tell users about these connections when relevant."
        return f"""You are Company Assistant, a VC intelligence system built for investment teams. Answer this question about your capabilities and features.

Question:
{question}

Answer based on what Company Assistant can do:
- Answer questions about uploaded documents (pitch decks, memos, meeting notes)
- Extract structured information from unstructured documents
- Track investment decisions and outcomes
- Provide insights from your fund's knowledge base
- Search across all uploaded sources semantically
- Help with due diligence by finding relevant information quickly
- Show company connections and suggest new ones based on portfolio analysis
- Map relationships between companies (BD, Investment, Knowledge, Partnership, Portfolio)
{connections_section}
Be helpful and specific. Explain what you can do and how you help investment teams.
"""
    else:
        # Document questions: use sources only
        comprehensive_instruction = ""
        if is_comprehensive:
            comprehensive_instruction = "\n\n🚨 CRITICAL: The user is asking for a COMPREHENSIVE answer. This means:\n- Provide ALL available information from the sources about this topic\n- Be thorough, detailed, and include all relevant details\n- Don't summarize or be brief - give the FULL picture\n- Include all sections, data points, qualifications, experiences, and any other information\n- If asked about a person (e.g., 'all you know about George'), include everything: background, education, experience, role, responsibilities, etc.\n- If asked about a document (e.g., 'what's inside source 1'), provide a complete overview of all content\n- Do NOT say 'limited information' or 'very limited' - extract and present EVERYTHING that exists in the sources\n- Be exhaustive, not defensive"
        raw_text_instruction = ""
        if is_raw_text:
            raw_text_instruction = "\n\nIMPORTANT: The user is asking for RAW/EXACT TEXT. Provide the source snippets verbatim (no paraphrasing). If the text is truncated, say so explicitly. Preserve formatting and line breaks when possible."
        
        # Build source reference instruction if user mentions a specific source
        source_ref_instruction = ""
        if source_ref:
            source_ref_instruction = f"\n\nIMPORTANT: The user is asking about SOURCE [{source_ref}]. Focus primarily on that source and provide a COMPREHENSIVE overview of everything in that document. Include all key details, sections, data points, and information from source [{source_ref}]."
        
        # Extract company/entity names from conversation history for explicit highlighting
        company_highlight = ""
        if previous_messages and len(previous_messages) > 0:
            all_conv_text = " ".join([m.content for m in previous_messages[-10:]])
            # Look for company names (capitalized words, common company patterns)
            import re
            company_matches = re.findall(r'\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+)?)\b', all_conv_text)
            # Filter out common words
            common_words = {'The', 'This', 'That', 'Here', 'There', 'What', 'When', 'Where', 'Which', 'Could', 'Would', 'Should', 'Based', 'Found', 'Sorry', 'Please', 'User', 'Assistant'}
            company_names = [m for m in company_matches if m not in common_words and len(m) > 2]
            # Get unique company names, prioritizing those in the last user question
            if company_names:
                unique_companies = list(dict.fromkeys(company_names))  # Preserve order, remove duplicates
                # Check last user question for the most relevant company
                last_user_q = next((m.content for m in reversed(previous_messages) if m.role == "user"), "")
                if last_user_q:
                    companies_in_last_q = [c for c in unique_companies if c.lower() in last_user_q.lower()]
                    if companies_in_last_q:
                        primary_company = companies_in_last_q[0]
                        company_highlight = f"\n\n🎯 **PRIMARY COMPANY/ENTITY FROM CONVERSATION: {primary_company}**\n⚠️ The user's question is about **{primary_company}**. Make sure you answer about **{primary_company}**, NOT about other companies mentioned in the sources (e.g., Giga Energy, Ridelink, etc.). If sources mention other companies, ignore them unless they are directly related to {primary_company}.\n"
        
        # Build the prompt with conversation history at the top
        history_section = conversation_context if conversation_context else "\n\n=== PREVIOUS CONVERSATION HISTORY ===\n(No previous conversation history available)\n=== END OF CONVERSATION HISTORY ===\n"
        
        return f"""You are Company Assistant, a VC intelligence system. You answer questions based on the provided sources and conversation history.

{history_section}
{company_highlight}
🚨 CRITICAL: READ THE CONVERSATION HISTORY ABOVE CAREFULLY BEFORE ANSWERING!

CRITICAL RULES:
1. **MANDATORY: CHECK CONVERSATION HISTORY FIRST**. If the user uses pronouns like "him", "her", "it", "they", "them", "his", "her", "their", "this", "that", "these", "those", you MUST look in the conversation history above to find what they're referring to. The conversation history shows the full context of what was discussed previously.
2. **IF YOU SEE "tell me more about him" AND THE HISTORY SHOWS A PREVIOUS QUESTION ABOUT "George Goloborodkin"**, then "him" = "George Goloborodkin". Use the conversation history to resolve ALL pronouns. **NEVER say "I cannot determine who 'him' refers to" if there is conversation history above - ALWAYS check it first!**
3. **COMPANY NAME RESOLUTION**: If the conversation history shows the user asked about a specific company (e.g., "Weego"), and the current question is vague (e.g., "you have it, just give the answer"), you MUST answer about that specific company (Weego), NOT about other companies mentioned in the sources (e.g., Giga Energy). **DO NOT confuse different companies - if the user asked about Weego, answer about Weego, not Giga Energy!**
3. If the user asks "what's inside", "what is in source X", "all you know", or similar questions about document contents, provide a COMPREHENSIVE and DETAILED answer covering ALL information in the relevant source(s). Do NOT be brief or defensive - give the FULL picture.
4. If the user references a specific source (e.g., "source 1", "source [1]", "document 1"), focus on that source and provide comprehensive details from it. Recognize that [1] refers to the first source, [2] to the second, etc.
5. The sources provided may NOT be relevant to the question. You MUST verify relevance before answering.
6. If the sources DO contain relevant details that DIRECTLY answer the question, provide a thorough, well-structured answer using those details. Be comprehensive and include all relevant information from the sources.
7. If the sources do NOT contain relevant information about the question topic, still try to be helpful. You can say something like "I didn't find specific documents about this topic in your uploaded sources" and then offer to help in other ways, suggest what to search for, or answer based on general knowledge if appropriate. NEVER show bullet-point instructions about uploading documents — that's unhelpful and annoying.
8. Do NOT answer with information that is tangentially related but doesn't actually address the question.
9. If a source talks about a completely different topic (e.g., trading/ATR when asked about a person's resume), you MUST reject it and say you don't have information.
10. Cite sources using [1], [2], etc. for every claim.
11. Do NOT be overly apologetic. If you have information, present it confidently and thoroughly. Only apologize if you truly have no relevant information.
12. When the user asks about something mentioned in the conversation (e.g., "tell me more about him" after discussing George), search the sources for information about that person/entity, even if the current question is vague.

🔗 **CONNECTIONS & PARTNERSHIP RULES**:
13. **ALWAYS CHECK THE CONNECTIONS GRAPH** below when the user mentions ANY company name. If the company has connections in the graph, list them with type and status.
14. If the user explicitly asks about a SPECIFIC company (e.g., "help me understand Ridelink") and the sources do NOT contain information about that company, DO NOT talk about unrelated companies like they are the answer. Instead: (a) state clearly what you found in the Connections Graph for that company, (b) briefly mention that no detailed documents about this company were found, (c) suggest the user enable Web Search or upload relevant documents. Do NOT ramble about other companies as if they were the answer.
15. **Sources labeled [Portfolio company/document: ...]** represent companies in the user's portfolio. ONLY suggest partnerships or connections between them and the asked-about company when the user specifically asks about connections, partnerships, or synergies — not when they ask general questions like "what is this company about."
16. When suggesting connections, explain WHY they could partner based on their industries, stages, markets, or complementary capabilities. Use connection types: BD, Investment, Knowledge, Partnership, Portfolio.
17. If web search results are provided (marked as [WEB] sources), use them to answer the question. Web results are real-time internet data and should be treated as trustworthy supplementary context.{comprehensive_instruction}{raw_text_instruction}{source_ref_instruction}

Answer style:
- Prioritize comprehensive, coherent narrative answers grounded in sources.
- Prefer completeness over brevity when sources list multiple items.
- Do not force a sectioned structure; use paragraphs with bullets only when they improve clarity.
- If information is sparse, still provide the most complete answer possible from the available evidence.
- Use bullet points for responsibilities, qualifications, and scope when asked.
- For comprehensive questions, expand with all relevant details from sources in a flowing, human-like summary.
- For raw text requests, return verbatim snippets with source labels and no paraphrasing.

Question:
{question}

Sources:
{sources_block}

Decision history (optional context):
{decisions_block}

Company Connections Graph (use to answer questions about company relationships, partnerships, and portfolio connections):
{connections_block}

Remember: 
- ALWAYS check the conversation history above to understand pronouns and context.
- ALWAYS check the Connections Graph for the company being asked about. Report ALL known connections.
- If a company is NOT in the Connections Graph and NOT in the sources, say so honestly. Do NOT fabricate information or ramble about unrelated companies.
- ONLY suggest portfolio connections/partnerships when the user asks about connections, partnerships, or synergies — not when they ask "what is this company."
- If [WEB] sources are present, use them confidently to answer questions about companies not found in internal documents.
- Be helpful and concise. Answer the actual question asked.
"""

# Fast model for simple questions (3-5x faster)
HAIKU_MODEL = "claude-haiku-4-20250514"

def is_simple_question(question: str, sources: List[AskSource]) -> bool:
    """
    Detect if question is simple enough for Haiku (3-5x faster).
    Simple = short question, few sources, straightforward answer expected.
    """
    # Simple heuristics:
    # 1. Short question (< 15 words)
    word_count = len(question.split())
    if word_count > 15:
        return False
    
    # 2. Few sources (1-2)
    if len(sources) > 2:
        return False
    
    # 3. Not asking for analysis/comparison
    complex_keywords = ["compare", "analyze", "why", "explain", "evaluate", "assess", "strategy"]
    question_lower = question.lower()
    if any(kw in question_lower for kw in complex_keywords):
        return False
    
    # 4. Asking for facts (what, who, when, where, how much)
    simple_patterns = ["what is", "what are", "who is", "when", "where", "how much", "how many", "tell me about"]
    if any(pattern in question_lower for pattern in simple_patterns):
        return True
    
    # Default: use Sonnet for better quality
    return False


# ── Tool definitions for tool-augmented RAG ──
TOOLS_FOR_ANSWERS = [
    {
        "name": "query_kpis",
        "description": "Query structured KPIs (revenue, growth, valuations) for companies. Use this when users ask about metrics, financials, or comparisons.",
        "input_schema": {
            "type": "object",
            "properties": {
                "company_name": {"type": "string", "description": "Company name to query"},
                "metric_name": {"type": "string", "description": "Metric name (e.g., 'revenue', 'ARR', 'valuation')"},
                "period": {"type": "string", "description": "Time period (e.g., '2024', 'Q1 2024')"},
            },
            "required": ["company_name"],
        },
    },
    {
        "name": "search_graph",
        "description": "Search the knowledge graph for entities and relationships. Use this to find connections, investors, rounds, or related companies.",
        "input_schema": {
            "type": "object",
            "properties": {
                "entity_name": {"type": "string", "description": "Entity name to search for"},
                "entity_type": {"type": "string", "description": "Entity type (company, person, fund, round)"},
                "relation_type": {"type": "string", "description": "Filter by relation type (invested_in, founded, partner_with)"},
                "max_depth": {"type": "integer", "description": "Max traversal depth (default 2)"},
            },
            "required": ["entity_name"],
        },
    },
]

# ── Anthropic native Web Search tool (server-side, Claude handles searching) ──
ANTHROPIC_WEB_SEARCH_TOOL = {
    "type": "web_search_20250305",
    "name": "web_search",
    "max_uses": 5,
}

# Models that support Anthropic native web search
WEB_SEARCH_COMPATIBLE_MODELS = [
    "claude-sonnet-4-20250514",
    "claude-opus-4-20250514",
    "claude-haiku-4-20250514",
    "claude-3-7-sonnet-latest",
    "claude-3-7-sonnet-20250219",
]

def get_web_search_model_list() -> List[str]:
    """Return model list prioritizing web-search-compatible models."""
    # Start with the default model if it's compatible
    models = []
    if ANTHROPIC_MODEL in WEB_SEARCH_COMPATIBLE_MODELS:
        models.append(ANTHROPIC_MODEL)
    # Add other compatible models
    for m in WEB_SEARCH_COMPATIBLE_MODELS:
        if m not in models:
            models.append(m)
    return models


def _append_web_citations(content_blocks, text: str) -> str:
    """Extract web search citations from Claude's response and append source links."""
    citations_seen: Dict[str, str] = {}  # url -> title
    for block in content_blocks:
        block_citations = getattr(block, "citations", None)
        if not block_citations:
            continue
        for cite in block_citations:
            cite_type = getattr(cite, "type", "")
            if cite_type == "web_search_result_location":
                url = getattr(cite, "url", "")
                title = getattr(cite, "title", "")
                if url and url not in citations_seen:
                    citations_seen[url] = title
    if citations_seen:
        text += "\n\n**Web Sources:**"
        for i, (url, title) in enumerate(citations_seen.items(), 1):
            # Format: [number] [Title](url) - only title is visible, clicking opens URL
            text += f"\n[{i}] [{title}]({url})"
    return text


async def execute_tool_call(tool_name: str, tool_input: Dict[str, Any], event_id: Optional[str] = None) -> str:
    """
    Execute a tool call requested by Claude.
    Returns a JSON string with results (or error message).
    """
    if tool_name == "query_kpis":
        # TODO: Add Supabase client to query company_kpis table
        company = tool_input.get("company_name", "")
        metric = tool_input.get("metric_name")
        period = tool_input.get("period")
        return json.dumps({
            "status": "tool_not_implemented",
            "message": f"KPI query for {company} (metric: {metric}, period: {period}) requires Supabase connection. Coming soon.",
        })
    elif tool_name == "search_graph":
        # TODO: Add Supabase client to query kg_entities/kg_edges
        entity = tool_input.get("entity_name", "")
        entity_type = tool_input.get("entity_type")
        return json.dumps({
            "status": "tool_not_implemented",
            "message": f"Graph search for {entity} (type: {entity_type}) requires Supabase connection. Coming soon.",
        })
    return json.dumps({"status": "error", "message": f"Unknown tool: {tool_name}"})


async def call_anthropic_answer(
    prompt: str, 
    question: str = "", 
    sources: List[AskSource] = None, 
    event_id: Optional[str] = None,
    organization_id: Optional[str] = None,
    user_id: Optional[str] = None,
    web_search_enabled: bool = False
) -> str:
    """
    Call Claude to answer a user question with tool-augmented RAG.
    Uses the Anthropic SDK when available (faster, automatic retries, prompt caching) with httpx fallback.
    When web_search_enabled=True, adds Anthropic's native web search tool so Claude can search the internet.
    """
    if not ANTHROPIC_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="ANTHROPIC_API_KEY not set. Set it in the server environment to use Claude."
        )

    # Choose model based on question complexity (Haiku is 3-5x faster)
    use_haiku = question and sources and is_simple_question(question, sources)
    if web_search_enabled:
        # Web search requires compatible models — override model list
        model_list = get_web_search_model_list()
    else:
        model_list = ([HAIKU_MODEL] + ANTHROPIC_MODEL_FALLBACKS) if use_haiku else ANTHROPIC_MODEL_FALLBACKS
    max_tokens = 10000 if use_haiku else ASK_MAX_TOKENS

    # Build tools list — add native web search if enabled
    tools = list(TOOLS_FOR_ANSWERS)
    if web_search_enabled:
        tools.append(ANTHROPIC_WEB_SEARCH_TOOL)

    system_msg = (
        "You are Company Assistant, a VC intelligence system. You answer questions based on "
        "provided sources and the Company Connections Graph. Cite sources with [1], [2], etc. "
        "When a user asks about a company, check the Connections Graph for relationships. "
        "If the user asks WHAT a company IS or what it does, focus on answering that question — "
        "do NOT ramble about unrelated companies. Only suggest connections when the user asks about partnerships or connections. "
        + ("You have web search enabled. Use it to find up-to-date information about companies, markets, or topics "
           "not covered by the provided internal documents. CRITICAL: Prioritize the most recent information (2026, 2025) "
           "when searching. When performing web searches, include terms like '2026', 'latest', 'recent', or 'current' "
           "in your search queries to ensure you get the freshest data. Always cite web sources. " if web_search_enabled else "")
        + "Be helpful, concise, and answer the actual question asked."
    )

    # ── SDK path (preferred) — with tool calling ──
    if _anthropic_sdk_available:
        client = _get_anthropic_async_client()
        last_error: Optional[str] = None
        for model_name in model_list:
            try:
                # First call: Claude may request tools (+ native web search if enabled)
                message = await client.messages.create(
                    model=model_name,
                    max_tokens=max_tokens,
                    temperature=0.5,
                    system=system_msg,
                    messages=[{"role": "user", "content": prompt}],
                    tools=tools,
                )
                
                # Log usage for first call
                if hasattr(message, 'usage') and message.usage:
                    usage = message.usage
                    input_tokens = getattr(usage, 'input_tokens', 0)
                    output_tokens = getattr(usage, 'output_tokens', 0)
                    import asyncio
                    asyncio.create_task(log_api_usage(
                        provider="anthropic",
                        model=model_name,
                        endpoint="/ask",
                        input_tokens=input_tokens,
                        output_tokens=output_tokens,
                        organization_id=organization_id,
                        event_id=event_id,
                        user_id=user_id,
                        request_id=getattr(message, 'id', None),
                    ))
                
                # Handle tool use if present (skip server_tool_use — web search is handled by Anthropic)
                tool_results = []
                for content_block in message.content:
                    block_type = getattr(content_block, "type", "")
                    if block_type == "tool_use":
                        tool_id = getattr(content_block, "id", "")
                        tool_name = getattr(content_block, "name", "")
                        tool_input = getattr(content_block, "input", {})
                        result = await execute_tool_call(tool_name, tool_input, event_id)
                        tool_results.append({
                            "type": "tool_result",
                            "tool_use_id": tool_id,
                            "content": result,
                        })
                    # server_tool_use (web_search) and web_search_tool_result are handled
                    # automatically by the Anthropic API — no action needed from us
                
                # If tools were used, make a follow-up call with results
                if tool_results:
                    follow_up = await client.messages.create(
                        model=model_name,
                        max_tokens=max_tokens,
                        temperature=0.5,
                        system=system_msg,
                        messages=[
                            {"role": "user", "content": prompt},
                            {"role": "assistant", "content": message.content},
                            {"role": "user", "content": tool_results},
                        ],
                    )
                    # Log usage for follow-up call
                    if hasattr(follow_up, 'usage') and follow_up.usage:
                        usage = follow_up.usage
                        input_tokens = getattr(usage, 'input_tokens', 0)
                        output_tokens = getattr(usage, 'output_tokens', 0)
                        import asyncio
                        asyncio.create_task(log_api_usage(
                            provider="anthropic",
                            model=model_name,
                            endpoint="/ask",
                            input_tokens=input_tokens,
                            output_tokens=output_tokens,
                            organization_id=organization_id,
                            event_id=event_id,
                            user_id=user_id,
                            request_id=getattr(follow_up, 'id', None),
                        ))
                    text_parts = [b.text for b in follow_up.content if hasattr(b, "text")]
                    text = "\n".join(text_parts).strip()
                    if text:
                        return text
                else:
                    # No client-side tools used — return direct answer
                    # (web search results are already incorporated by Anthropic server)
                    text_parts = [b.text for b in message.content if hasattr(b, "text")]
                    text = "\n".join(text_parts).strip()
                    
                    # Append web search citation URLs if present
                    if web_search_enabled:
                        text = _append_web_citations(message.content, text)
                    
                    if text:
                        return text
                
                last_error = "Claude returned empty content."
            except anthropic.NotFoundError:
                last_error = f"Model not found: {model_name}"
                continue
            except Exception as e:
                last_error = str(e)
                continue
        raise HTTPException(status_code=503, detail=last_error or "No Claude model available.")

    # ── httpx fallback (tools not supported in httpx path — use SDK) ──
    headers = {
        "Content-Type": "application/json",
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": "2023-06-01",
    }
    url = get_anthropic_api_url()
    default_url = "https://api.anthropic.com/v1/messages"

    last_error = None
    async with httpx.AsyncClient(timeout=120.0) as http_client:
        for model_name in [m for m in model_list if m]:
            payload = {
                "model": model_name,
                "max_tokens": max_tokens,
                "temperature": 0.5,
                "system": system_msg,
                "messages": [{"role": "user", "content": prompt}],
                "tools": tools,
            }

            res = await http_client.post(url, headers=headers, json=payload)
            if res.status_code == 404 and url != default_url:
                res = await http_client.post(default_url, headers=headers, json=payload)
            if res.status_code == 404 and is_model_not_found(res):
                last_error = f"Model not found: {model_name}"
                continue
            if res.status_code >= 400:
                body = res.text[:400].strip()
                raise HTTPException(
                    status_code=502,
                    detail=f"Claude API error ({res.status_code}): {body or 'empty response'}",
                )

            data = res.json() or {}
            content = data.get("content") or []
            text = ""
            if isinstance(content, list) and content:
                # Handle tool_use blocks (would need follow-up in full implementation)
                text_blocks = [b.get("text", "") for b in content if b.get("type") == "text"]
                text = "\n".join(text_blocks).strip()
            elif isinstance(content, str):
                text = content
            if not text:
                raise HTTPException(status_code=502, detail="Claude returned empty content.")
            return text.strip()

    raise HTTPException(status_code=503, detail=last_error or "No Claude model available.")

def normalize_startup_data(data: Dict[str, Any]) -> StartupData:
    """Normalize extracted startup data to match schema"""
    def safe_str(val: Any) -> str:
        return val.strip() if isinstance(val, str) else (str(val).strip() if val is not None else "")
    def safe_int(val: Any, default: int = 0) -> int:
        try:
            if val is None:
                return default
            if isinstance(val, (int, float)):
                return int(val)
            if isinstance(val, str):
                cleaned = re.sub(r'[^\d.]', '', val)
                return int(float(cleaned)) if cleaned else default
            return default
        except Exception:
            return default

    # Handle geoMarkets (accept snake_case + common synonyms like region)
    geo_markets = data.get('geoMarkets', data.get('geo_markets', data.get('region', data.get('regions', data.get('geography', [])))))
    if isinstance(geo_markets, str):
        geo_markets = [g.strip() for g in re.split(r'[,;|]', geo_markets)]
    
    # Handle fundingTarget
    funding_target = data.get('fundingTarget', data.get('funding_target', 0))
    if isinstance(funding_target, str):
        # Extract number from string
        funding_target = re.sub(r'[^\d.]', '', funding_target)
        funding_target = int(float(funding_target)) if funding_target else 0
    funding_target = safe_int(funding_target, 0)
    
    return StartupData(
        companyName=safe_str(data.get('companyName', data.get('company_name', data.get('name', '')))),
        geoMarkets=geo_markets if isinstance(geo_markets, list) else [],
        industry=safe_str(data.get('industry', data.get('sector', data.get('startup_industry', '')))),
        fundingTarget=safe_int(funding_target, 0),
        fundingStage=safe_str(data.get('fundingStage', data.get('funding_stage', data.get('stage', '')))),
        availabilityStatus='present'
    )

def normalize_investor_data(data: Dict[str, Any]) -> InvestorData:
    """Normalize extracted investor data to match schema"""
    def safe_str(val: Any) -> str:
        return val.strip() if isinstance(val, str) else (str(val).strip() if val is not None else "")
    def safe_int(val: Any, default: int = 0) -> int:
        try:
            if val is None:
                return default
            if isinstance(val, (int, float)):
                return int(val)
            if isinstance(val, str):
                cleaned = re.sub(r'[^\d.]', '', val)
                return int(float(cleaned)) if cleaned else default
            return default
        except Exception:
            return default

    # Handle lists
    def parse_list(value):
        if isinstance(value, list):
            return [str(item).strip().strip('[]').strip() for item in value if str(item).strip()]
        if isinstance(value, str):
            # Strip outer ClickUp brackets: "[USA, Europe, SEA]" → "USA, Europe, SEA"
            cleaned = value.strip()
            if cleaned.startswith('[') and cleaned.endswith(']'):
                cleaned = cleaned[1:-1]
            return [item.strip() for item in re.split(r'[,;|]', cleaned) if item.strip()]
        return []
    
    # Handle numbers
    def parse_number(value):
        if value is None:
            return 0
        if isinstance(value, (int, float)):
            return int(value)
        if isinstance(value, str):
            # Handle currency and multipliers
            val = value.upper()
            multiplier = 1
            if 'M' in val or 'MILLION' in val:
                multiplier = 1000000
            elif 'K' in val or 'THOUSAND' in val:
                multiplier = 1000
            digits = re.sub(r'[^\d.]', '', val)
            try:
                return int(float(digits) * multiplier) if digits else 0
            except Exception:
                return 0
        # Fallback for any other type
        return safe_int(value, 0)
    
    geo_focus = parse_list(
        data.get('geoFocus') or
        data.get('geo_focus') or 
        data.get('Main Geographies Targeted (labels)') or  # Platform CSV format
        data.get('Location') or
        data.get('geoMarkets') or 
        data.get('region') or 
        data.get('regions') or 
        data.get('geography') or 
        []
    )
    # Clean ClickUp bracket wrappers from list items: "[USA, Europe]" → ["USA", "Europe"]
    geo_focus = [item.strip().strip('[]').strip() for item in geo_focus if item.strip().strip('[]').strip()]
    
    industry_prefs = parse_list(
        data.get('industryPreferences') or 
        data.get('industry_preferences') or 
        data.get('[BD] Vertical Interests / Vertical (labels)') or  # Platform CSV format
        data.get('[BD] Partner Industry (labels)') or  # Platform CSV format (alternative)
        data.get('Vertical Interests') or
        data.get('industries') or 
        []
    )
    # Clean ClickUp bracket wrappers from list items: "[Fintech, Agnostic]" → ["Fintech", "Agnostic"]
    industry_prefs = [item.strip().strip('[]').strip() for item in industry_prefs if item.strip().strip('[]').strip()]
    
    stage_prefs = parse_list(
        data.get('stagePreferences') or 
        data.get('stage_preferences') or 
        data.get('stages') or 
        []
    )
    
    # Handle cheque/ticket size - may be comma-separated ranges like ">1M, 100K-500K"
    cheque_size_raw = data.get('[INV] Cheque Size (labels)') or data.get('Cheque Size') or data.get('Check Size') or ''
    # Strip ClickUp outer brackets: "[>1M, 100K-500K]" → ">1M, 100K-500K"
    if isinstance(cheque_size_raw, str):
        cheque_size_raw = cheque_size_raw.strip().strip('[]').strip()
    
    if cheque_size_raw and isinstance(cheque_size_raw, str):
        # Split on comma first to handle multi-value entries like ">1M, 100K-500K"
        all_mins = []
        all_maxs = []
        for part in re.split(r'[,;]', cheque_size_raw):
            part = part.strip()
            if not part:
                continue
            if '-' in part:
                range_parts = part.split('-', 1)
                pmin = parse_number(range_parts[0])
                pmax = parse_number(range_parts[1]) if len(range_parts) > 1 else pmin * 10
                if pmin: all_mins.append(pmin)
                if pmax: all_maxs.append(pmax)
            elif '>' in part:
                val = parse_number(part.replace('>', ''))
                if val: all_mins.append(val)
            elif '<' in part:
                val = parse_number(part.replace('<', ''))
                if val: all_maxs.append(val)
            else:
                val = parse_number(part)
                if val:
                    all_mins.append(val)
                    all_maxs.append(val)
        min_ticket = min(all_mins) if all_mins else 0
        max_ticket = max(all_maxs) if all_maxs else (min_ticket * 10 if min_ticket else 0)
    else:
        min_ticket = parse_number(data.get('minTicketSize') or data.get('min_ticket_size') or data.get('minInvestment') or 0)
        max_ticket = parse_number(data.get('maxTicketSize') or data.get('max_ticket_size') or data.get('maxInvestment') or 10000000)
    
    total_slots = safe_int(data.get('totalSlots') or data.get('total_slots') or data.get('slots') or 3, 3)

    # Helper: strip ClickUp bracket wrappers like "[Nikita Ponomarev]" → "Nikita Ponomarev"
    def strip_brackets(val: str) -> str:
        if not val:
            return val
        val = val.strip()
        if val.startswith('[') and val.endswith(']'):
            val = val[1:-1].strip()
        # Also handle multiple names: "[Name1, Name2]" → take first
        if ',' in val:
            parts = [p.strip() for p in val.split(',') if p.strip()]
            val = parts[0] if parts else val
        return val

    # Handle various column name formats from different sources
    firm_name = safe_str(
        data.get('firmName') or 
        data.get('firm_name') or 
        data.get('Investor name') or  # Platform CSV format
        data.get('Task Name') or  # ClickUp export format
        data.get('name') or 
        data.get('firm') or 
        data.get('Company Name') or
        ''
    )
    # Clean up firm name: remove ClickUp task ID artifacts, parenthetical notes
    if firm_name:
        # Strip leading/trailing quotes and whitespace
        firm_name = firm_name.strip().strip('"').strip()
    
    member_name_raw = safe_str(
        data.get('memberName') or 
        data.get('member_name') or 
        data.get('🦅 [INV] Team Member (users)') or  # Platform CSV format (with emoji)
        data.get('[INV] Team Member (users)') or  # Platform CSV format (without emoji)
        data.get('Team Member') or
        data.get('investment_member') or 
        data.get('investorMemberName') or 
        data.get('contactName') or 
        data.get('partnerName') or 
        data.get('personName') or 
        ''
    )
    # ClickUp wraps names in brackets: "[Nikita Ponomarev]" → "Nikita Ponomarev"
    member_name = strip_brackets(member_name_raw)
    
    # If member_name is empty, try Assignee column (ClickUp fallback)
    if not member_name:
        assignee_raw = safe_str(data.get('Assignee') or '')
        member_name = strip_brackets(assignee_raw)
    
    return InvestorData(
        firmName=firm_name,
        memberName=member_name,
        geoFocus=geo_focus,
        industryPreferences=industry_prefs,
        stagePreferences=stage_prefs,
        minTicketSize=min_ticket,
        maxTicketSize=max_ticket,
        totalSlots=total_slots,
        tableNumber=data.get('tableNumber', data.get('table_number', data.get('table', None))),
        availabilityStatus='present'
    )

def normalize_mentor_data(data: Dict[str, Any]) -> MentorData:
    """Normalize extracted mentor data to match schema"""
    def safe_str(val: Any) -> str:
        return val.strip() if isinstance(val, str) else (str(val).strip() if val is not None else "")
    
    def parse_list(value):
        if isinstance(value, list):
            return value
        if isinstance(value, str):
            return [item.strip() for item in re.split(r'[,;|]', value) if item.strip()]
        return []
    
    full_name = safe_str(
        data.get('fullName') or 
        data.get('full_name') or 
        data.get('Full Name') or 
        data.get('name') or 
        ''
    )
    
    email = safe_str(
        data.get('email') or 
        data.get('Email') or 
        ''
    )
    
    linkedin_url = safe_str(
        data.get('linkedinUrl') or 
        data.get('linkedin_url') or 
        data.get('LinkedIn URL') or 
        data.get('LinkedIn') or 
        None
    )
    
    geo_focus = parse_list(
        data.get('geoFocus') or
        data.get('geo_focus') or 
        data.get('Location') or
        data.get('region') or 
        []
    )
    
    industry_prefs = parse_list(
        data.get('industryPreferences') or 
        data.get('industry_preferences') or 
        data.get('Industry Preferences') or
        data.get('industries') or 
        []
    )
    
    expertise_areas = parse_list(
        data.get('expertiseAreas') or 
        data.get('expertise_areas') or 
        data.get('Expertise Areas') or
        data.get('expertise') or 
        []
    )
    
    total_slots = int(data.get('totalSlots') or data.get('total_slots') or data.get('Total Slots') or 3)
    
    return MentorData(
        fullName=full_name,
        email=email,
        linkedinUrl=linkedin_url,
        geoFocus=geo_focus,
        industryPreferences=industry_prefs,
        expertiseAreas=expertise_areas,
        totalSlots=total_slots,
        availabilityStatus='present'
    )

def normalize_corporate_data(data: Dict[str, Any]) -> CorporateData:
    """Normalize extracted corporate data to match schema"""
    def safe_str(val: Any) -> str:
        return val.strip() if isinstance(val, str) else (str(val).strip() if val is not None else "")
    
    def parse_list(value):
        if isinstance(value, list):
            return value
        if isinstance(value, str):
            return [item.strip() for item in re.split(r'[,;|]', value) if item.strip()]
        return []
    
    firm_name = safe_str(
        data.get('firmName') or 
        data.get('firm_name') or 
        data.get('Company Name') or 
        data.get('companyName') or 
        data.get('name') or 
        ''
    )
    
    contact_name = safe_str(
        data.get('contactName') or 
        data.get('contact_name') or 
        data.get('Contact Name') or 
        ''
    )
    
    email = safe_str(
        data.get('email') or 
        data.get('Email') or 
        None
    )
    
    geo_focus = parse_list(
        data.get('geoFocus') or
        data.get('geo_focus') or 
        data.get('Location') or
        data.get('region') or 
        []
    )
    
    industry_prefs = parse_list(
        data.get('industryPreferences') or 
        data.get('industry_preferences') or 
        data.get('Industry Preferences') or
        data.get('industries') or 
        []
    )
    
    partnership_types = parse_list(
        data.get('partnershipTypes') or 
        data.get('partnership_types') or 
        data.get('Partnership Types') or
        []
    )
    
    stages = parse_list(
        data.get('stages') or 
        data.get('Stages') or 
        []
    )
    
    total_slots = int(data.get('totalSlots') or data.get('total_slots') or data.get('Total Slots') or 3)
    
    return CorporateData(
        firmName=firm_name,
        contactName=contact_name,
        email=email,
        geoFocus=geo_focus,
        industryPreferences=industry_prefs,
        partnershipTypes=partnership_types,
        stages=stages,
        totalSlots=total_slots,
        availabilityStatus='present'
    )

async def extract_text_content(file: UploadFile) -> Tuple[str, str]:
    """
    Shared helper to read an uploaded file and extract text_content with best-effort parsing.
    Returns (file_ext, text_content).
    """
    # Read the uploaded file bytes.
    # On some setups, UploadFile may be at EOF (e.g. if something already read the stream),
    # so we retry once after seeking back to the start.
    content = await file.read()
    if not content or len(content) == 0:
        try:
            await file.seek(0)
            content = await file.read()
        except Exception:
            # If seek/read fails, we fall through to the empty-upload guard below.
            pass
    file_ext = file.filename.split('.')[-1].lower() if file.filename else ""
    text_content = None  # Initialize to None

    # Guard: empty upload (common when the browser upload failed or the file is zero bytes)
    if not content or len(content) == 0:
        raise HTTPException(
            status_code=400,
            detail=(
                "Uploaded file is empty (0 bytes). Re-upload the file. "
                "If this keeps happening, the browser may be sending an empty payload or the file may be corrupt. "
                f"filename={file.filename!r}, content_type={getattr(file, 'content_type', None)!r}"
            )
        )

    # Normalize extension and detect by magic bytes only if we don't already recognize a handled type.
    handled_exts = {'pdf', 'xlsx', 'xls', 'csv', 'txt', 'json', 'doc', 'docx'}
    # Fallback detection by magic bytes (override when extension is missing or wrong)
    # NOTE: some PDFs may have leading bytes before the %PDF header, so we search the first chunk.
    head = content[:2048]
    if b'%PDF' in head and file_ext not in handled_exts:
        file_ext = 'pdf'
    elif head[:2] == b'PK' and b'[Content_Types].xml' in head:
        # Peek inside the zip to disambiguate docx vs xlsx
        try:
            import zipfile
            # BytesIO is already imported at top of file
            with zipfile.ZipFile(BytesIO(content)) as z:
                names = set(z.namelist())
                if 'word/document.xml' in names:
                    file_ext = 'docx'
                elif 'xl/workbook.xml' in names or any(n.startswith('xl/') for n in names):
                    file_ext = 'xlsx'
                elif file_ext not in handled_exts:
                    # default fallback
                    file_ext = 'xlsx'
        except Exception:
            # fall back to extension if zip probe fails
            if file_ext not in handled_exts and file_ext not in ['docx', 'xlsx']:
                file_ext = 'xlsx'
    elif content.startswith(b'\xd0\xcf\x11\xe0') and file_ext not in handled_exts:
        # Old Office formats (could be doc or xls); if extension says doc, keep doc, else assume xls
        file_ext = 'doc' if file_ext == 'doc' else 'xls'
    
    # Handle Excel files (XLSX, XLS)
    if file_ext in ['xlsx', 'xls']:
        try:
            if file_ext == 'xlsx':
                try:
                    import openpyxl
                    # BytesIO is already imported at top of file
                    excel_file = BytesIO(content)
                    workbook = openpyxl.load_workbook(excel_file, data_only=True)
                    text_content = ""

                    def _cell_str(val) -> str:
                        if val is None:
                            return ""
                        if hasattr(val, "isoformat"):  # datetime/date
                            return val.isoformat()
                        return str(val).strip()

                    # Extract all sheets so content is never empty when any sheet has data
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        if sheet.max_row == 0 and sheet.max_column == 0:
                            continue
                        text_content += f"\n--- Sheet: {sheet_name} ---\n"
                        # Headers (first row)
                        if sheet.max_row > 0:
                            first_row = next(sheet.iter_rows(min_row=1, max_row=1, values_only=True), None)
                            if first_row:
                                text_content += ",".join(_cell_str(v) for v in first_row) + "\n"
                            # Data rows
                            for row in sheet.iter_rows(min_row=2, values_only=True):
                                if row is None:
                                    continue
                                text_content += ",".join(_cell_str(v) for v in row) + "\n"

                    if not text_content.strip():
                        raise HTTPException(status_code=400, detail="Excel file appears to be empty.")
                except ImportError:
                    raise HTTPException(
                        status_code=500,
                        detail="XLSX support requires openpyxl. Install with: pip install openpyxl"
                    )
            else:  # XLS
                try:
                    import xlrd
                    
                    workbook = xlrd.open_workbook(file_contents=content)
                    text_content = ""
                    
                    # Get the first sheet
                    sheet = workbook.sheet_by_index(0)
                    
                    # Extract headers
                    if sheet.nrows > 0:
                        headers = [str(sheet.cell_value(0, col)) for col in range(sheet.ncols)]
                        text_content += ",".join(headers) + "\n"
                        
                        # Extract data rows
                        for row_idx in range(1, sheet.nrows):
                            row_data = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                            text_content += ",".join(row_data) + "\n"
                    
                    if not text_content.strip():
                        raise HTTPException(status_code=400, detail="Excel file appears to be empty.")
                except ImportError:
                    raise HTTPException(
                        status_code=500,
                        detail="XLS support requires xlrd. Install with: pip install xlrd"
                    )
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse Excel file: {str(e)}")
        
        # After Excel parsing, text_content should be set
        if text_content is None or not text_content.strip():
            raise HTTPException(status_code=500, detail="Excel file parsing completed but no content extracted.")
    
    # Handle DOCX files (prefer python-docx for tables; fallback to raw XML)
    elif file_ext == 'docx':
        try:
            # BytesIO is already imported at top of file
            try:
                from docx import Document  # type: ignore
                doc = Document(BytesIO(content))
                parts = []
                for p in doc.paragraphs:
                    if p.text and p.text.strip():
                        parts.append(p.text)
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text is not None)
                        if row_text.strip():
                            parts.append(row_text)
                text_content = "\n".join(parts)
            except ImportError:
                # Fallback: manual XML strip
                import zipfile
                with zipfile.ZipFile(BytesIO(content)) as z:
                    with z.open('word/document.xml') as doc_xml:
                        raw_xml = doc_xml.read().decode('utf-8', errors='ignore')
                        # Replace paragraph boundaries with newline
                        raw_xml = raw_xml.replace('</w:p>', '\n')
                        # Strip XML tags
                        text_content = re.sub(r'<[^>]+>', '', raw_xml)
            if not text_content or not text_content.strip():
                raise HTTPException(status_code=400, detail="DOCX appears to have no extractable text. If this is a scanned/image DOCX, re-save as PDF or CSV.")
        except KeyError:
            raise HTTPException(status_code=400, detail="DOCX is missing document.xml. Please re-save the file or export to PDF.")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not extract text from DOCX. Please export to PDF or CSV. Error: {str(e)}")

    # Handle DOC files (legacy binary) — best-effort text extraction; if empty, ask user to convert
    elif file_ext == 'doc':
        try:
            # DOC is legacy OLE; we don't depend on heavy converters here.
            # Best-effort: decode as latin-1 ignoring errors and strip control chars.
            raw = content.decode('latin-1', errors='ignore')
            # Remove nulls and most control chars
            cleaned = re.sub(r'[\x00-\x08\x0B-\x1F\x7F]', ' ', raw)
            # Collapse whitespace
            cleaned = re.sub(r'\s+', ' ', cleaned).strip()
            if len(cleaned) < 20:
                raise HTTPException(status_code=400, detail="DOC (legacy Word) has no extractable text. Please re-save as DOCX or PDF and re-upload.")
            text_content = cleaned
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"DOC (legacy Word) could not be read. Please re-save as DOCX or PDF and re-upload. Error: {str(e)}")

    # Handle PDF files — Claude-native ingestion pipeline (V2)
    elif file_ext == 'pdf':
        # Quick sanity-check: real PDFs should contain a %PDF header near the beginning.
        if b"%PDF" not in content[:8192]:
            head_hex = content[:32].hex()
            raise HTTPException(
                status_code=400,
                detail=f'File has ".pdf" extension but does not look like a valid PDF (missing %PDF header). First bytes (hex): {head_hex}'
            )

        # ── Strategy 1: Claude-native PDF document block (best quality) ──
        # Sends the entire PDF as a base64 document content block.
        # Claude "sees" layouts, charts, tables, scanned text — no OCR needed.
        if ANTHROPIC_API_KEY:
            try:
                claude_text = await extract_pdf_with_claude_native(content, max_pages=MAX_PDF_PAGES)
                if claude_text and len(claude_text.strip()) >= 50:
                    return file_ext, claude_text
            except Exception as claude_err:
                pass

        # ── Strategy 2: PyMuPDF fast text extraction (digital PDFs) ──
        import concurrent.futures
        pymupdf_error = None
        try:
            import fitz  # PyMuPDF
            # BytesIO is already imported at top of file
            pdf_stream = BytesIO(content)
            doc = fitz.open(stream=pdf_stream.getvalue(), filetype="pdf")
            page_limit = min(doc.page_count, MAX_PDF_PAGES)

            def _extract_page(i: int) -> str:
                try:
                    page = doc.load_page(i)
                    return f"\n--- Page {i + 1} ---\n{page.get_text('text') or ''}"
                except Exception as e:
                    return f"\n--- Page {i + 1} (error: {e}) ---\n"

            parts: list[str] = []
            with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_PARALLEL_PAGES) as pool:
                futures = [pool.submit(_extract_page, i) for i in range(page_limit)]
                for f in concurrent.futures.as_completed(futures):
                    parts.append(f.result())

            parts.sort(
                key=lambda x: int(m.group(1)) if (m := re.search(r"Page (\d+)", x)) else 0
            )
            text_content = "\n".join(parts).strip()
        except Exception as e:
            pymupdf_error = e
            text_content = None

        if text_content and len(text_content.strip()) >= 50:
            return file_ext, text_content

        # ── Strategy 3: PyPDF2 / pdfplumber fallback ──
        py_pdf2_error = None
        try:
            import PyPDF2
            # BytesIO is already imported at top of file
            pdf_stream = BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_stream)
            text_content = ""
            page_limit = min(len(pdf_reader.pages), MAX_PDF_PAGES)
            for pn in range(page_limit):
                text_content += f"\n--- Page {pn + 1} ---\n"
                extracted = pdf_reader.pages[pn].extract_text()
                if extracted:
                    text_content += extracted

            if not text_content.strip():
                raise ValueError("empty")
        except Exception as e:
            py_pdf2_error = e
            try:
                import pdfplumber
                # BytesIO is already imported at top of file
                pdf_stream = BytesIO(content)
                text_content = ""
                with pdfplumber.open(pdf_stream) as pdf:
                    page_limit = min(len(pdf.pages), MAX_PDF_PAGES)
                    for pn in range(page_limit):
                        text_content += f"\n--- Page {pn + 1} ---\n"
                        pt = pdf.pages[pn].extract_text()
                        if pt:
                            text_content += pt
                if not text_content.strip():
                    raise ValueError("empty after pdfplumber")
            except Exception as plumber_error:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        "Could not extract text from PDF. Claude-native, PyMuPDF, PyPDF2 and pdfplumber all failed.\n"
                        f"Errors: PyMuPDF={pymupdf_error}; PyPDF2={py_pdf2_error}; pdfplumber={plumber_error}\n"
                        "Fix: upload the original XLSX/CSV, or re-export as a searchable PDF."
                    ),
                )
    elif file_ext in ['csv', 'txt', 'json']:
        # Regular text files (CSV, TXT, JSON)
        try:
            text_content = content.decode('utf-8')
            # For JSON files, validate and pretty-print if needed
            if file_ext == 'json':
                try:
                    parsed = json.loads(text_content)
                    # Re-encode as pretty JSON for better readability
                    text_content = json.dumps(parsed, indent=2, ensure_ascii=False)
                except json.JSONDecodeError as e:
                    # Invalid JSON - keep original but log warning
                    pass
        except UnicodeDecodeError:
            # Try other encodings
            try:
                text_content = content.decode('latin-1')
                if file_ext == 'json':
                    try:
                        parsed = json.loads(text_content)
                        text_content = json.dumps(parsed, indent=2, ensure_ascii=False)
                    except json.JSONDecodeError:
                        pass
            except:
                raise HTTPException(status_code=400, detail="Could not decode file. Please ensure it's a text-based file (CSV, TXT, JSON).")
    else:
        # Unsupported format - explicitly prevent binary files from being decoded as text
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file format: {file_ext or 'unknown'}. Supported formats: CSV, TXT, JSON, PDF, XLSX, XLS. If you uploaded an Excel file, make sure openpyxl (for XLSX) or xlrd (for XLS) is installed: pip install openpyxl xlrd"
        )
    
    # Ensure text_content is set before creating request
    if text_content is None:
        raise HTTPException(status_code=500, detail="Internal error: text_content was not set during file processing.")
    
    return file_ext, text_content

@app.get("/")
async def root():
    return {
        "message": "Company Second Brain V2 API",
        "version": "2.0.0",
        "architecture": "Anthropic-native (Claude 3.5/3.7 Sonnet)",
        "endpoints": {
            "/convert": "POST - Convert unstructured data (structured output via tool_choice)",
            "/convert-file": "POST - Convert uploaded file",
            "/ask": "POST - Ask questions about fund documents",
            "/ask/stream": "POST - Streaming Q&A (SSE)",
            "/embed/query": "POST - Generate embeddings",
            "/rerank": "POST - Cross-encoder reranking",
            "/multi-query": "POST - [V2] Multi-query expansion for better retrieval recall",
            "/contextualize-chunk": "POST - [V2] Contextual Retrieval header generation",
            "/agentic-chunk": "POST - [V2] LLM-driven document chunking (agentic splitting by topic)",
            "/graphrag/retrieve": "POST - [V2] LazyGraphRAG retrieval pipeline",
            "/ingest/document-stream": "POST - [V2] SSE streaming document ingestion",
            "/health": "GET - Health check",
            "/models": "GET - List available models",
        },
    }

@app.get("/health")
async def health_check():
    """Check which converter provider is available"""
    # Always include embedding config in health check
    embedding_config = {
        "embeddings_provider": EMBEDDINGS_PROVIDER,
        "voyage_model": VOYAGE_EMBEDDING_MODEL if EMBEDDINGS_PROVIDER == "voyage" else None,
        "voyage_api_key_set": bool(VOYAGE_API_KEY),
        "openai_model": OPENAI_EMBEDDING_MODEL if EMBEDDINGS_PROVIDER == "openai" else None,
        "openai_api_key_set": bool(OPENAI_API_KEY),
        "embedding_dim": EMBEDDING_DIM,
        "rerank_model": RERANK_MODEL,
        "voyage_rerank_model": VOYAGE_RERANK_MODEL,
        "cohere_api_key_set": bool(COHERE_API_KEY),
    }
    
    if ANTHROPIC_API_KEY:
        return {
            "status": "healthy",
            "available": True,
            "provider": "claude",
            "models": [ANTHROPIC_MODEL],
            "error": None,
            "embedding_config": embedding_config,
        }
    if CONVERTER_PROVIDER == "claude":
        return {
            "status": "unhealthy",
            "available": False,
            "provider": "claude",
            "models": [ANTHROPIC_MODEL],
            "error": "ANTHROPIC_API_KEY not set",
            "embedding_config": embedding_config,
        }

    try:
        model_list = await fetch_ollama_model_names()
        return {
            "status": "healthy",
            "available": True,
            "provider": "ollama",
            "models": model_list,
            "ollama_host": OLLAMA_HOST,
            "preferred_model": PREFERRED_OLLAMA_MODEL,
            "embedding_config": embedding_config,
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "available": False,
            "provider": "ollama",
            "error": str(e),
            "ollama_host": OLLAMA_HOST,
            "preferred_model": PREFERRED_OLLAMA_MODEL,
            "embedding_config": embedding_config,
        }

@app.get("/embedding-config")
async def get_embedding_config():
    """Show current embedding configuration for debugging"""
    return {
        "provider": EMBEDDINGS_PROVIDER,
        "voyage": {
            "model": VOYAGE_EMBEDDING_MODEL,
            "api_key_set": bool(VOYAGE_API_KEY),
            "api_key_prefix": VOYAGE_API_KEY[:8] + "..." if VOYAGE_API_KEY else None,
        },
        "openai": {
            "model": OPENAI_EMBEDDING_MODEL,
            "api_key_set": bool(OPENAI_API_KEY),
        },
        "ollama": {
            "model": OLLAMA_EMBEDDING_MODEL,
        },
        "embedding_dim": EMBEDDING_DIM,
        "reranking": {
            "model": VOYAGE_RERANK_MODEL,
            "cohere_fallback_model": RERANK_MODEL,
            "voyage_api_key_set": bool(VOYAGE_API_KEY),
            "cohere_api_key_set": bool(COHERE_API_KEY),
        },
        "supported_voyage_models": [
            "voyage-large-2",        # General purpose, 1536d (current default)
            "voyage-4-large",        # Latest MoE model, 1024d (no 1536 support)
            "voyage-finance-2",      # Best for finance (recommended for VC)
            "voyage-3",              # General model, 1024d
            "voyage-3-lite",         # Faster, cheaper, 512d
            "voyage-code-3",         # Best for code
            "voyage-law-2",          # Best for legal
            "voyage-multilingual-2", # Best for non-English
        ],
    }

@app.get("/models")
async def list_models():
    """List available Ollama models"""
    try:
        model_names = await fetch_ollama_model_names()
        return {
            "models": [{"name": n} for n in model_names],
            "ollama_host": OLLAMA_HOST,
            "preferred_model": PREFERRED_OLLAMA_MODEL,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to list models: {str(e)}")

def try_direct_csv_parse(text_data: str, data_type: Optional[str]) -> Optional[ConversionResponse]:
    """
    Try to parse CSV directly without Ollama if headers are clear.
    Returns ConversionResponse if successful, None if uncertain.
    """
    def normalize_header(value: str) -> str:
        if not value:
            return ""
        lowered = value.lower()
        lowered = re.sub(r'\[.*?\]', ' ', lowered)
        lowered = re.sub(r'[^a-z0-9\s]', ' ', lowered)
        lowered = re.sub(r'\s+', ' ', lowered).strip()
        return lowered

    try:
        # Parse rows with csv.reader to detect header row even if file starts with blank lines
        raw_rows = list(csv.reader(StringIO(text_data)))
        if not raw_rows:
            return None

        header_idx = None
        normalized_rows = []
        for idx, row in enumerate(raw_rows):
            normalized = [normalize_header(cell) for cell in row]
            normalized_rows.append(normalized)
            if not any(normalized):
                continue
            # Look for known header signals
            if (
                ("investor name" in normalized or "firm name" in normalized) and
                any("team member" in h or "member" == h for h in normalized)
            ):
                header_idx = idx
                break
            # ClickUp/Platform format: "task name" + investor signals
            if (
                "task name" in normalized and 
                any("team member" in h or "cheque size" in h or "platform relationship" in h for h in normalized)
            ):
                header_idx = idx
                break
            if ("company name" in normalized and any("funding" in h or "stage" == h for h in normalized)):
                header_idx = idx
                break
            if ("full name" in normalized and "email" in normalized):
                header_idx = idx
                break
            if (("contact name" in normalized or "contact" in normalized) and ("firm name" in normalized or "company name" in normalized)):
                header_idx = idx
                break

        if header_idx is None:
            # Fallback to DictReader using first line as headers
            reader = csv.DictReader(StringIO(text_data))
            rows = list(reader)
        else:
            header = raw_rows[header_idx]
            data_rows = raw_rows[header_idx + 1 :]
            rows = []
            for row in data_rows:
                if not any(cell.strip() for cell in row if isinstance(cell, str)) and not any(row):
                    continue
                record = {header[i]: (row[i] if i < len(row) else "") for i in range(len(header))}
                rows.append(record)

        if not rows:
            return None

        # Check first row to determine type
        first_row = rows[0]
        headers_lower = {normalize_header(k): k for k in first_row.keys() if k}
        
        
        # Helper: check if ANY normalized header contains a given substring
        def any_header_contains(substring: str) -> bool:
            return any(substring in h for h in headers_lower)
        
        # Detect type based on headers (use both exact and substring matching)
        has_mentor_headers = any(h in headers_lower for h in ['full name', 'fullname']) and any(h in headers_lower for h in ['email'])
        has_corporate_headers = any(h in headers_lower for h in ['contact name', 'contactname']) and any(h in headers_lower for h in ['firm name', 'firmname', 'company name', 'companyname'])
        
        # Investor detection: standard headers OR ClickUp/Platform format
        has_investor_headers_standard = (
            any(h in headers_lower for h in ['investor name', 'firm name', 'firmname']) and 
            any(h in headers_lower for h in ['member name', 'membername', 'team member'])
        )
        # ClickUp/Platform CSV: "Task Name" column + investor signals (cheque size, team member, platform relationship type with INV)
        has_investor_headers_clickup = (
            any(h in headers_lower for h in ['task name']) and 
            (any_header_contains('cheque size') or any_header_contains('team member') or any_header_contains('syndicates'))
        )
        # Also detect if headers contain "platform relationship" — strong signal for Platform ClickUp export
        has_platform_relationship = any_header_contains('platform relationship') or any_header_contains('relationship type')
        has_investor_headers = has_investor_headers_standard or has_investor_headers_clickup or (
            any(h in headers_lower for h in ['task name']) and has_platform_relationship
        )
        
        has_startup_headers = any(h in headers_lower for h in ['company name', 'companyname']) and any(h in headers_lower for h in ['funding', 'stage'])
        
        
        startups = []
        investors = []
        mentors = []
        corporates = []
        warnings = []
        
        if has_mentor_headers:
            for row in rows:
                try:
                    mentor = normalize_mentor_data(row)
                    if mentor.fullName and mentor.email:
                        mentors.append(mentor)
                except Exception as e:
                    warnings.append(f"Error parsing mentor row: {str(e)}")
                
            if mentors:
                return ConversionResponse(
                    startups=[],
                    investors=[],
                    mentors=mentors,
                    corporates=[],
                    detectedType="mentor",
                    confidence=0.95,
                    warnings=warnings,
                    errors=[]
                )
        
        elif has_corporate_headers:
            for row in rows:
                try:
                    corp = normalize_corporate_data(row)
                    if corp.firmName and corp.contactName:
                        corporates.append(corp)
                except Exception as e:
                    warnings.append(f"Error parsing corporate row: {str(e)}")
            
            if corporates:
                return ConversionResponse(
                    startups=[],
                    investors=[],
                    mentors=[],
                    corporates=corporates,
                    detectedType="corporate",
                    confidence=0.95,
                    warnings=warnings,
                    errors=[]
                )
        
        elif has_investor_headers:
            skipped_no_name = 0
            for row_idx, row in enumerate(rows):
                try:
                    inv = normalize_investor_data(row)
                    if inv.firmName:
                        if not inv.memberName:
                            inv.memberName = "UNKNOWN"
                        investors.append(inv)
                    else:
                        skipped_no_name += 1
                except Exception as e:
                    warnings.append(f"Error parsing investor row {row_idx}: {str(e)}")
            
            
            if investors:
                return ConversionResponse(
                    startups=[],
                    investors=investors,
                    mentors=[],
                    corporates=[],
                    detectedType="investor",
                    confidence=0.95,
                    warnings=warnings,
                    errors=[]
                )
        
        elif has_startup_headers:
            for row in rows:
                try:
                    startup = normalize_startup_data(row)
                    if startup.companyName:
                        startups.append(startup)
                except Exception as e:
                    warnings.append(f"Error parsing startup row: {str(e)}")
            
            if startups:
                return ConversionResponse(
                    startups=startups,
                    investors=[],
                    mentors=[],
                    corporates=[],
                    detectedType="startup",
                    confidence=0.95,
                    warnings=warnings,
                    errors=[]
                )
        
        # If we couldn't determine type or no data, return None to fall back to Ollama
        return None
        
    except Exception as e:
        return None

@app.post("/convert", response_model=ConversionResponse)
async def convert_data(request: ConversionRequest):
    """
    Convert unstructured data to structured format using Ollama
    """
    # Try direct CSV parsing first if format is CSV
    if request.format == 'csv':
        try:
            direct_result = try_direct_csv_parse(request.data, request.dataType)
            if direct_result:
                return direct_result
        except Exception as e:
            # If direct CSV parsing fails, fall through to Ollama
            pass

    try:
        # Create prompt
        prompt = create_conversion_prompt(request.data, request.dataType)

        # Decide which provider to use
        use_claude = ANTHROPIC_API_KEY is not None and ANTHROPIC_API_KEY.strip() != ""
        if CONVERTER_PROVIDER == "claude" or use_claude:
            # ── V2: Prefer strict structured output (tool_choice) ──
            # This guarantees valid JSON matching our Pydantic schema,
            # eliminating the need for regex/bracket-balancing parsing.
            if _anthropic_sdk_available:
                try:
                    structured = await call_anthropic_structured(prompt)
                    # structured is already a dict matching StructuredConversionResult
                    parsed_data = structured
                except Exception as struct_err:
                    response_text = await call_anthropic(prompt)
                    parsed_data = parse_ollama_response(response_text)
            else:
                response_text = await call_anthropic(prompt)
                try:
                    parsed_data = parse_ollama_response(response_text)
                except Exception:
                    retry_prompt = (
                        "Return ONLY valid JSON. Do not include markdown or explanations. "
                        "Restart the JSON from scratch and ensure all brackets are closed.\n\n"
                        + create_conversion_prompt(request.data, request.dataType)
                    )
                    retry_text = await call_anthropic(retry_prompt)
                    parsed_data = parse_ollama_response(retry_text)
        else:
            # Check models via HTTP API (more reliable than python ollama.list on some setups)
            try:
                available_models = await fetch_ollama_model_names()
            except Exception as e:
                raise HTTPException(
                    status_code=503,
                    detail=f"Ollama not reachable at {OLLAMA_HOST}. Error: {str(e)}"
                )

            if not available_models:
                # Final fallback: attempt python client list directly
                try:
                    client = get_ollama_client()
                    models = client.list()
                    if isinstance(models, dict):
                        for m in models.get("models", []) or []:
                            name = None
                            if isinstance(m, dict) and m.get("name"):
                                name = m["name"]
                            elif isinstance(m, str):
                                name = m
                            if name:
                                available_models.append(name)
                except Exception:
                    pass

            if not available_models:
                raise HTTPException(
                    status_code=503,
                    detail=f"No Ollama models available at {OLLAMA_HOST}. Run: ollama pull llama3.1"
                )

            model_name = pick_model(available_models)

            # Call Ollama
            client = get_ollama_client()
            # Prefer JSON mode if supported by the client/version, otherwise fall back.
            chat_kwargs = dict(
                model=model_name,
                messages=[
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": prompt},
                ],
                options={
                    "temperature": 0.1,  # Low temperature for consistent extraction
                    "num_predict": 4096,  # More headroom to avoid truncated JSON
                },
            )
            try:
                response = client.chat(**chat_kwargs, format="json")
            except TypeError:
                response = client.chat(**chat_kwargs)

            # Extract response content
            response_text = response.get('message', {}).get('content')
            if not isinstance(response_text, str):
                raise HTTPException(status_code=502, detail="Ollama returned empty content. Ensure the model is available and retry.")

            # Parse JSON from response (retry once if model output is truncated/non-JSON)
            try:
                parsed_data = parse_ollama_response(response_text)
            except Exception:
                # Retry with a stricter prompt and higher output budget
                retry_prompt = (
                    "Return ONLY valid JSON. Do not include markdown or explanations. "
                    "Restart the JSON from scratch and ensure all brackets are closed.\n\n"
                    + create_conversion_prompt(request.data, request.dataType)
                )
                retry_kwargs = dict(
                    model=model_name,
                    messages=[
                        {"role": "system", "content": SYSTEM_PROMPT},
                        {"role": "user", "content": retry_prompt},
                    ],
                    options={
                        "temperature": 0.0,
                        "num_predict": 8192,
                    },
                )
                try:
                    retry_res = client.chat(**retry_kwargs, format="json")
                except TypeError:
                    retry_res = client.chat(**retry_kwargs)
                retry_text = retry_res.get("message", {}).get("content")
                if not isinstance(retry_text, str):
                    raise HTTPException(status_code=502, detail="Ollama returned empty content on retry.")
                parsed_data = parse_ollama_response(retry_text)
        
        # If the model returns a wrapper object (common), unwrap it.
        # Supported shapes:
        # - { startups: [...], investors: [...] }
        # - { data: [...] }
        # - { detectedType: "...", investors: [...] } etc.
        if isinstance(parsed_data, dict):
            wrapper = parsed_data

            # Direct "data" wrapper
            if isinstance(wrapper.get("data"), list):
                parsed_data = wrapper["data"]
            # Direct "startups"/"investors" wrapper: bypass generic detection loop
            elif isinstance(wrapper.get("startups"), (list, dict)) or isinstance(wrapper.get("investors"), (list, dict)):
                startups = []
                investors = []
                mentors = []
                corporates = []
                warnings = []
                errors = []

                def ensure_list(v: Any) -> List[Any]:
                    if v is None:
                        return []
                    if isinstance(v, list):
                        return v
                    return [v]

                for item in ensure_list(wrapper.get("startups")):
                    if isinstance(item, dict):
                        try:
                            s = normalize_startup_data(item)
                            if s.companyName:
                                startups.append(s)
                        except Exception as e:
                            errors.append(f"Error processing startup item: {str(e)}")

                for item in ensure_list(wrapper.get("investors")):
                    if isinstance(item, dict):
                        try:
                            inv = normalize_investor_data(item)
                            if inv.firmName:
                                if not inv.memberName:
                                    inv.memberName = "UNKNOWN"
                                    warnings.append(
                                        f"Investor missing memberName; using placeholder 'UNKNOWN' for firm '{inv.firmName}'."
                                    )
                                investors.append(inv)
                        except Exception as e:
                            errors.append(f"Error processing investor item: {str(e)}")

                for item in ensure_list(wrapper.get("mentors")):
                    if isinstance(item, dict):
                        try:
                            mentor = normalize_mentor_data(item)
                            if mentor.fullName and mentor.email:
                                mentors.append(mentor)
                        except Exception as e:
                            errors.append(f"Error processing mentor item: {str(e)}")

                for item in ensure_list(wrapper.get("corporates")):
                    if isinstance(item, dict):
                        try:
                            corp = normalize_corporate_data(item)
                            if corp.firmName and corp.contactName:
                                corporates.append(corp)
                        except Exception as e:
                            errors.append(f"Error processing corporate item: {str(e)}")

                # Determine detected type
                types_found = []
                if startups:
                    types_found.append("startup")
                if investors:
                    types_found.append("investor")
                if mentors:
                    types_found.append("mentor")
                if corporates:
                    types_found.append("corporate")
                
                detected_type = "+".join(types_found) if types_found else "unknown"
                
                if not (startups or investors or mentors or corporates):
                    errors.append("No valid data extracted. Please check the input format and column names.")

                return ConversionResponse(
                    startups=startups,
                    investors=investors,
                    mentors=mentors,
                    corporates=corporates,
                    detectedType=detected_type,
                    confidence=0.8 if (startups or investors) else 0.0,
                    warnings=warnings,
                    errors=errors,
                )

            # Fallback: treat wrapper as a single item
            else:
                parsed_data = [wrapper]
        # Normalize to list if single object
        elif isinstance(parsed_data, dict):
            parsed_data = [parsed_data]
        
        # Convert to structured format
        startups = []
        investors = []
        mentors = []
        corporates = []
        warnings = []
        errors = []
        detected_type = request.dataType or "unknown"
        
        for item in parsed_data:
            # Skip non-dict items to avoid type errors
            if not isinstance(item, dict):
                warnings.append(f"Skipping non-dict item: {item}")
                continue
            try:
                # Auto-detect type if not specified
                if not request.dataType:
                    has_startup_fields = any(k in item for k in ['companyName', 'fundingTarget', 'fundingStage'])
                    has_investor_fields = any(k in item for k in ['firmName', 'minTicketSize', 'maxTicketSize', 'memberName'])
                    has_mentor_fields = any(k in item for k in ['fullName', 'Full Name', 'expertiseAreas', 'Expertise Areas']) and any(k in item for k in ['email', 'Email'])
                    has_corporate_fields = any(k in item for k in ['contactName', 'Contact Name', 'partnershipTypes', 'Partnership Types'])
                    
                    if has_mentor_fields:
                        detected_type = "mentor"
                    elif has_corporate_fields:
                        detected_type = "corporate"
                    elif has_startup_fields and not has_investor_fields:
                        detected_type = "startup"
                    elif has_investor_fields and not has_startup_fields:
                        detected_type = "investor"
                    elif has_startup_fields and has_investor_fields:
                        # Ambiguous - check more indicators
                        if 'companyName' in item and 'fundingTarget' in item:
                            detected_type = "startup"
                        else:
                            detected_type = "investor"
                
                # Convert based on detected type
                if detected_type == "startup" or (not request.dataType and 'companyName' in item):
                    startup = normalize_startup_data(item)
                    if startup.companyName:
                        startups.append(startup)
                elif detected_type == "mentor" or (not request.dataType and any(k in item for k in ['fullName', 'Full Name'])):
                    mentor = normalize_mentor_data(item)
                    if mentor.fullName and mentor.email:
                        mentors.append(mentor)
                    elif mentor.fullName:
                        warnings.append(f"Mentor '{mentor.fullName}' missing email, skipping")
                elif detected_type == "corporate" or (not request.dataType and any(k in item for k in ['contactName', 'Contact Name'])):
                    corporate = normalize_corporate_data(item)
                    if corporate.firmName and corporate.contactName:
                        corporates.append(corporate)
                    elif corporate.firmName:
                        warnings.append(f"Corporate '{corporate.firmName}' missing contact name, skipping")
                elif detected_type == "investor" or (not request.dataType and 'firmName' in item):
                    investor = normalize_investor_data(item)
                    if investor.firmName:
                        # Some sources (esp. PDFs) list only firm names without a specific person.
                        # Don't hard-fail the whole conversion; fill a placeholder and warn.
                        if not investor.memberName:
                            investor.memberName = "UNKNOWN"
                            warnings.append(
                                f"Investor missing memberName; using placeholder 'UNKNOWN' for firm '{investor.firmName}'."
                            )
                        investors.append(investor)
                else:
                    warnings.append(f"Could not determine type for item: {item}")
            except Exception as e:
                errors.append(f"Error processing item: {str(e)}")
        
        if not (startups or investors or mentors or corporates):
            errors.append("No valid data extracted. Please check the input format and column names.")
        
        return ConversionResponse(
            startups=startups,
            investors=investors,
            mentors=mentors,
            corporates=corporates,
            detectedType=detected_type,
            confidence=0.8 if (startups or investors) else 0.0,
            warnings=warnings,
            errors=errors
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

def validate_structured_rows(
    startups: List[StartupData], 
    investors: List[InvestorData],
    mentors: List[MentorData] = None,
    corporates: List[CorporateData] = None
) -> List[str]:
    """
    Row-level validation with explicit row numbers and missing fields.
    """
    errors: List[str] = []

    for idx, s in enumerate(startups, start=1):
        missing = []
        if not s.companyName:
            missing.append("companyName")
        # NOTE: We intentionally do NOT hard-require every field here.
        # Many PDFs / unstructured sources omit fields like stage/geo/ticket size.
        # The UI supports editing later; blocking imports is worse UX.
        if missing:
            errors.append(f"Startup row {idx}: missing {', '.join(missing)}")

    for idx, inv in enumerate(investors, start=1):
        missing = []
        if not inv.firmName:
            missing.append("firmName")
        if not inv.memberName:
            missing.append("memberName")
        # Do not require geoFocus/industryPreferences/stagePreferences/ticket sizes here.
        if missing:
            errors.append(f"Investor row {idx}: missing {', '.join(missing)}")

    for idx, mentor in enumerate(mentors or [], start=1):
        missing = []
        if not mentor.fullName:
            missing.append("fullName")
        if not mentor.email:
            missing.append("email")
        if missing:
            errors.append(f"Mentor row {idx}: missing {', '.join(missing)}")

    for idx, corp in enumerate(corporates or [], start=1):
        missing = []
        if not corp.firmName:
            missing.append("firmName")
        if not corp.contactName:
            missing.append("contactName")
        if missing:
            errors.append(f"Corporate row {idx}: missing {', '.join(missing)}")

    return errors

def build_startup_csv(startups: List[StartupData]) -> str:
    headers = ["company_name", "geo_markets", "industry", "funding_target", "funding_stage"]
    output = StringIO()
    writer = csv.DictWriter(output, fieldnames=headers)
    writer.writeheader()
    for s in startups:
        writer.writerow({
            "company_name": s.companyName or "",
            "geo_markets": "; ".join(s.geoMarkets) if s.geoMarkets else "",
            "industry": s.industry or "",
            "funding_target": s.fundingTarget if s.fundingTarget is not None else "",
            "funding_stage": s.fundingStage or "",
        })
    return output.getvalue()

def build_investor_csv(investors: List[InvestorData]) -> str:
    headers = [
        "firm_name",
        "investment_member",
        "geo_focus",
        "industry_preferences",
        "min_ticket_size",
        "max_ticket_size",
        "total_slots",
        "table_number",
    ]
    output = StringIO()
    writer = csv.DictWriter(output, fieldnames=headers)
    writer.writeheader()
    for inv in investors:
        writer.writerow({
            "firm_name": inv.firmName or "",
            "investment_member": inv.memberName or "",
            "geo_focus": "; ".join(inv.geoFocus) if inv.geoFocus else "",
            "industry_preferences": "; ".join(inv.industryPreferences) if inv.industryPreferences else "",
            "min_ticket_size": inv.minTicketSize if inv.minTicketSize is not None else "",
            "max_ticket_size": inv.maxTicketSize if inv.maxTicketSize is not None else "",
            "total_slots": inv.totalSlots if inv.totalSlots is not None else "",
            "table_number": inv.tableNumber or "",
        })
    return output.getvalue()

@app.post("/convert-file")
async def convert_file(file: UploadFile = File(...), dataType: Optional[str] = None):
    """Convert uploaded file (CSV, text, PDF, XLSX, etc.). Always returns raw_content when text was extracted."""
    file_ext = ""
    text_content = ""
    try:
        file_ext, text_content = await extract_text_content(file)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to extract text from file: {str(e)}")

    # Try structured conversion (startups/investors/mentors/corporates)
    try:
        request = ConversionRequest(
            data=text_content,
            dataType=dataType,
            format=file_ext
        )
        conversion_result = await convert_data(request)
        # Include extracted text for downstream indexing (truncate to control payload size)
        conversion_result.raw_content = text_content[:MAX_MODEL_INPUT_CHARS]

        # Validate critical identifiers, but don't hard-fail if optional fields are missing.
        row_errors = validate_structured_rows(
            conversion_result.startups,
            conversion_result.investors,
            conversion_result.mentors,
            conversion_result.corporates
        )

        # Block only if nothing was extracted
        has_any_data = (
            conversion_result.startups or
            conversion_result.investors or
            conversion_result.mentors or
            conversion_result.corporates
        )
        if not has_any_data:
            # Still return 200 with raw_content so the document can be indexed; add a warning
            conversion_result.errors = (conversion_result.errors or []) + [
                "No valid structured data extracted. Document text was extracted and will be searchable."
            ]
            return conversion_result

        # Surface missing critical fields as warnings so users can import and edit in the UI.
        if row_errors:
            conversion_result.warnings = (conversion_result.warnings or []) + row_errors

        return conversion_result
    except HTTPException:
        raise
    except Exception as e:
        # Structured conversion failed (e.g. Claude timeout), but we have extracted text.
        # Return raw_content so the frontend can store it and index for search.
        return ConversionResponse(
            startups=[],
            investors=[],
            mentors=[],
            corporates=[],
            detectedType=file_ext or "file",
            confidence=0.0,
            warnings=[],
            errors=[f"Structured conversion failed: {str(e)}. Document text was extracted and will be searchable."],
            raw_content=text_content[:MAX_MODEL_INPUT_CHARS],
        )

def parse_google_drive_url(url: str) -> Tuple[str, str]:
    patterns = [
        ("document", r"https?://docs\.google\.com/document/d/([^/]+)"),
        ("presentation", r"https?://docs\.google\.com/presentation/d/([^/]+)"),
        ("spreadsheet", r"https?://docs\.google\.com/spreadsheets/d/([^/]+)"),
        ("drive", r"https?://drive\.google\.com/file/d/([^/]+)"),
    ]
    for kind, pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return kind, match.group(1)
    # Alternate Drive URL pattern: open?id=FILE_ID
    match = re.search(r"[?&]id=([^&]+)", url)
    if match:
        return "drive", match.group(1)
    raise HTTPException(status_code=400, detail="Unsupported Google Drive URL format.")

@app.post("/ingest/clickup", response_model=ClickUpIngestResponse)
async def ingest_clickup(request: ClickUpIngestRequest):
    if not CLICKUP_API_TOKEN:
        raise HTTPException(status_code=503, detail="CLICKUP_API_TOKEN not set on the server.")

    list_id = request.list_id.strip()
    if not list_id:
        raise HTTPException(status_code=400, detail="list_id is required.")

    url = f"https://api.clickup.com/api/v2/list/{list_id}/task"
    params = {"include_closed": "true" if request.include_closed else "false"}
    headers = {"Authorization": CLICKUP_API_TOKEN}

    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(url, headers=headers, params=params)
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code, detail=res.text[:400])
        data = res.json()

    tasks = []
    for task in data.get("tasks", []):
        assignees = [a.get("username") for a in (task.get("assignees") or []) if a.get("username")]
        tasks.append({
            "id": task.get("id"),
            "name": task.get("name"),
            "url": task.get("url"),
            "status": (task.get("status") or {}).get("status"),
            "assignees": assignees,
            "description": task.get("description"),
        })

    return ClickUpIngestResponse(tasks=tasks)

@app.post("/ingest/clickup/lists", response_model=ClickUpListsResponse)
async def list_clickup_lists(request: ClickUpListsRequest):
    if not CLICKUP_API_TOKEN:
        raise HTTPException(status_code=503, detail="CLICKUP_API_TOKEN not set on the server.")

    team_id = request.team_id.strip()
    if not team_id:
        raise HTTPException(status_code=400, detail="team_id is required.")

    url = f"https://api.clickup.com/api/v2/team/{team_id}/list"
    headers = {"Authorization": CLICKUP_API_TOKEN}

    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(url, headers=headers, params={"archived": "false"})
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code, detail=res.text[:400])
        data = res.json()

    lists = []
    for item in data.get("lists", []):
        lists.append({
            "id": item.get("id"),
            "name": item.get("name"),
        })

    return ClickUpListsResponse(lists=lists)

async def stream_anthropic_answer(prompt: str, question: str = "", sources: List[AskSource] = None, event_id: Optional[str] = None, web_search_enabled: bool = False) -> AsyncGenerator[str, None]:
    """
    Stream Claude's response token by token for ChatGPT-like experience with tool-augmented RAG.
    Uses Anthropic SDK streaming when available, httpx SSE fallback otherwise.
    When web_search_enabled=True, adds Anthropic's native web search tool.
    """
    if not ANTHROPIC_API_KEY:
        yield json.dumps({"error": "ANTHROPIC_API_KEY not set"})
        return

    is_comp = is_comprehensive_question(question)
    use_haiku = question and sources and is_simple_question(question, sources) and not is_comp
    if web_search_enabled:
        model_list = get_web_search_model_list()
    else:
        model_list = ([HAIKU_MODEL] + ANTHROPIC_MODEL_FALLBACKS) if use_haiku else ANTHROPIC_MODEL_FALLBACKS
    max_tokens = 8000 if is_comp else (2000 if use_haiku else ASK_MAX_TOKENS)

    # Build tools list — add native web search if enabled
    tools = list(TOOLS_FOR_ANSWERS)
    if web_search_enabled:
        tools.append(ANTHROPIC_WEB_SEARCH_TOOL)

    system_msg = (
        "You are Company Assistant, a VC intelligence system. You answer questions based on "
        "provided sources and the Company Connections Graph. Cite sources with [1], [2], etc. "
        "When a user asks about a company, check the Connections Graph for relationships. "
        "If the user asks WHAT a company IS or what it does, focus on answering that question — "
        "do NOT ramble about unrelated companies. Only suggest connections when the user asks about partnerships or connections. "
        + ("You have web search enabled. Use it to find up-to-date information about companies, markets, or topics "
           "not covered by the provided internal documents. CRITICAL: Prioritize the most recent information (2026, 2025) "
           "when searching. When performing web searches, include terms like '2026', 'latest', 'recent', or 'current' "
           "in your search queries to ensure you get the freshest data. Always cite web sources. " if web_search_enabled else "")
        + "Be helpful, concise, and answer the actual question asked."
    )

    # ── SDK streaming (preferred) — with tool support + native web search ──
    if _anthropic_sdk_available:
        client = _get_anthropic_async_client()
        for model_name in model_list:
            try:
                # Stream: Claude may use tools or web search (native search is handled server-side)
                async with client.messages.stream(
                    model=model_name,
                    max_tokens=max_tokens,
                    temperature=0.1,
                    system=system_msg,
                    messages=[{"role": "user", "content": prompt}],
                    tools=tools,
                ) as stream:
                    tool_uses = []
                    web_search_citations: Dict[str, str] = {}  # url -> title
                    text_content_blocks = []  # Collect full text blocks for citation extraction
                    async for event in stream:
                        # Handle text deltas (streaming text)
                        if event.type == "content_block_delta" and hasattr(event.delta, "type"):
                            if event.delta.type == "text_delta" and hasattr(event.delta, "text"):
                                yield json.dumps({"text": event.delta.text})
                        # Handle content block starts (tools, web search, etc.)
                        elif event.type == "content_block_start" and hasattr(event.content_block, "type"):
                            block_type = getattr(event.content_block, "type", "")
                            if block_type == "tool_use":
                                tool_uses.append({
                                    "id": getattr(event.content_block, "id", ""),
                                    "name": getattr(event.content_block, "name", ""),
                                    "input": getattr(event.content_block, "input", {}),
                                })
                            elif block_type == "server_tool_use":
                                # Native web search is executing — notify frontend
                                yield json.dumps({"status": "🌐 Searching the web..."})
                            elif block_type == "web_search_tool_result":
                                # Collect web search result citations
                                result_content = getattr(event.content_block, "content", [])
                                if isinstance(result_content, list):
                                    for item in result_content:
                                        if hasattr(item, "type") and getattr(item, "type", "") == "web_search_result":
                                            url = getattr(item, "url", "")
                                            title = getattr(item, "title", "")
                                            if url:
                                                web_search_citations[url] = title
                            elif block_type == "text":
                                # Track text blocks for citation extraction
                                text_content_blocks.append(event.content_block)
                                # Also check if this text block has citations attached
                                block_citations = getattr(event.content_block, "citations", None)
                                if block_citations:
                                    for cite in block_citations:
                                        cite_type = getattr(cite, "type", "")
                                        if cite_type == "web_search_result_location":
                                            url = getattr(cite, "url", "")
                                            title = getattr(cite, "title", "")
                                            if url and url not in web_search_citations:
                                                web_search_citations[url] = title
                        # Handle message delta (final message content)
                        elif event.type == "message_delta" and hasattr(event.delta, "stop_reason"):
                            # Message is completing
                            pass
                    
                    # After stream completes, extract citations from text blocks and append web sources
                    # Note: Web search citations are already collected above, but we also check text blocks for citations
                    if web_search_citations:
                        sources_text = "\n\n**Web Sources:**"
                        for i, (url, title) in enumerate(web_search_citations.items(), 1):
                            # Format: [number] [Title](url) - only title is visible, clicking opens URL
                            sources_text += f"\n[{i}] [{title}]({url})"
                        yield json.dumps({"text": sources_text})
                    
                    # If client-side tools were used, execute and stream follow-up
                    if tool_uses:
                        yield json.dumps({"status": "Searching for more information...", "tools": len(tool_uses)})
                        tool_results = []
                        for tool_use in tool_uses:
                            result = await execute_tool_call(tool_use["name"], tool_use["input"], event_id)
                            tool_results.append({
                                "type": "tool_result",
                                "tool_use_id": tool_use["id"],
                                "content": result,
                            })
                        
                        # Stream follow-up response with tool results
                        async with client.messages.stream(
                            model=model_name,
                            max_tokens=max_tokens,
                            temperature=0.1,
                            system=system_msg,
                            messages=[
                                {"role": "user", "content": prompt},
                                {"role": "assistant", "content": [{"type": "tool_use", **tu} for tu in tool_uses]},
                                {"role": "user", "content": tool_results},
                            ],
                        ) as follow_stream:
                            async for event in follow_stream:
                                if event.type == "content_block_delta" and event.delta.type == "text_delta":
                                    yield json.dumps({"text": event.delta.text})
                return  # Success
            except anthropic.NotFoundError:
                continue
            except Exception as e:
                if model_name == model_list[-1]:
                    yield json.dumps({"error": f"All models failed: {str(e)}"})
                continue
        return

    # ── httpx SSE fallback ──
    headers = {
        "Content-Type": "application/json",
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": "2023-06-01",
    }
    url = get_anthropic_api_url()
    default_url = "https://api.anthropic.com/v1/messages"

    http_timeout = 120.0 if web_search_enabled else 60.0
    async with httpx.AsyncClient(timeout=http_timeout) as http_client:
        for model_name in [m for m in model_list if m]:
            payload = {
                "model": model_name,
                "max_tokens": max_tokens,
                "temperature": 0.1,
                "stream": True,
                "system": system_msg,
                "messages": [{"role": "user", "content": prompt}],
                "tools": tools,
            }

            try:
                async with http_client.stream("POST", url, headers=headers, json=payload) as response:
                    if response.status_code == 404 and url != default_url:
                        async with http_client.stream("POST", default_url, headers=headers, json=payload) as retry_response:
                            async for line in retry_response.aiter_lines():
                                if line.startswith("data: "):
                                    data_str = line[6:]
                                    if data_str == "[DONE]":
                                        return
                                    try:
                                        data = json.loads(data_str)
                                        if "delta" in data and "text" in data["delta"]:
                                            yield json.dumps({"text": data["delta"]["text"]})
                                    except json.JSONDecodeError:
                                        continue
                        return

                    if response.status_code >= 400:
                        error_text = await response.aread()
                        error_str = error_text[:200].decode() if isinstance(error_text, bytes) else str(error_text)[:200]
                        if response.status_code == 404 and ("not_found_error" in error_str.lower()):
                            if model_name == model_list[-1]:
                                yield json.dumps({"error": f"All models failed. Error: {error_str}"})
                                return
                            continue
                        yield json.dumps({"error": f"Claude API error ({response.status_code}): {error_str}"})
                        return

                    async for line in response.aiter_lines():
                        if line.startswith("data: "):
                            data_str = line[6:]
                            if data_str == "[DONE]":
                                return
                            try:
                                data = json.loads(data_str)
                                delta = data.get("delta", {})
                                if delta.get("type") == "text_delta" and "text" in delta:
                                    yield json.dumps({"text": delta["text"]})
                                elif "delta" in data and "text" in data["delta"]:
                                    yield json.dumps({"text": data["delta"]["text"]})
                                # Handle web search events in httpx fallback
                                elif data.get("type") == "content_block_start":
                                    cb = data.get("content_block", {})
                                    if cb.get("type") == "server_tool_use" and cb.get("name") == "web_search":
                                        yield json.dumps({"status": "🌐 Searching the web..."})
                                elif "error" in data:
                                    yield json.dumps({"error": str(data["error"])})
                                    return
                            except json.JSONDecodeError:
                                continue
                    return
            except Exception as e:
                if model_name == model_list[-1]:
                    yield json.dumps({"error": f"All models failed: {str(e)}"})
                continue


# ---------------------------------------------------------------------------
#  Multi-Agent RAG — Orchestrator, Graph/KPI Retrieval, Critic
# ---------------------------------------------------------------------------

class OrchestrateRequest(BaseModel):
    question: str
    previous_messages: List[ChatMessage] = Field(default_factory=list, alias="previousMessages")
    model_config = {"populate_by_name": True}

class OrchestrateResponse(BaseModel):
    use_vector: bool = True
    use_graph: bool = False
    use_kpis: bool = False
    use_web: bool = False
    reasoning: str = ""
    sub_queries: Dict[str, str] = Field(default_factory=dict)


@app.post("/orchestrate-query", response_model=OrchestrateResponse)
async def orchestrate_query(request: OrchestrateRequest):
    """
    Multi-Agent RAG — ORCHESTRATOR (Router Agent).
    Analyzes the user question and decides which retrieval agents to activate.
    Returns a routing plan: use_vector, use_graph, use_kpis, use_web.
    """
    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        return OrchestrateResponse(use_vector=True, reasoning="No AI available, defaulting to vector search")

    question = (request.question or "").strip()
    if not question:
        return OrchestrateResponse(use_vector=True, reasoning="Empty question")

    # Build conversation context for the router
    conv_context = ""
    if request.previous_messages:
        recent = request.previous_messages[-5:]
        conv_context = "\nRecent conversation:\n" + "\n".join(
            f"{'User' if m.role == 'user' else 'Assistant'}: {m.content[:200]}"
            for m in recent
        )

    prompt = (
        "You are the router for a VC intelligence RAG system. You do NOT answer the question. "
        "You only decide which retrieval agents to use.\n\n"
        "Available agents:\n"
        "- **vector**: semantic search over uploaded documents (pitch decks, memos, PDFs). "
        "Use for: company description, product, problem/solution, narrative, 'what does X do', 'what did the deck say'.\n"
        "- **graph**: knowledge graph of entities and relationships (companies, people, funds, rounds). "
        "Use for: 'who invested in X', 'connections', 'partners', 'portfolio', 'relationships', 'who is connected to X'.\n"
        "- **kpis**: structured numbers (ARR, valuation, burn, runway, metrics). "
        "Use for: 'ARR of X', 'valuation', 'burn rate', 'metrics', 'numbers', 'financials', comparisons.\n"
        "- **web**: live web search. "
        "Use for: 'latest', 'recent', 'news', 'today', 'current', or when the question is about "
        "public companies/funding/news and internal docs might be stale.\n\n"
        "Rules:\n"
        "- If the question is vague or conversational (greeting, meta), use ONLY vector.\n"
        "- If the question asks about a company's info, use vector.\n"
        "- If the question asks about connections or relationships, use vector + graph.\n"
        "- If the question asks for numbers or metrics, use vector + kpis.\n"
        "- If the question asks to 'compare X and Y', or 'how X differs from Y', or 'difference between X and Y', or 'business model of X vs Y', use vector + graph + kpis.\n"
        "- If the question explicitly asks for 'latest', 'news', or 'current', include web.\n"
        "- For each enabled agent, optionally provide a sub_query (a more targeted query for that agent).\n\n"
        "Output ONLY valid JSON. No markdown, no explanation.\n"
        "Format: {\"use_vector\": bool, \"use_graph\": bool, \"use_kpis\": bool, \"use_web\": bool, "
        "\"reasoning\": \"one sentence\", \"sub_queries\": {\"graph\": \"optional targeted query\", \"kpis\": \"optional targeted query\"}}\n\n"
        f"Question: \"{question}\""
        f"{conv_context}"
    )

    try:
        client = _get_anthropic_async_client()
        message = await client.messages.create(
            model=HAIKU_MODEL,
            max_tokens=300,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text")).strip()

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            json_match = re.search(r'\{[\s\S]*\}', raw)
            if json_match:
                data = json.loads(json_match.group())
            else:
                return OrchestrateResponse(use_vector=True, reasoning="Could not parse router output")

        return OrchestrateResponse(
            use_vector=data.get("use_vector", True),
            use_graph=data.get("use_graph", False),
            use_kpis=data.get("use_kpis", False),
            use_web=data.get("use_web", False),
            reasoning=data.get("reasoning", ""),
            sub_queries=data.get("sub_queries", {}),
        )
    except Exception as e:
        return OrchestrateResponse(use_vector=True, reasoning=f"Router error: {str(e)[:100]}")


class CriticRequest(BaseModel):
    question: str
    answer: str
    context_vector: str = ""
    context_graph: str = ""
    context_kpis: str = ""
    context_web: str = ""


class CriticResponse(BaseModel):
    issues: List[str] = Field(default_factory=list)
    is_grounded: bool = True
    confidence: float = 1.0


@app.post("/critic-check", response_model=CriticResponse)
async def critic_check(request: CriticRequest):
    """
    Multi-Agent RAG — CRITIC (Verifier Agent).
    Checks whether the generated answer is grounded in the provided context.
    Returns a list of unsupported claims.
    """
    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        return CriticResponse(issues=[], is_grounded=True, confidence=1.0)

    answer = (request.answer or "").strip()
    if not answer:
        return CriticResponse(issues=[], is_grounded=True, confidence=1.0)

    context_parts = []
    if request.context_vector:
        context_parts.append(f"[DOCUMENT SOURCES]\n{request.context_vector[:4000]}")
    if request.context_graph:
        context_parts.append(f"[GRAPH / RELATIONSHIPS]\n{request.context_graph[:2000]}")
    if request.context_kpis:
        context_parts.append(f"[KPIs / METRICS]\n{request.context_kpis[:2000]}")
    if request.context_web:
        context_parts.append(f"[WEB SEARCH RESULTS]\n{request.context_web[:2000]}")

    if not context_parts:
        return CriticResponse(issues=[], is_grounded=True, confidence=0.5)

    context_text = "\n\n".join(context_parts)

    prompt = (
        "You are a verifier. You do NOT answer the question. "
        "You only check whether each factual claim in the ASSISTANT ANSWER is supported by the PROVIDED CONTEXT.\n\n"
        "Rules:\n"
        "- If a claim is directly stated in the context, it is supported.\n"
        "- If a claim is a reasonable inference from the context, it is supported.\n"
        "- If a claim has no basis in the context, flag it.\n"
        "- Generic/conversational statements don't need support.\n\n"
        f"PROVIDED CONTEXT:\n{context_text}\n\n"
        f"ASSISTANT ANSWER:\n{answer}\n\n"
        "Output ONLY valid JSON: {\"issues\": [\"string for each unsupported claim\"], "
        "\"is_grounded\": bool, \"confidence\": float 0.0-1.0}.\n"
        "If everything is supported, output {\"issues\": [], \"is_grounded\": true, \"confidence\": 0.95}."
    )

    try:
        client = _get_anthropic_async_client()
        message = await client.messages.create(
            model=HAIKU_MODEL,
            max_tokens=500,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text")).strip()
        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            json_match = re.search(r'\{[\s\S]*\}', raw)
            data = json.loads(json_match.group()) if json_match else {"issues": [], "is_grounded": True, "confidence": 0.5}

        return CriticResponse(
            issues=data.get("issues", []),
            is_grounded=data.get("is_grounded", True),
            confidence=data.get("confidence", 0.5),
        )
    except Exception as e:
        return CriticResponse(issues=[], is_grounded=True, confidence=0.5)


# ---------------------------------------------------------------------------
#  System 2 RAG — Test-Time Compute (Iterative Think → Search → Refine)
# ---------------------------------------------------------------------------

class System2ReflectRequest(BaseModel):
    question: str
    draft_answer: str
    vector_context: str = ""
    graph_context: str = ""
    kpi_context: str = ""
    iteration: int = 0
    max_iterations: int = 3

class System2ReflectResponse(BaseModel):
    needs_more_data: bool = False
    missing_data_types: List[str] = Field(default_factory=list)
    follow_up_queries: List[str] = Field(default_factory=list)
    reasoning: str = ""
    refined_answer: str = ""
    confidence: float = 0.0


@app.post("/system2-reflect", response_model=System2ReflectResponse)
async def system2_reflect(request: System2ReflectRequest):
    """
    System 2 RAG — REFLECTOR (Test-Time Compute).
    Analyzes a draft answer, identifies gaps, and generates follow-up queries.
    The frontend calls this in a loop: draft → reflect → search → draft → reflect ...
    """
    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        return System2ReflectResponse(
            needs_more_data=False,
            refined_answer=request.draft_answer,
            confidence=0.8,
            reasoning="No AI available for reflection",
        )

    context_summary = []
    if request.vector_context:
        context_summary.append(f"[DOCUMENTS]: {request.vector_context[:3000]}")
    if request.graph_context:
        context_summary.append(f"[GRAPH]: {request.graph_context[:1500]}")
    if request.kpi_context:
        context_summary.append(f"[KPIS]: {request.kpi_context[:1500]}")
    context_block = "\n\n".join(context_summary) or "No context provided."

    prompt = f"""You are a reflective VC intelligence analyst performing test-time compute.
You must evaluate a DRAFT ANSWER and decide if more information is needed.

ITERATION: {request.iteration + 1} of {request.max_iterations}

ORIGINAL QUESTION:
{request.question}

AVAILABLE CONTEXT:
{context_block}

DRAFT ANSWER:
{request.draft_answer}

TASK: Analyze the draft and decide:
1. Is the answer complete and well-supported by the context? 
2. Are there specific facts, metrics, or relationships that are mentioned but lack supporting data?
3. If data is missing, what specific follow-up search queries would fill the gaps?

OUTPUT ONLY valid JSON:
{{
  "needs_more_data": bool,
  "missing_data_types": ["vector" | "graph" | "kpis"],
  "follow_up_queries": ["specific search query 1", "specific search query 2"],
  "reasoning": "one paragraph explaining what's missing and why",
  "confidence": float 0.0-1.0 (how confident you are the draft fully answers the question)
}}

RULES:
- If confidence >= 0.85 OR this is the final iteration, set needs_more_data to false
- follow_up_queries should be specific, actionable search queries (not vague)
- missing_data_types tells the frontend WHICH agents to re-query
- Maximum 3 follow-up queries per reflection
- If the draft is good enough, just return needs_more_data: false with high confidence"""

    try:
        client = _get_anthropic_async_client()
        message = await client.messages.create(
            model=HAIKU_MODEL,
            max_tokens=600,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text")).strip()

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            json_match = re.search(r'\{[\s\S]*\}', raw)
            if json_match:
                data = json.loads(json_match.group())
            else:
                return System2ReflectResponse(
                    needs_more_data=False,
                    refined_answer=request.draft_answer,
                    confidence=0.7,
                    reasoning="Could not parse reflection output",
                )

        # Force stop if at max iterations
        if request.iteration + 1 >= request.max_iterations:
            data["needs_more_data"] = False

        return System2ReflectResponse(
            needs_more_data=data.get("needs_more_data", False),
            missing_data_types=data.get("missing_data_types", []),
            follow_up_queries=data.get("follow_up_queries", [])[:3],
            reasoning=data.get("reasoning", ""),
            refined_answer=data.get("refined_answer", ""),
            confidence=data.get("confidence", 0.5),
        )
    except Exception as e:
        return System2ReflectResponse(
            needs_more_data=False,
            refined_answer=request.draft_answer,
            confidence=0.5,
            reasoning=f"Reflection error: {str(e)[:100]}",
        )


class System2RefineRequest(BaseModel):
    question: str
    draft_answer: str
    original_context: str = ""
    additional_context: str = ""
    reflection_reasoning: str = ""
    previous_messages: List[ChatMessage] = Field(default_factory=list, alias="previousMessages")
    model_config = {"populate_by_name": True}


@app.post("/system2-refine/stream")
async def system2_refine_stream(request: System2RefineRequest):
    """
    System 2 RAG — REFINER (streaming).
    Takes the draft answer + additional context from follow-up searches
    and produces an improved, refined final answer via streaming SSE.
    """
    prompt = f"""You are Company Assistant, a VC intelligence synthesis agent performing iterative refinement.

You previously produced a draft answer, but upon reflection, you identified gaps.
Additional data has been retrieved. Your job: produce an IMPROVED, REFINED answer
that integrates ALL available context.

REFLECTION NOTES (what was missing):
{request.reflection_reasoning}

ORIGINAL QUESTION:
{request.question}

ORIGINAL CONTEXT:
{request.original_context[:6000] if request.original_context else "None"}

ADDITIONAL CONTEXT (from follow-up searches):
{request.additional_context[:4000] if request.additional_context else "None"}

PREVIOUS DRAFT:
{request.draft_answer[:3000]}

CITATION RULES:
- Cite document sources with [1], [2], etc.
- Cite graph relationships with [G]
- Cite KPI/metrics data with [K]
- Every factual claim MUST have a citation

REFINEMENT RULES:
- DO NOT repeat the draft verbatim — integrate the new data
- If new data contradicts the draft, prefer the new data
- Be thorough but concise
- Ensure the answer is complete and addresses the original question fully"""

    async def generate():
        try:
            yield f"data: {json.dumps({'status': 'Refining answer with additional context...'})}\n\n"
            async for chunk in stream_anthropic_answer(prompt, question=request.question):
                yield f"data: {chunk}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f'data: {{"error": "{str(e)[:200]}"}}\n\n'

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
            "Access-Control-Allow-Origin": "*",
        },
    )


# ---------------------------------------------------------------------------
#  ColBERT Late Interaction Reranking
# ---------------------------------------------------------------------------

class ColBERTRerankRequest(BaseModel):
    query: str
    documents: List[str]
    doc_ids: List[str] = Field(default_factory=list)
    top_k: int = 10

class ColBERTRerankResponse(BaseModel):
    results: List[dict] = Field(default_factory=list)
    method: str = "colbert_late_interaction"


@app.post("/colbert-rerank", response_model=ColBERTRerankResponse)
async def colbert_rerank(request: ColBERTRerankRequest):
    """
    ColBERT-style Late Interaction Reranking.
    
    Instead of a single embedding per document, this computes token-level
    embeddings for both query and documents, then uses MaxSim (maximum
    similarity) scoring for fine-grained relevance matching.
    
    Flow:
    1. Embed query tokens individually (each word/subword gets its own vector)
    2. Embed document tokens individually  
    3. For each query token, find its max similarity across all doc tokens
    4. Sum the MaxSim scores = final ColBERT score
    
    This captures partial matches that dense single-vector retrieval misses.
    """
    if not request.documents:
        return ColBERTRerankResponse(results=[])

    query = request.query.strip()
    if not query:
        return ColBERTRerankResponse(results=[])

    # Strategy: Use Voyage reranker as the cross-encoder backbone (it already
    # does token-level attention internally), then augment with our own
    # token-level similarity signal for late-interaction scoring.
    
    try:
        # Phase 1: Cross-encoder reranking via Voyage rerank-2.5
        voyage_results = await rerank_with_voyage(
            query=query,
            documents=request.documents,
            top_k=min(request.top_k * 2, len(request.documents)),
        )
        
        # Phase 2: Token-level late-interaction scoring
        # Split query into tokens, embed each, compute MaxSim against doc chunks
        query_tokens = _tokenize_for_colbert(query)
        
        # Embed query tokens as a batch
        query_token_embeddings = await _batch_embed_tokens(query_tokens)
        
        scored_results = []
        for vr in voyage_results:
            doc_idx = vr.get("index", 0)
            doc_text = request.documents[doc_idx] if doc_idx < len(request.documents) else ""
            voyage_score = vr.get("relevance_score", 0.0)
            
            # Compute token-level MaxSim score
            doc_tokens = _tokenize_for_colbert(doc_text[:1000])  # Cap doc length
            if doc_tokens and query_token_embeddings:
                doc_token_embeddings = await _batch_embed_tokens(doc_tokens)
                maxsim_score = _compute_maxsim(query_token_embeddings, doc_token_embeddings)
            else:
                maxsim_score = 0.0
            
            # Combine: weighted fusion of cross-encoder + late-interaction
            combined_score = 0.6 * voyage_score + 0.4 * maxsim_score
            
            scored_results.append({
                "index": doc_idx,
                "doc_id": request.doc_ids[doc_idx] if doc_idx < len(request.doc_ids) else "",
                "relevance_score": round(combined_score, 4),
                "voyage_score": round(voyage_score, 4),
                "maxsim_score": round(maxsim_score, 4),
                "document": doc_text[:200],
            })
        
        # Sort by combined score, take top_k
        scored_results.sort(key=lambda x: x["relevance_score"], reverse=True)
        scored_results = scored_results[:request.top_k]
        
        return ColBERTRerankResponse(results=scored_results, method="colbert_late_interaction")
        
    except Exception as e:
        # Fallback: just use Voyage reranker scores
        try:
            fallback = await rerank_with_voyage(
                query=query,
                documents=request.documents,
                top_k=request.top_k,
            )
            return ColBERTRerankResponse(
                results=[{
                    "index": r.get("index", i),
                    "doc_id": request.doc_ids[r.get("index", i)] if r.get("index", i) < len(request.doc_ids) else "",
                    "relevance_score": r.get("relevance_score", 0.0),
                    "voyage_score": r.get("relevance_score", 0.0),
                    "maxsim_score": 0.0,
                    "document": request.documents[r.get("index", i)][:200] if r.get("index", i) < len(request.documents) else "",
                } for i, r in enumerate(fallback)],
                method="voyage_fallback",
            )
        except Exception as e2:
            return ColBERTRerankResponse(results=[], method="error")


def _tokenize_for_colbert(text: str) -> List[str]:
    """
    Tokenize text into meaningful chunks for ColBERT late interaction.
    Uses word-level + bigram tokens for richer representation.
    """
    text = text.strip()
    if not text:
        return []
    
    # Clean and split into words
    words = re.findall(r'\b\w+\b', text.lower())
    if not words:
        return []
    
    tokens = []
    # Individual words (skip very common stop words for efficiency)
    stop_words = {"the", "a", "an", "is", "are", "was", "were", "be", "been",
                  "being", "have", "has", "had", "do", "does", "did", "will",
                  "would", "could", "should", "may", "might", "shall", "can",
                  "to", "of", "in", "for", "on", "with", "at", "by", "from",
                  "as", "into", "through", "during", "before", "after", "and",
                  "but", "or", "not", "no", "so", "if", "than", "too", "very",
                  "just", "about", "it", "its", "this", "that", "these", "those"}
    
    meaningful_words = [w for w in words if w not in stop_words and len(w) > 1]
    tokens.extend(meaningful_words[:30])  # Cap at 30 word tokens
    
    # Add bigrams for phrase-level matching
    for i in range(len(meaningful_words) - 1):
        if i >= 15:
            break
        tokens.append(f"{meaningful_words[i]} {meaningful_words[i+1]}")
    
    return tokens


async def _batch_embed_tokens(tokens: List[str]) -> List[List[float]]:
    """
    Embed a batch of tokens using Voyage API's batch endpoint.
    Returns a list of embedding vectors, one per token.
    """
    if not tokens or not VOYAGE_API_KEY:
        return []
    
    # Voyage supports batch embedding
    payload = {
        "model": VOYAGE_EMBEDDING_MODEL,
        "input": tokens[:50],  # Cap batch size
        "input_type": "query",
    }
    
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                "https://api.voyageai.com/v1/embeddings",
                headers={
                    "Authorization": f"Bearer {VOYAGE_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
            
            if response.status_code >= 400:
                return []
            
            data = response.json() or {}
            embeddings = []
            for item in data.get("data", []):
                emb = item.get("embedding", [])
                if emb:
                    embeddings.append(emb)
            return embeddings
    except Exception as e:
        return []


def _compute_maxsim(
    query_embeddings: List[List[float]],
    doc_embeddings: List[List[float]],
) -> float:
    """
    Compute ColBERT MaxSim score: for each query token embedding,
    find max cosine similarity with any document token embedding,
    then average across query tokens.
    """
    if not query_embeddings or not doc_embeddings:
        return 0.0
    
    import math
    
    def cosine_sim(a: List[float], b: List[float]) -> float:
        dot = sum(x * y for x, y in zip(a, b))
        norm_a = math.sqrt(sum(x * x for x in a))
        norm_b = math.sqrt(sum(x * x for x in b))
        if norm_a < 1e-9 or norm_b < 1e-9:
            return 0.0
        return dot / (norm_a * norm_b)
    
    total_maxsim = 0.0
    for q_emb in query_embeddings:
        max_sim = max(cosine_sim(q_emb, d_emb) for d_emb in doc_embeddings)
        total_maxsim += max_sim
    
    return total_maxsim / len(query_embeddings)


# ---------------------------------------------------------------------------
#  Multi-Agent RAG — Synthesis prompt builder
# ---------------------------------------------------------------------------

def build_multiagent_answer_prompt(
    question: str,
    vector_context: str,
    graph_context: str,
    kpi_context: str,
    web_context: str,
    decisions: List[AskDecision] = None,
    previous_messages: List[ChatMessage] = None,
    connections: List[AskConnection] = None,
) -> str:
    """
    Build the synthesis prompt for multi-agent RAG.
    Takes pre-retrieved context from each agent and produces a unified prompt.
    """
    # Decisions
    decision_lines: List[str] = []
    for d in decisions or []:
        summary = " | ".join([part for part in [d.startup_name, d.action_type, d.outcome, d.notes] if part])
        if summary:
            decision_lines.append(f"- {summary}")
    decisions_block = "\n".join(decision_lines) if decision_lines else "No decision history."

    # Connections graph
    connection_lines: List[str] = []
    for conn in connections or []:
        parts = [
            f"{conn.source_company_name} -> {conn.target_company_name}",
            f"type={conn.connection_type}" if conn.connection_type else None,
            f"status={conn.connection_status}" if conn.connection_status else None,
            f"reason: {conn.ai_reasoning[:120]}" if conn.ai_reasoning else None,
        ]
        connection_lines.append("- " + " | ".join(p for p in parts if p))
    connections_block = "\n".join(connection_lines) if connection_lines else "No connections."

    # Conversation history
    conversation_context = ""
    if previous_messages:
        recent = previous_messages[-10:]
        conversation_context = (
            "\n\n=== CONVERSATION HISTORY ===\n"
            + "\n".join(
                f"{'User' if m.role == 'user' else 'Assistant'}: {m.content[:500]}"
                for m in recent
            )
            + "\n=== END HISTORY ===\n"
        )

    return f"""You are Company Assistant, a VC intelligence synthesis agent. You received pre-retrieved context from multiple retrieval agents. Your job is to produce a single, coherent, well-cited answer.

CITATION RULES:
- Cite document sources with [1], [2], etc.
- Cite graph relationships with [G]
- Cite KPI/metrics data with [K]
- Cite web search results with [W]
- Every factual claim MUST have a citation
- If no context answers the question, say so clearly

ANSWER RULES:
- Prioritize document sources (highest confidence) over web results
- Integrate graph relationships naturally into the answer
- Present KPIs/metrics with exact numbers when available
- Be thorough but concise — no unnecessary filler
- For comprehensive questions, be exhaustive with all available details
{conversation_context}
Question:
{question}

[DOCUMENTS] (semantic search results from uploaded documents):
{vector_context if vector_context else "No document results."}

[GRAPH] (knowledge graph — entities and relationships):
{graph_context if graph_context else "No graph results."}

[KPIS] (structured metrics and numbers):
{kpi_context if kpi_context else "No KPI data."}

[WEB] (live web search results):
{web_context if web_context else "No web results."}

Decision history:
{decisions_block}

Company Connections:
{connections_block}

Provide the final answer with proper citations."""


@app.post("/ask", response_model=AskResponse)
async def ask_fund(request: AskRequest, auth: AuthContext = Depends(get_auth_context)):
    """Ask a question about fund documents.  ACL: user sees only their scoped sources."""
    question = (request.question or "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="question is required.")
    # ACL filter is available as `acl_metadata_filter(auth)` for upstream vector search

    # Use LLM-based query rewriting for better pronoun resolution
    resolved_question = await rewrite_query_with_llm(question, request.previous_messages or [])
    
    # ALWAYS forward to Claude — let the AI decide if it can answer.
    # The overlap check is kept only as a hint, not a blocker.
    has_overlap = has_question_overlap(
        resolved_question, request.sources or [], request.previous_messages or [], request.decisions or [], request.connections or []
    )
    
    # ALWAYS forward to Claude - even with empty sources, it can answer greetings and general questions
    # Only block if it's explicitly not a meta question AND no overlap AND sources is None (not just empty)
    if not is_meta_question(resolved_question) and not has_overlap and request.sources is None:
        # Only block if sources is explicitly None (not provided at all)
        # Empty array [] should still go through - Claude can answer greetings/conversational questions
        return AskResponse(answer="I couldn't find relevant information for your question. Could you provide more details or try rephrasing?")

    prompt = build_answer_prompt(resolved_question, request.sources, request.decisions, request.previous_messages, request.connections)
    # Extract event_id from sources if available (first source's metadata)
    event_id = None
    if request.sources and len(request.sources) > 0:
        # Try to extract from first source's metadata or document_id prefix
        first_source = request.sources[0]
        if hasattr(first_source, "metadata") and isinstance(first_source.metadata, dict):
            event_id = first_source.metadata.get("event_id")
    # Extract auth context for cost tracking
    organization_id = auth.org_id if auth.org_id else None
    user_id = auth.user_id if auth.user_id != "anonymous" else None
    answer = await call_anthropic_answer(
        prompt, 
        question=resolved_question, 
        sources=request.sources, 
        event_id=event_id,
        organization_id=organization_id,
        user_id=user_id,
        web_search_enabled=request.web_search_enabled
    )
    return AskResponse(answer=answer)


@app.post("/ask/stream")
async def ask_fund_stream(request: AskRequest, auth: AuthContext = Depends(get_auth_context)):
    """
    Streaming endpoint for ChatGPT-like gradual text typing.
    Returns Server-Sent Events (SSE) stream.
    ACL: user sees only their scoped sources.
    """
    try:
        question = (request.question or "").strip()
        if not question:
            raise HTTPException(status_code=400, detail="question is required.")

        # Use LLM-based query rewriting for better pronoun resolution
        
        resolved_question = await rewrite_query_with_llm(question, request.previous_messages or [])
        
        # Check overlap - but ONLY block if there are literally no sources at all
        has_overlap = has_question_overlap(
            resolved_question, request.sources or [], request.previous_messages or [], request.decisions or [], request.connections or []
        )
        
        # ALWAYS forward to Claude - even with empty sources, it can answer greetings and general questions
        # Only block if it's explicitly not a meta question AND no overlap AND sources is None (not just empty)
        if not is_meta_question(resolved_question) and not has_overlap and request.sources is None:
            no_info_message = "I couldn't find relevant information for your question. Could you provide more details or try rephrasing?"
            async def generate_empty():
                yield f"data: {json.dumps({'text': no_info_message})}\n\n"
                yield "data: [DONE]\n\n"
            return StreamingResponse(
                generate_empty(),
                media_type="text/event-stream",
                headers={
                    "Cache-Control": "no-cache",
                    "Connection": "keep-alive",
                    "X-Accel-Buffering": "no",
                    "Access-Control-Allow-Origin": "*",
                    "Access-Control-Allow-Methods": "POST, OPTIONS",
                    "Access-Control-Allow-Headers": "*",
                },
            )

        # Handle previous_messages safely (default to empty list if None)
        previous_messages = request.previous_messages or []
        
        # Debug logging - DETAILED
        if previous_messages:
            # Print all messages for debugging
            pass

        prompt = build_answer_prompt(resolved_question, request.sources or [], request.decisions or [], previous_messages, request.connections or [])
        
        async def generate():
            try:
                # ── Real-Time SSE Status Updates ──
                # Push status updates so the UI doesn't time out during 30s+ reasoning
                yield f"data: {json.dumps({'ping': True})}\n\n"
                yield f"data: {json.dumps({'status': 'Analyzing your question...'})}\n\n"

                if request.sources:
                    n = len(request.sources)
                    # Show a friendly count: distinguish company cards from documents
                    card_count = sum(1 for s in request.sources if s.title and "Company card:" in s.title)
                    doc_count = n - card_count
                    parts = []
                    if card_count > 0:
                        parts.append(f"{card_count} company card{'s' if card_count != 1 else ''}")
                    if doc_count > 0:
                        parts.append(f"{doc_count} document{'s' if doc_count != 1 else ''}")
                    status = f"Reading {' and '.join(parts)}..." if parts else f"Reading {n} source(s)..."
                    yield f"data: {json.dumps({'status': status})}\n\n"

                yield f"data: {json.dumps({'status': 'Generating response...'})}\n\n"

                # Extract event_id from sources if available
                event_id = None
                if request.sources and len(request.sources) > 0:
                    first_source = request.sources[0]
                    if hasattr(first_source, "metadata") and isinstance(first_source.metadata, dict):
                        event_id = first_source.metadata.get("event_id")

                async for chunk in stream_anthropic_answer(prompt, question=question, sources=request.sources or [], event_id=event_id, web_search_enabled=request.web_search_enabled):
                    yield f"data: {chunk}\n\n"
                yield "data: [DONE]\n\n"
            except Exception as e:
                import traceback
                error_trace = traceback.format_exc()
                error_msg = str(e)[:500]
                yield f'data: {{"error": "{error_msg}"}}\n\n'
        
        return StreamingResponse(
            generate(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",  # Disable nginx buffering
                "Access-Control-Allow-Origin": "*",  # Explicit CORS header for streaming
                "Access-Control-Allow-Methods": "POST, OPTIONS",
                "Access-Control-Allow-Headers": "*",
            }
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        raise HTTPException(
            status_code=500,
            detail=f"Streaming error: {str(e)[:200]}"
        )

def normalize_embedding(embedding: List[float]) -> List[float]:
    """Ensure embeddings are a fixed size for pgvector."""
    if not embedding:
        return []
    if len(embedding) == EMBEDDING_DIM:
        return embedding
    if len(embedding) > EMBEDDING_DIM:
        return embedding[:EMBEDDING_DIM]
    # Pad with zeros if smaller than expected
    return embedding + [0.0] * (EMBEDDING_DIM - len(embedding))



@app.post("/rerank", response_model=RerankResponse)
async def rerank_documents(request: RerankRequest):
    """
    Optional cross-encoder reranking for hybrid search results.
    If no reranker API key is configured, returns the original order with score=0.
    """
    if not request.documents:
        return RerankResponse(results=[])

    if not COHERE_API_KEY:
        return RerankResponse(
            results=[RerankResult(id=d.id, score=0.0) for d in request.documents]
        )

    payload = {
        "model": RERANK_MODEL,
        "query": request.query,
        "documents": [d.text for d in request.documents],
        "top_n": request.top_n or len(request.documents),
    }
    headers = {
        "Authorization": f"Bearer {COHERE_API_KEY}",
        "Content-Type": "application/json",
    }

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post("https://api.cohere.ai/v1/rerank", json=payload, headers=headers)
        response.raise_for_status()
        data = response.json()
        results = []
        for item in data.get("results", []):
            idx = item.get("index")
            score = item.get("relevance_score", 0.0)
            if idx is None or idx >= len(request.documents):
                continue
            results.append(RerankResult(id=request.documents[idx].id, score=score))
        return RerankResponse(results=results)
    except Exception:
        # Fail closed: return original order if rerank fails
        return RerankResponse(
            results=[RerankResult(id=d.id, score=0.0) for d in request.documents]
        )


async def generate_embedding_openai(text: str) -> List[float]:
    """Generate embedding using OpenAI API."""
    if not OPENAI_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="OPENAI_API_KEY not set. Set it in the server environment to use OpenAI embeddings.",
        )
    payload = {
        "model": OPENAI_EMBEDDING_MODEL,
        "input": text,
    }
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json",
    }
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.post(
            "https://api.openai.com/v1/embeddings",
            headers=headers,
            json=payload,
        )
        if response.status_code >= 400:
            error_text = response.text[:400]
            raise HTTPException(
                status_code=502,
                detail=f"OpenAI embedding API error ({response.status_code}): {error_text}",
            )
        data = response.json()
        embedding_data = data.get("data", [{}])[0] if data.get("data") else {}
        embedding = embedding_data.get("embedding")
        if not embedding:
            raise HTTPException(status_code=502, detail="No embedding returned from OpenAI.")
        return normalize_embedding(embedding)


async def generate_embedding_voyage(text: str, input_type: str) -> List[float]:
    """Generate embedding using VoyageAI API."""
    if not VOYAGE_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="VOYAGE_API_KEY not set. Set it in the server environment to use VoyageAI embeddings."
        )

    # Log which model is being used for debugging

    payload = {
        "model": VOYAGE_EMBEDDING_MODEL,
        "input": [text],
        "input_type": input_type,
    }

    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.post(
            "https://api.voyageai.com/v1/embeddings",
            headers={
                "Authorization": f"Bearer {VOYAGE_API_KEY}",
                "Content-Type": "application/json",
            },
            json=payload,
        )

        if response.status_code >= 400:
            error_text = response.text[:400]
            raise HTTPException(
                status_code=502,
                detail=f"VoyageAI embedding API error ({response.status_code}): {error_text}"
            )

        data = response.json() or {}
        embedding_data = data.get("data", [{}])[0] if data.get("data") else {}
        embedding = embedding_data.get("embedding")

        if not embedding:
            raise HTTPException(status_code=502, detail="No embedding returned from VoyageAI.")

        # Log dimension for debugging
        
        return normalize_embedding(embedding)

async def generate_embedding_ollama(text: str) -> List[float]:
    """Generate embedding using Ollama."""
    if ollama is None:
        raise HTTPException(status_code=503, detail="ollama package not installed.")
    try:
        response = ollama.embeddings(model=OLLAMA_EMBEDDING_MODEL, prompt=text)
    except Exception as e:
        raise HTTPException(status_code=503, detail=f"Ollama embedding failed: {str(e)}")

    embedding = response.get("embedding") if isinstance(response, dict) else None
    if not embedding:
        raise HTTPException(status_code=502, detail="No embedding returned from Ollama.")
    return normalize_embedding(embedding)

# ---------------------------------------------------------------------------
#  Voyage Reranking — cross-encoder reranker for RAG accuracy
# ---------------------------------------------------------------------------

VOYAGE_RERANK_MODEL = os.getenv("VOYAGE_RERANK_MODEL", "rerank-2.5")

async def rerank_with_voyage(
    query: str,
    documents: list[str],
    top_k: int = 5,
    model: str | None = None,
) -> list[dict]:
    """
    Call Voyage AI rerank API to re-score documents by cross-encoder relevance.
    Returns list of {index, relevance_score, document} sorted by descending score.
    Falls back gracefully: returns original order if API fails or no key.
    """
    if not VOYAGE_API_KEY:
        return [{"index": i, "relevance_score": 1.0, "document": d} for i, d in enumerate(documents[:top_k])]
    if not documents:
        return []

    rerank_model = model or VOYAGE_RERANK_MODEL
    payload = {
        "query": query[:4000],
        "documents": [d[:8000] for d in documents],
        "model": rerank_model,
        "top_k": min(top_k, len(documents)),
    }

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.post(
                "https://api.voyageai.com/v1/rerank",
                headers={
                    "Authorization": f"Bearer {VOYAGE_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
            if response.status_code >= 400:
                return [{"index": i, "relevance_score": 1.0, "document": d} for i, d in enumerate(documents[:top_k])]

            data = response.json() or {}
            results = data.get("data", [])
            total_tokens = data.get("usage", {}).get("total_tokens", 0)
            return [
                {
                    "index": r.get("index", 0),
                    "relevance_score": r.get("relevance_score", 0),
                    "document": documents[r["index"]] if r.get("index", 0) < len(documents) else "",
                }
                for r in results
            ]
    except Exception as e:
        return [{"index": i, "relevance_score": 1.0, "document": d} for i, d in enumerate(documents[:top_k])]


# ---------------------------------------------------------------------------
#  Real-Time SSE Streaming — Long-running agentic workflows
# ---------------------------------------------------------------------------

@app.post("/ingest/document-stream")
async def ingest_document_stream(file: UploadFile = File(...), dataType: Optional[str] = None):
    """
    SSE-streaming document ingestion.
    Pushes real-time status updates (e.g., {"status": "Reading PDF..."})
    so the UI doesn't time out during the 30s+ processing pipeline.
    """
    async def generate():
        try:
            yield f"data: {json.dumps({'status': 'Uploading file...'})}\n\n"
            yield f"data: {json.dumps({'status': f'Processing {file.filename}...'})}\n\n"

            # Step 1: Extract text
            yield f"data: {json.dumps({'status': 'Extracting text content...'})}\n\n"
            file_ext, text_content = await extract_text_content(file)
            yield f"data: {json.dumps({'status': f'Extracted {len(text_content)} characters from {file_ext.upper()}'})}\n\n"

            # Step 2: Convert to structured data
            yield f"data: {json.dumps({'status': 'Converting to structured data...'})}\n\n"
            conversion_request = ConversionRequest(data=text_content, dataType=dataType, format=file_ext)
            conversion_result = await convert_data(conversion_request)
            conversion_result.raw_content = text_content[:MAX_MODEL_INPUT_CHARS]

            # Step 3: Validate
            yield f"data: {json.dumps({'status': 'Validating extracted data...'})}\n\n"
            row_errors = validate_structured_rows(
                conversion_result.startups,
                conversion_result.investors,
                conversion_result.mentors,
                conversion_result.corporates,
            )
            if row_errors:
                conversion_result.warnings = (conversion_result.warnings or []) + row_errors

            # Step 4: Return result
            result_dict = conversion_result.model_dump() if hasattr(conversion_result, 'model_dump') else conversion_result.dict()
            yield f"data: {json.dumps({'status': 'Complete!'})}\n\n"
            yield f"data: {json.dumps({'result': result_dict})}\n\n"
            yield "data: [DONE]\n\n"

        except HTTPException as he:
            yield f"data: {json.dumps({'error': he.detail})}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)[:500]})}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
            "Access-Control-Allow-Origin": "*",
        },
    )


# ---------------------------------------------------------------------------
#  Web Search — DuckDuckGo HTML search (LEGACY — kept as fallback)
#  Primary web search now uses Anthropic native web_search_20250305 tool
#  which is passed directly in /ask and /ask/stream when web_search_enabled=true
# ---------------------------------------------------------------------------

class WebSearchRequest(BaseModel):
    query: str
    max_results: int = Field(default=5, ge=1, le=10)

class WebSearchResult(BaseModel):
    title: str
    snippet: str
    url: str

class WebSearchResponse(BaseModel):
    results: List[WebSearchResult]
    query: str


@app.post("/web-search", response_model=WebSearchResponse)
async def web_search_endpoint(request: WebSearchRequest):
    """
    Search the web using DuckDuckGo HTML (no API key needed).
    Returns titles, snippets, and URLs for the top results.
    """
    query = (request.query or "").strip()
    if not query:
        raise HTTPException(status_code=400, detail="query is required.")

    results: List[WebSearchResult] = []
    try:
        import re as _re
        import html as _html

        # Use DuckDuckGo HTML search (no API key needed)
        async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as client:
            resp = await client.get(
                "https://html.duckduckgo.com/html/",
                params={"q": query},
                headers={
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                },
            )
            resp.raise_for_status()
            body = resp.text

            # Parse results from HTML
            # DuckDuckGo HTML results are in <div class="result"> blocks
            result_blocks = _re.findall(
                r'<div[^>]*class="[^"]*result__body[^"]*"[^>]*>(.*?)</div>\s*</div>',
                body,
                _re.DOTALL,
            )

            if not result_blocks:
                # Alternative parsing: look for result titles and snippets separately
                titles_raw = _re.findall(
                    r'<a[^>]*class="result__a"[^>]*>(.*?)</a>',
                    body,
                    _re.DOTALL,
                )
                snippets_raw = _re.findall(
                    r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
                    body,
                    _re.DOTALL,
                )
                urls_raw = _re.findall(
                    r'<a[^>]*class="result__a"[^>]*href="([^"]*)"',
                    body,
                )

                for i in range(min(len(titles_raw), len(snippets_raw), request.max_results)):
                    title = _re.sub(r'<[^>]+>', '', titles_raw[i]).strip()
                    snippet = _re.sub(r'<[^>]+>', '', snippets_raw[i]).strip()
                    url = urls_raw[i] if i < len(urls_raw) else ""

                    # DuckDuckGo wraps URLs in a redirect — extract the actual URL
                    url_match = _re.search(r'uddg=([^&]+)', url)
                    if url_match:
                        from urllib.parse import unquote
                        url = unquote(url_match.group(1))

                    title = _html.unescape(title)
                    snippet = _html.unescape(snippet)

                    if title and snippet:
                        results.append(WebSearchResult(
                            title=title,
                            snippet=snippet,
                            url=url,
                        ))
            else:
                for block in result_blocks[:request.max_results]:
                    title_m = _re.search(r'<a[^>]*class="result__a"[^>]*>(.*?)</a>', block, _re.DOTALL)
                    snippet_m = _re.search(r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>', block, _re.DOTALL)
                    url_m = _re.search(r'<a[^>]*class="result__a"[^>]*href="([^"]*)"', block)

                    title = _re.sub(r'<[^>]+>', '', title_m.group(1)).strip() if title_m else ""
                    snippet = _re.sub(r'<[^>]+>', '', snippet_m.group(1)).strip() if snippet_m else ""
                    url = url_m.group(1) if url_m else ""

                    url_match = _re.search(r'uddg=([^&]+)', url)
                    if url_match:
                        from urllib.parse import unquote
                        url = unquote(url_match.group(1))

                    title = _html.unescape(title)
                    snippet = _html.unescape(snippet)

                    if title and snippet:
                        results.append(WebSearchResult(
                            title=title,
                            snippet=snippet,
                            url=url,
                        ))


    except Exception as e:
        # Return empty results rather than failing the whole request
        pass

    return WebSearchResponse(results=results, query=query)


@app.post("/rewrite-query", response_model=RewriteQueryResponse)
async def rewrite_query_endpoint(request: RewriteQueryRequest):
    """
    Rewrite a query using LLM to resolve pronouns and contextualize vague questions.
    This should be called BEFORE document search to ensure proper query understanding.
    """
    question = (request.question or "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="question is required.")
    
    rewritten = await rewrite_query_with_llm(question, request.previous_messages or [])
    return RewriteQueryResponse(rewritten_question=rewritten)


# ---------------------------------------------------------------------------
#  Multi-Query Generation — produce 2-3 query variants for better recall
# ---------------------------------------------------------------------------

class MultiQueryRequest(BaseModel):
    question: str = Field(..., description="The primary user question (already rewritten if applicable)")
    max_variants: int = Field(default=3, ge=2, le=5, description="Number of query variants to generate")


class MultiQueryResponse(BaseModel):
    queries: List[str] = Field(..., description="List of query variants including the original")
    model_used: str = ""


@app.post("/multi-query", response_model=MultiQueryResponse)
async def multi_query_endpoint(request: MultiQueryRequest):
    """
    Multi-Query RAG: Generate 2-3 diverse reformulations of a user question.
    Each variant emphasizes a different angle (synonyms, specificity, broader scope)
    so parallel retrieval covers more relevant chunks.

    Always includes the original query as the first variant.
    """
    question = (request.question or "").strip()
    if not question:
        return MultiQueryResponse(queries=[question], model_used="")

    # Short or simple queries (< 4 words) — don't bother generating variants
    if len(question.split()) < 4:
        return MultiQueryResponse(queries=[question], model_used="")

    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        return MultiQueryResponse(queries=[question], model_used="")

    n_variants = min(request.max_variants, 5)
    prompt = (
        f"You are a search query expansion assistant for a venture capital platform.\n\n"
        f"Original question: {question}\n\n"
        f"Generate exactly {n_variants - 1} alternative search queries that would help find relevant information "
        f"to answer the original question. Each variant should:\n"
        f"1. Approach the topic from a different angle or use different keywords/synonyms\n"
        f"2. Be a standalone search query (not a follow-up)\n"
        f"3. Be concise (under 20 words)\n"
        f"4. Together with the original, cover the full scope of what the user might want\n\n"
        f"Examples of good variants:\n"
        f'- Original: "What is their burn rate?" → "monthly cash expenditure and runway" , "operating expenses vs revenue"\n'
        f'- Original: "Tell me about the team" → "founders background and experience" , "key hires and advisors"\n\n'
        f"Return ONLY a JSON array of strings, no markdown fences, no explanation.\n"
        f'Example: ["variant 1", "variant 2"]\n'
    )

    try:
        client = _get_anthropic_async_client()
        message = await _call_claude_with_fallback(
            client,
            preferred_model=HAIKU_MODEL,
            max_tokens=300,
            temperature=0.3,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text")).strip()
        model_used = getattr(message, "model", HAIKU_MODEL)

        # Parse JSON
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[1] if "\n" in raw else raw[3:]
            if raw.endswith("```"):
                raw = raw[:-3]
            raw = raw.strip()

        import json as _json
        variants = _json.loads(raw)

        if not isinstance(variants, list):
            raise ValueError("LLM returned non-list")

        # Clean and dedupe
        clean_variants = []
        seen = {question.lower().strip()}
        for v in variants:
            v_str = str(v).strip()
            if v_str and v_str.lower() not in seen and len(v_str) < 200:
                clean_variants.append(v_str)
                seen.add(v_str.lower())
            if len(clean_variants) >= n_variants - 1:
                break

        all_queries = [question] + clean_variants
        return MultiQueryResponse(queries=all_queries, model_used=str(model_used))

    except Exception as e:
        return MultiQueryResponse(queries=[question], model_used="")


# ---------------------------------------------------------------------------
#  Adaptive Query Router — entity extraction, intent, complexity, routing
# ---------------------------------------------------------------------------

class QueryEntity(BaseModel):
    name: str
    type: str = Field(description="company, person, fund, metric, sector, unknown")

class QueryAnalysisResponse(BaseModel):
    entities: List[QueryEntity] = Field(default_factory=list)
    intent: str = Field(default="factual", description="factual, compare, summarize, diligence, forecast, relationship, meta, conversational")
    complexity: float = Field(default=0.3, ge=0.0, le=1.0)
    retrieval_strategy: str = Field(default="vector", description="vector, vector+graph, vector+graph+structured, none")
    rewritten_query: str = ""

class AnalyzeQueryRequest(BaseModel):
    question: str
    previous_messages: List[ChatMessage] = Field(default_factory=list)


@app.post("/analyze-query", response_model=QueryAnalysisResponse)
async def analyze_query_endpoint(request: AnalyzeQueryRequest):
    """
    Adaptive retrieval router: extracts entities, classifies intent,
    scores complexity, and decides which retrieval stores to hit.
    Uses Haiku for speed (~200ms).
    """
    question = (request.question or "").strip()
    if not question:
        return QueryAnalysisResponse(rewritten_query=question)

    # First do query rewriting
    rewritten = await rewrite_query_with_llm(question, request.previous_messages or [])

    # Check if it's a meta/greeting question (fast path — no LLM needed)
    if is_meta_question(question):
        return QueryAnalysisResponse(
            entities=[],
            intent="meta" if len(question.split()) > 3 else "conversational",
            complexity=0.1,
            retrieval_strategy="none",
            rewritten_query=rewritten,
        )

    # Use Claude Haiku for entity extraction + intent classification in one call
    if _anthropic_sdk_available and ANTHROPIC_API_KEY:
        try:
            client = _get_anthropic_async_client()
            analysis_prompt = (
                f"Analyze this VC intelligence query:\n"
                f"Query: {question}\n\n"
                "Return JSON with:\n"
                '- "entities": [{{"name": "...", "type": "company|person|fund|metric|sector|unknown"}}]\n'
                '- "intent": one of "factual", "compare", "summarize", "diligence", "forecast", "relationship", "meta", "conversational"\n'
                '- "complexity": 0.0-1.0 (how much retrieval is needed)\n'
                '- "retrieval_strategy": one of "vector", "vector+graph", "vector+graph+structured", "none"\n\n'
                "Rules:\n"
                '- "compare" = comparing 2+ entities\n'
                '- "diligence" = risk assessment, red flags, deep analysis\n'
                '- "relationship" = company connections, partnerships, who knows who\n'
                '- "forecast" = growth projections, market potential, forward-looking\n'
                '- "factual" = simple lookup, what is X\n'
                '- complexity > 0.6 means multiple entities or multi-hop reasoning\n'
                '- "vector+graph" when relationships matter\n'
                '- "vector+graph+structured" when numbers/KPIs matter\n'
                "Return ONLY valid JSON."
            )

            message = await _call_claude_with_fallback(
                client,
                preferred_model=HAIKU_MODEL,
                max_tokens=300,
                temperature=0.0,
                messages=[{"role": "user", "content": analysis_prompt}],
            )
            raw = "".join(b.text for b in message.content if hasattr(b, "text"))

            # Parse JSON from response
            try:
                # Try direct parse
                data = json.loads(raw)
            except json.JSONDecodeError:
                # Try extracting JSON from markdown code block
                json_match = re.search(r'\{[\s\S]*\}', raw)
                if json_match:
                    data = json.loads(json_match.group())
                else:
                    data = {}

            entities = [
                QueryEntity(name=e.get("name", ""), type=e.get("type", "unknown"))
                for e in data.get("entities", [])
                if e.get("name")
            ]

            return QueryAnalysisResponse(
                entities=entities,
                intent=data.get("intent", "factual"),
                complexity=float(data.get("complexity", 0.3)),
                retrieval_strategy=data.get("retrieval_strategy", "vector"),
                rewritten_query=rewritten,
            )
        except Exception as e:
            pass

    # Fallback: keyword-based heuristic
    q = question.lower()
    intent = "factual"
    strategy = "vector"
    complexity = 0.3

    connection_words = ["connect", "partner", "introduce", "relationship", "linked", "network"]
    compare_words = ["compare", "vs", "versus", "difference", "better"]
    diligence_words = ["risk", "diligence", "red flag", "concern", "weakness", "threat"]
    forecast_words = ["growth", "potential", "forecast", "predict", "future", "projection"]

    if any(w in q for w in connection_words):
        intent, strategy = "relationship", "vector+graph"
    elif any(w in q for w in compare_words):
        intent, strategy, complexity = "compare", "vector+graph", 0.7
    elif any(w in q for w in diligence_words):
        intent, strategy, complexity = "diligence", "vector+graph+structured", 0.8
    elif any(w in q for w in forecast_words):
        intent, strategy, complexity = "forecast", "vector+graph+structured", 0.7

    return QueryAnalysisResponse(
        entities=[],
        intent=intent,
        complexity=complexity,
        retrieval_strategy=strategy,
        rewritten_query=rewritten,
    )


# ---------------------------------------------------------------------------
#  RAG Eval Logging — track retrieval quality for regression detection
# ---------------------------------------------------------------------------

class RAGEvalLogEntry(BaseModel):
    question: str = ""
    retrieval_strategy: str = ""
    chunks_retrieved: int = 0
    chunks_cited: int = 0
    model_used: str = ""
    latency_ms: float = 0.0
    user_feedback: Optional[str] = None  # "helpful" | "not_helpful"


# In-memory buffer (flush to DB or file periodically)
_rag_eval_buffer: List[Dict[str, Any]] = []


@app.post("/rag-eval/log")
async def log_rag_eval(entry: RAGEvalLogEntry):
    """Log a RAG evaluation entry for quality tracking and regression detection."""
    entry_dict = entry.dict()
    entry_dict["timestamp"] = time.time()
    _rag_eval_buffer.append(entry_dict)

    # Keep buffer manageable (flush oldest if > 1000)
    if len(_rag_eval_buffer) > 1000:
        _rag_eval_buffer.pop(0)

    return {"status": "logged", "buffer_size": len(_rag_eval_buffer)}


@app.get("/rag-eval/stats")
async def get_rag_eval_stats():
    """Get aggregated RAG quality stats from the eval buffer."""
    if not _rag_eval_buffer:
        return {"total_queries": 0, "message": "No eval data yet"}

    total = len(_rag_eval_buffer)
    avg_chunks_retrieved = sum(e.get("chunks_retrieved", 0) for e in _rag_eval_buffer) / total
    avg_chunks_cited = sum(e.get("chunks_cited", 0) for e in _rag_eval_buffer) / total
    avg_latency = sum(e.get("latency_ms", 0) for e in _rag_eval_buffer) / total

    helpful = sum(1 for e in _rag_eval_buffer if e.get("user_feedback") == "helpful")
    not_helpful = sum(1 for e in _rag_eval_buffer if e.get("user_feedback") == "not_helpful")

    strategies = {}
    for e in _rag_eval_buffer:
        s = e.get("retrieval_strategy", "unknown")
        strategies[s] = strategies.get(s, 0) + 1

    return {
        "total_queries": total,
        "avg_chunks_retrieved": round(avg_chunks_retrieved, 1),
        "avg_chunks_cited": round(avg_chunks_cited, 1),
        "avg_latency_ms": round(avg_latency, 0),
        "feedback": {"helpful": helpful, "not_helpful": not_helpful, "no_feedback": total - helpful - not_helpful},
        "retrieval_strategies": strategies,
        "citation_rate": round(avg_chunks_cited / max(avg_chunks_retrieved, 1), 2),
    }


@app.post("/embed/query", response_model=EmbedResponse)
async def embed_query(request: EmbedRequest):
    text = (request.text or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="text is required.")

    try:
        input_type = (request.input_type or "document").strip().lower()
        if input_type not in ["document", "query"]:
            input_type = "document"

        # Try the configured provider first, then fall through to alternatives
        providers_to_try = []
        if EMBEDDINGS_PROVIDER == "voyage":
            providers_to_try = ["voyage", "openai", "ollama"]
        elif EMBEDDINGS_PROVIDER == "openai":
            providers_to_try = ["openai", "voyage", "ollama"]
        elif EMBEDDINGS_PROVIDER == "ollama":
            providers_to_try = ["ollama", "voyage", "openai"]
        else:
            providers_to_try = ["voyage", "openai", "ollama"]

        embedding = None
        last_error = None
        for provider in providers_to_try:
            try:
                if provider == "voyage" and VOYAGE_API_KEY:
                    embedding = await generate_embedding_voyage(text, input_type)
                    break
                elif provider == "openai" and OPENAI_API_KEY:
                    embedding = await generate_embedding_openai(text)
                    break
                elif provider == "ollama":
                    embedding = await generate_embedding_ollama(text)
                    break
            except Exception as provider_err:
                last_error = provider_err
                continue

        if embedding is None:
            detail = (
                f"No embedding provider available. "
                f"Set one of: VOYAGE_API_KEY, OPENAI_API_KEY, or run Ollama locally. "
                f"Configured provider: '{EMBEDDINGS_PROVIDER}'. "
                f"Last error: {last_error}"
            )
            raise HTTPException(status_code=503, detail=detail)
        
        return EmbedResponse(embedding=embedding)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=503,
            detail=f"Embedding generation failed: {str(e)}"
        )


# ---------------------------------------------------------------------------
#  Entity Extraction — auto-populate knowledge graph from documents
# ---------------------------------------------------------------------------

class ExtractedEntity(BaseModel):
    name: str
    type: str = Field(description="company, person, fund, round, sector, metric, location")
    properties: Dict[str, Any] = Field(default_factory=dict)
    confidence: float = 0.8

class ExtractedRelationship(BaseModel):
    source_name: str
    target_name: str
    relation_type: str
    properties: Dict[str, Any] = Field(default_factory=dict)
    confidence: float = 0.8

class ExtractedKPI(BaseModel):
    company_name: str
    metric_name: str
    value: float
    unit: str = "USD"
    period: str = ""
    category: str = "financial"
    confidence: float = 0.8

class EntityExtractionRequest(BaseModel):
    document_title: str = ""
    document_text: str = ""
    document_type: str = ""  # pitch_deck, memo, email, report
    pdf_base64: Optional[str] = Field(default=None, description="Base64-encoded PDF bytes for visual extraction")

class EntityExtractionResponse(BaseModel):
    entities: List[ExtractedEntity] = []
    relationships: List[ExtractedRelationship] = []
    kpis: List[ExtractedKPI] = []


@app.post("/extract-entities", response_model=EntityExtractionResponse)
async def extract_entities(request: EntityExtractionRequest):
    """
    Extract entities, relationships, and KPIs from a document using Claude.
    
    When `pdf_base64` is provided, Claude reads the PDF visually for much better
    extraction of companies, people, metrics from pitch decks.
    """
    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        return EntityExtractionResponse()

    has_pdf = bool(request.pdf_base64 and len(request.pdf_base64) > 100)
    text = (request.document_text or "").strip()[:8000]
    
    if not text and not has_pdf:
        return EntityExtractionResponse()

    # Build extraction instructions — add email-specific guidance when document_type is 'email'
    email_hints = ""
    if request.document_type == "email":
        email_hints = (
            "\nEMAIL-SPECIFIC INSTRUCTIONS:\n"
            "- Extract the sender and all recipients as 'person' entities.\n"
            "- Infer company names from email domains (e.g. jane@acme.com → company 'Acme').\n"
            "- Look for deal-related language: 'term sheet', 'Series A', 'valuation', 'cap table', 'closing'.\n"
            "- Extract any mentioned meeting dates, deadlines, or milestones.\n"
            "- Create 'works_at' relationships between people and their domain-inferred companies.\n"
            "- Create 'invested_in' or 'partner_of' relationships if investment or partnership context is present.\n\n"
        )

    extraction_instructions = (
        "Extract ALL of the following from this VC/investment document:\n\n"
        "1. ENTITIES — people, companies, funds, funding rounds, sectors, locations:\n"
        '   Format: [{{"name": "...", "type": "company|person|fund|round|sector|location", '
        '"properties": {{"industry": "...", "role": "...", "email": "...", etc.}}, "confidence": 0.0-1.0}}]\n\n'
        "2. RELATIONSHIPS between entities:\n"
        '   Format: [{{"source_name": "...", "target_name": "...", '
        '"relation_type": "founded|works_at|invested_in|raised|led_round|partner_of|'
        'competitor_of|acquired|operates_in|located_in|board_member|advisor|portfolio_company", '
        '"properties": {{}}, "confidence": 0.0-1.0}}]\n\n'
        "3. KPIs — any numbers, metrics, financial data:\n"
        '   Format: [{{"company_name": "...", "metric_name": "revenue|arr|mrr|valuation|'
        'burn_rate|headcount|users|growth_rate|raise_amount|etc.", '
        '"value": 123.0, "unit": "USD|%|count", "period": "2024-Q3", '
        '"category": "financial|growth|fundraising|operational|market|tokenomics", '
        '"confidence": 0.0-1.0}}]\n\n'
        f"{email_hints}"
        'Return JSON with keys: "entities", "relationships", "kpis". Return ONLY valid JSON.'
    )

    try:
        client = _get_anthropic_async_client()
        
        # Build content blocks — PDF visual or text-only
        if has_pdf:
            content_blocks = [
                {
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": request.pdf_base64,
                    },
                },
                {
                    "type": "text",
                    "text": (
                        f"Document title: {request.document_title}\n"
                        f"Document type: {request.document_type or 'pitch_deck'}\n\n"
                        f"{extraction_instructions}"
                    ),
                },
            ]
        else:
            content_blocks = [
                {
                    "type": "text",
                    "text": (
                        f"Document title: {request.document_title}\n"
                        f"Document type: {request.document_type or 'unknown'}\n\n"
                        f"Text:\n{text}\n\n"
                        f"{extraction_instructions}"
                    ),
                },
            ]
        
        # Use Sonnet for higher quality extraction — with 404 fallback
        message = await _call_claude_with_fallback(
            client,
            max_tokens=4000,
            temperature=0.0,
            messages=[{"role": "user", "content": content_blocks}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text"))

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            json_match = re.search(r'\{[\s\S]*\}', raw)
            if json_match:
                data = json.loads(json_match.group())
            else:
                return EntityExtractionResponse()

        entities = [
            ExtractedEntity(**e)
            for e in data.get("entities", [])
            if e.get("name")
        ]
        relationships = [
            ExtractedRelationship(**r)
            for r in data.get("relationships", [])
            if r.get("source_name") and r.get("target_name")
        ]
        kpis = [
            ExtractedKPI(**k)
            for k in data.get("kpis", [])
            if k.get("company_name") and k.get("metric_name")
        ]

        return EntityExtractionResponse(
            entities=entities,
            relationships=relationships,
            kpis=kpis,
        )

    except Exception as e:
        return EntityExtractionResponse()


# ---------------------------------------------------------------------------
#  Company Property Extraction — auto-populate company cards from documents
# ---------------------------------------------------------------------------

class CompanyPropertyExtractionRequest(BaseModel):
    """Extract structured company properties from a document."""
    raw_content: str = Field(default="", description="Document text content")
    document_title: str = ""
    document_type: str = ""  # pitch_deck, investment_memo, data_room, email
    existing_properties: Dict[str, Any] = Field(default_factory=dict, description="Current card properties (for context)")
    pdf_base64: Optional[str] = Field(default=None, description="Base64-encoded PDF bytes — Claude reads the PDF directly for much better extraction")


class CompanyPropertyExtractionResponse(BaseModel):
    properties: Dict[str, Any] = Field(default_factory=dict)
    confidence: Dict[str, float] = Field(default_factory=dict)
    document_type_detected: str = ""


@app.post("/extract-company-properties", response_model=CompanyPropertyExtractionResponse)
async def extract_company_properties(request: CompanyPropertyExtractionRequest):
    """
    Extract structured company card properties from document text or PDF.
    
    When `pdf_base64` is provided, Claude reads the PDF visually (native document block)
    for dramatically better extraction of pitch decks with tables, charts, and layouts.
    Falls back to text-based extraction when only `raw_content` is provided.
    """
    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        return CompanyPropertyExtractionResponse()

    has_pdf = bool(request.pdf_base64 and len(request.pdf_base64) > 100)
    text = (request.raw_content or "").strip()
    
    if not text and not has_pdf:
        return CompanyPropertyExtractionResponse()

    # Truncate text to fit context window (only used as fallback or supplement)
    text = text[:MAX_MODEL_INPUT_CHARS]

    # Step 1: Detect document type — skip LLM call if we have a PDF (it's almost certainly a pitch deck)
    doc_type = (request.document_type or "").strip()
    if not doc_type:
        if has_pdf:
            doc_type = "pitch_deck"  # PDFs uploaded to VC platform are almost always pitch decks
        elif text:
            doc_type_prompt = (
                f"Document title: {request.document_title}\n"
                f"First 500 chars:\n{text[:500]}\n\n"
                "Classify this document as ONE of: pitch_deck, investment_memo, data_room, email, report, other\n"
                "Return ONLY the classification word, nothing else."
            )
            try:
                client = _get_anthropic_async_client()
                msg = await _call_claude_with_fallback(
                    client,
                    preferred_model=HAIKU_MODEL,
                    max_tokens=20,
                    temperature=0.0,
                    messages=[{"role": "user", "content": doc_type_prompt}],
                )
                doc_type = "".join(b.text for b in msg.content if hasattr(b, "text")).strip().lower()
                doc_type = doc_type.replace('"', '').replace("'", "").strip()
            except Exception:
                doc_type = "unknown"

    # Step 2: Build extraction prompt based on document type
    field_guidance = ""
    if doc_type == "pitch_deck":
        field_guidance = (
            "This is a pitch deck. Extract EVERYTHING available: bio, problem, solution, "
            "TAM/SAM/SOM, market_growth_rate, competitive_edge, competitors, business_model, "
            "revenue_model, pricing, gtm_strategy, traction, unit economics (cac, ltv, "
            "ltv_cac_ratio, payback_period, gross_margin, net_margin, churn_rate), "
            "founders (with LinkedIn if shown), funding_stage, amount_seeking, use_of_funds, "
            "mrr, arr, customers_count, growth_rate, team_size, founded_year, headquarters, "
            "industry, geo_markets, website, email, phone, linkedin_url, twitter_url, "
            "key_partnerships, awards_recognition. Extract EVERY data point you can find."
        )
    elif doc_type == "investment_memo":
        field_guidance = (
            "This is an investment memo. Extract EVERYTHING: valuation, amount_seeking, "
            "funding_stage, arr, mrr, burn_rate, runway_months, competitive_edge, competitors, "
            "business_model, revenue_model, traction, unit economics (cac, ltv, ltv_cac_ratio, "
            "gross_margin, net_margin, churn_rate), founders, tam, sam, som, market_growth_rate, "
            "gtm_strategy, use_of_funds, customers_count, growth_rate, industry, geo_markets, "
            "headquarters, team_size, key_partnerships."
        )
    elif doc_type == "data_room":
        field_guidance = (
            "This is a data room document. Extract: arr, mrr, burn_rate, runway_months, "
            "valuation, founders, industry, unit economics (cac, ltv, gross_margin, "
            "net_margin, churn_rate), customers_count, growth_rate, traction, team_size."
        )
    else:
        field_guidance = "Extract ALL available company properties from this document — every data point matters."

    existing_context = ""
    if request.existing_properties:
        non_empty = {k: v for k, v in request.existing_properties.items()
                     if v and not k.startswith("_") and k not in ("auto_created", "source", "first_seen_document", "last_seen_document", "document_count")}
        if non_empty:
            existing_context = (
                f"\n\nExisting card data (for context — extract NEW information not already here):\n"
                f"{json.dumps(non_empty, indent=2, default=str)}\n"
            )

    json_schema_instructions = (
        "Extract ALL available company properties from this document into JSON.\n"
        "For each field, also provide a confidence score (0.0-1.0) indicating how sure you are.\n"
        "Leave fields as empty string or empty array if the information is NOT in the document.\n"
        "Do NOT make up information — only extract what's explicitly stated or strongly implied.\n"
        "Extract as MUCH as possible — every data point matters for investment analysis.\n\n"
        "IMPORTANT: The 'company_name' field is CRITICAL. It must be the SHORT, official name of "
        "the PRIMARY company this document is about (e.g. 'TBE', 'UniToni', 'Weego'). "
        "Do NOT use the document title, folder name, or investor name. "
        "If the document covers multiple companies, pick the MAIN subject company.\n\n"
        "Required JSON format:\n"
        "{\n"
        '  "properties": {\n'
        '    "company_name": "Short official company name (e.g. TBE, UniToni, Payd)",\n'
        '    "bio": "1-3 sentence company description",\n'
        '    "funding_stage": "Pre-seed|Seed|Series A|Series B+|etc",\n'
        '    "amount_seeking": "$X.XM or $XXK",\n'
        '    "valuation": "$X.XM pre/post-money",\n'
        '    "arr": "$X.XM or $XXK annual recurring revenue",\n'
        '    "mrr": "$XXK monthly recurring revenue",\n'
        '    "burn_rate": "$XXK/month",\n'
        '    "runway_months": "XX months",\n'
        '    "problem": "the problem being solved (detailed)",\n'
        '    "solution": "how the company solves it (detailed)",\n'
        '    "tam": "total addressable market size with $ figure",\n'
        '    "sam": "serviceable addressable market size",\n'
        '    "som": "serviceable obtainable market",\n'
        '    "market_growth_rate": "market CAGR or growth rate %",\n'
        '    "competitive_edge": "moat / differentiation / why they win",\n'
        '    "business_model": "how the company makes money (subscription, marketplace, SaaS, etc.)",\n'
        '    "revenue_model": "revenue streams and pricing model details",\n'
        '    "pricing": "pricing tiers or strategy details",\n'
        '    "gtm_strategy": "go-to-market strategy, channels, customer acquisition approach",\n'
        '    "traction": "key traction summary — users, revenue milestones, growth highlights",\n'
        '    "customers_count": "number of customers or users",\n'
        '    "growth_rate": "user or revenue growth rate (MoM/YoY)",\n'
        '    "gmv": "gross merchandise value if applicable",\n'
        '    "cac": "customer acquisition cost",\n'
        '    "ltv": "customer lifetime value",\n'
        '    "ltv_cac_ratio": "LTV:CAC ratio",\n'
        '    "payback_period": "months to recover CAC",\n'
        '    "gross_margin": "gross margin percentage",\n'
        '    "net_margin": "net margin percentage",\n'
        '    "churn_rate": "monthly or annual churn rate",\n'
        '    "use_of_funds": "how the raised capital will be allocated",\n'
        '    "founded_year": "year the company was founded",\n'
        '    "headquarters": "HQ city, country",\n'
        '    "team_size": "number of employees",\n'
        '    "founders": [{"name": "...", "role": "CEO/CTO/etc", "background": "notable experience, education", "linkedin": "LinkedIn URL if found"}],\n'
        '    "competitors": [{"name": "competitor name", "differentiator": "how this company differs from them"}],\n'
        '    "website": "https://...",\n'
        '    "email": "contact email address",\n'
        '    "phone": "contact phone number",\n'
        '    "linkedin_url": "company LinkedIn page URL",\n'
        '    "twitter_url": "company Twitter/X URL or handle",\n'
        '    "industry": "Fintech|SaaS|AI/ML|Healthcare|Logistics|etc",\n'
        '    "geo_markets": ["market region 1", "market region 2", ...],\n'
        '    "key_partnerships": ["partner 1", "partner 2", ...],\n'
        '    "awards_recognition": "notable awards, accelerator programs, press mentions"\n'
        "  },\n"
        '  "confidence": {\n'
        '    "bio": 0.9,\n'
        '    "funding_stage": 0.8,\n'
        "    ...(one entry per non-empty field)\n"
        "  }\n"
        "}\n\n"
        "Return ONLY valid JSON. No markdown, no explanation."
    )

    try:
        client = _get_anthropic_async_client()
        
        # ── Build message content: PDF visual reading vs text-only ──
        if has_pdf:
            # BEST PATH: Send PDF as a native document block — Claude reads it visually
            # This captures tables, charts, multi-column layouts, logos, etc.
            content_blocks = [
                {
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": request.pdf_base64,
                    },
                },
                {
                    "type": "text",
                    "text": (
                        f"Document title: {request.document_title}\n"
                        f"Document type: {doc_type}\n"
                        f"{field_guidance}\n{existing_context}\n\n"
                        f"{json_schema_instructions}"
                    ),
                },
            ]
        else:
            # FALLBACK: Text-only extraction
            content_blocks = [
                {
                    "type": "text",
                    "text": (
                        f"Document title: {request.document_title}\n"
                        f"Document type: {doc_type}\n"
                        f"{field_guidance}\n{existing_context}\n"
                        f"Document text:\n---\n{text}\n---\n\n"
                        f"{json_schema_instructions}"
                    ),
                },
            ]
        
        message = await _call_claude_with_fallback(
            client,
            max_tokens=8192,
            temperature=0.0,
            messages=[{"role": "user", "content": content_blocks}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text"))

        # Parse JSON from response
        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            # Try to extract JSON block
            json_match = re.search(r'\{[\s\S]*\}', raw)
            if json_match:
                data = json.loads(json_match.group())
            else:
                return CompanyPropertyExtractionResponse(document_type_detected=doc_type)

        properties = data.get("properties", data)
        confidence = data.get("confidence", {})

        # Normalize: remove None values, convert to strings where expected
        clean_props: Dict[str, Any] = {}
        string_fields = [
            "company_name",
            "bio", "funding_stage", "amount_seeking", "valuation", "arr", "mrr",
            "burn_rate", "runway_months", "problem", "solution", "tam", "sam", "som",
            "market_growth_rate", "competitive_edge", "business_model", "revenue_model",
            "pricing", "gtm_strategy", "traction", "customers_count", "growth_rate",
            "gmv", "cac", "ltv", "ltv_cac_ratio", "payback_period", "gross_margin",
            "net_margin", "churn_rate", "use_of_funds", "founded_year", "headquarters",
            "team_size", "website", "email", "phone", "linkedin_url", "twitter_url",
            "industry", "awards_recognition",
        ]
        list_fields = ["geo_markets", "key_partnerships"]
        json_fields = ["founders", "competitors"]

        for field in string_fields:
            val = properties.get(field, "")
            if val is None:
                val = ""
            if isinstance(val, (int, float)):
                val = str(val)
            if isinstance(val, str) and val.strip():
                clean_props[field] = val.strip()

        for field in list_fields:
            val = properties.get(field, [])
            if isinstance(val, list) and val:
                clean_props[field] = [str(v).strip() for v in val if v]

        for field in json_fields:
            val = properties.get(field, [])
            if isinstance(val, list) and val:
                clean_props[field] = val
            elif isinstance(val, str) and val.strip():
                try:
                    parsed = json.loads(val)
                    if isinstance(parsed, list):
                        clean_props[field] = parsed
                except json.JSONDecodeError:
                    pass

        # Clean confidence: only keep entries for fields we actually extracted
        clean_confidence = {}
        for field, score in confidence.items():
            if field in clean_props:
                try:
                    clean_confidence[field] = float(score)
                except (ValueError, TypeError):
                    clean_confidence[field] = 0.5

        return CompanyPropertyExtractionResponse(
            properties=clean_props,
            confidence=clean_confidence,
            document_type_detected=doc_type,
        )

    except Exception as e:
        return CompanyPropertyExtractionResponse(document_type_detected=doc_type)


# ---------------------------------------------------------------------------
#  SSE Streaming Extraction — keeps Render.com connection alive during long
#  Claude API calls (Render free tier has ~30s request timeout; SSE avoids it
#  because data is continuously streamed).
# ---------------------------------------------------------------------------

@app.post("/extract-company-properties/stream")
async def extract_company_properties_stream(request: CompanyPropertyExtractionRequest):
    """
    Streaming version of /extract-company-properties.
    Sends keepalive pings every 5s while Claude processes the PDF,
    then sends the final JSON result.
    Avoids Render.com's 30-second request timeout.
    """
    import asyncio

    async def _do_extraction() -> dict:
        """Run the actual extraction and return a dict."""
        if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
            return {"properties": {}, "confidence": {}, "document_type_detected": ""}

        has_pdf = bool(request.pdf_base64 and len(request.pdf_base64) > 100)
        text = (request.raw_content or "").strip()[:MAX_MODEL_INPUT_CHARS]

        if not text and not has_pdf:
            return {"properties": {}, "confidence": {}, "document_type_detected": ""}

        doc_type = (request.document_type or "").strip() or ("pitch_deck" if has_pdf else "unknown")

        field_guidance = ""
        if doc_type == "pitch_deck":
            field_guidance = (
                "This is a pitch deck. Extract EVERYTHING available: bio, problem, solution, "
                "TAM/SAM/SOM, market_growth_rate, competitive_edge, competitors, business_model, "
                "revenue_model, pricing, gtm_strategy, traction, unit economics (cac, ltv, "
                "ltv_cac_ratio, payback_period, gross_margin, net_margin, churn_rate), "
                "founders (with LinkedIn if shown), funding_stage, amount_seeking, use_of_funds, "
                "mrr, arr, customers_count, growth_rate, team_size, founded_year, headquarters, "
                "industry, geo_markets, website, email, phone, linkedin_url, twitter_url, "
                "key_partnerships, awards_recognition. Extract EVERY data point you can find."
            )
        elif doc_type == "investment_memo":
            field_guidance = (
                "This is an investment memo. Extract EVERYTHING: valuation, amount_seeking, "
                "funding_stage, arr, mrr, burn_rate, runway_months, competitive_edge, competitors, "
                "business_model, revenue_model, traction, unit economics (cac, ltv, ltv_cac_ratio, "
                "gross_margin, net_margin, churn_rate), founders, tam, sam, som, market_growth_rate, "
                "gtm_strategy, use_of_funds, customers_count, growth_rate, industry, geo_markets, "
                "headquarters, team_size, key_partnerships."
            )
        else:
            field_guidance = "Extract ALL available company properties from this document — every data point matters."

        existing_context = ""
        if request.existing_properties:
            non_empty = {k: v for k, v in request.existing_properties.items()
                         if v and not k.startswith("_") and k not in ("auto_created", "source", "first_seen_document", "last_seen_document", "document_count")}
            if non_empty:
                existing_context = f"\n\nExisting card data:\n{json.dumps(non_empty, indent=2, default=str)}\n"

        json_instructions = (
            "Extract ALL company properties into JSON. Extract every data point you can find:\n"
            '{"properties": {"bio": "...", "funding_stage": "...", "amount_seeking": "...", '
            '"valuation": "...", "arr": "...", "mrr": "...", "burn_rate": "...", "runway_months": "...", '
            '"problem": "...", "solution": "...", "tam": "...", "sam": "...", "som": "...", '
            '"market_growth_rate": "...", "competitive_edge": "...", '
            '"business_model": "...", "revenue_model": "...", "pricing": "...", '
            '"gtm_strategy": "...", "traction": "...", '
            '"customers_count": "...", "growth_rate": "...", "gmv": "...", '
            '"cac": "...", "ltv": "...", "ltv_cac_ratio": "...", "payback_period": "...", '
            '"gross_margin": "...", "net_margin": "...", "churn_rate": "...", '
            '"use_of_funds": "...", "founded_year": "...", "headquarters": "...", "team_size": "...", '
            '"founders": [{"name": "...", "role": "...", "background": "...", "linkedin": "..."}], '
            '"competitors": [{"name": "...", "differentiator": "..."}], '
            '"website": "...", "email": "...", "phone": "...", '
            '"linkedin_url": "...", "twitter_url": "...", '
            '"industry": "...", "geo_markets": [...], '
            '"key_partnerships": [...], "awards_recognition": "..."}, '
            '"confidence": {"bio": 0.9, ...}}\n'
            "Leave fields empty if not in document. Return ONLY valid JSON."
        )

        client = _get_anthropic_async_client()

        if has_pdf:
            content_blocks = [
                {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": request.pdf_base64}},
                {"type": "text", "text": f"Document title: {request.document_title}\nType: {doc_type}\n{field_guidance}\n{existing_context}\n{json_instructions}"},
            ]
        else:
            content_blocks = [
                {"type": "text", "text": f"Document title: {request.document_title}\nType: {doc_type}\n{field_guidance}\n{existing_context}\nText:\n---\n{text}\n---\n\n{json_instructions}"},
            ]

        message = await _call_claude_with_fallback(
            client,
            max_tokens=8192,
            temperature=0.0,
            messages=[{"role": "user", "content": content_blocks}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text"))

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            json_match = re.search(r'\{[\s\S]*\}', raw)
            data = json.loads(json_match.group()) if json_match else {}

        properties = data.get("properties", data)
        confidence = data.get("confidence", {})

        # Normalize
        clean_props: Dict[str, Any] = {}
        for f in [
            "bio", "funding_stage", "amount_seeking", "valuation", "arr", "mrr",
            "burn_rate", "runway_months", "problem", "solution", "tam", "sam", "som",
            "market_growth_rate", "competitive_edge", "business_model", "revenue_model",
            "pricing", "gtm_strategy", "traction", "customers_count", "growth_rate",
            "gmv", "cac", "ltv", "ltv_cac_ratio", "payback_period", "gross_margin",
            "net_margin", "churn_rate", "use_of_funds", "founded_year", "headquarters",
            "team_size", "website", "email", "phone", "linkedin_url", "twitter_url",
            "industry", "awards_recognition",
        ]:
            val = properties.get(f, "")
            if val and isinstance(val, (str, int, float)):
                clean_props[f] = str(val).strip()
        for f in ["geo_markets", "key_partnerships"]:
            val = properties.get(f, [])
            if isinstance(val, list) and val:
                clean_props[f] = [str(v).strip() for v in val if v]
        for f in ["founders", "competitors"]:
            val = properties.get(f, [])
            if isinstance(val, list) and val:
                clean_props[f] = val
            elif isinstance(val, str) and val.strip():
                try:
                    parsed = json.loads(val)
                    if isinstance(parsed, list):
                        clean_props[f] = parsed
                except json.JSONDecodeError:
                    pass

        clean_confidence = {}
        for f, s in confidence.items():
            if f in clean_props:
                try:
                    clean_confidence[f] = float(s)
                except (ValueError, TypeError):
                    clean_confidence[f] = 0.5

        return {"properties": clean_props, "confidence": clean_confidence, "document_type_detected": doc_type}

    async def generate():
        """SSE generator: keepalive pings + final result."""
        task = asyncio.create_task(_do_extraction())
        try:
            while not task.done():
                yield f"data: {json.dumps({'status': 'processing'})}\n\n"
                await asyncio.sleep(5)

            result = await task
            yield f"data: {json.dumps({'status': 'done', 'result': result})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'status': 'error', 'error': str(e)})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Access-Control-Allow-Origin": "*"},
    )


@app.post("/extract-entities/stream")
async def extract_entities_stream(request: EntityExtractionRequest):
    """
    Streaming version of /extract-entities.
    Sends keepalive pings every 5s while Claude processes the document.
    """
    import asyncio

    async def _do_extraction() -> dict:
        if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
            return {"entities": [], "relationships": [], "kpis": []}

        has_pdf = bool(request.pdf_base64 and len(request.pdf_base64) > 100)
        text = (request.document_text or "").strip()[:8000]

        if not text and not has_pdf:
            return {"entities": [], "relationships": [], "kpis": []}

        extraction_instructions = (
            "Extract ALL of the following from this VC/investment document:\n\n"
            "1. ENTITIES — people, companies, funds, funding rounds, sectors, locations:\n"
            '   Format: [{"name": "...", "type": "company|person|fund|round|sector|location", '
            '"properties": {"industry": "...", "role": "...", etc.}, "confidence": 0.0-1.0}]\n\n'
            "2. RELATIONSHIPS between entities:\n"
            '   Format: [{"source_name": "...", "target_name": "...", '
            '"relation_type": "founded|works_at|invested_in|raised|led_round|partner_of|'
            'competitor_of|acquired|operates_in|located_in|board_member|advisor|portfolio_company", '
            '"properties": {}, "confidence": 0.0-1.0}]\n\n'
            "3. KPIs — any numbers, metrics, financial data:\n"
            '   Format: [{"company_name": "...", "metric_name": "revenue|arr|mrr|valuation|'
            'burn_rate|headcount|users|growth_rate|raise_amount|etc.", '
            '"value": 123.0, "unit": "USD|%|count", "period": "2024-Q3", '
            '"category": "financial|growth|fundraising|operational|market|tokenomics", '
            '"confidence": 0.0-1.0}]\n\n'
            'Return JSON with keys: "entities", "relationships", "kpis". Return ONLY valid JSON.'
        )

        client = _get_anthropic_async_client()

        if has_pdf:
            content_blocks = [
                {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": request.pdf_base64}},
                {"type": "text", "text": f"Document title: {request.document_title}\nType: {request.document_type or 'pitch_deck'}\n\n{extraction_instructions}"},
            ]
        else:
            content_blocks = [
                {"type": "text", "text": f"Document title: {request.document_title}\nType: {request.document_type or 'unknown'}\n\nText:\n{text}\n\n{extraction_instructions}"},
            ]

        message = await _call_claude_with_fallback(
            client,
            max_tokens=4000,
            temperature=0.0,
            messages=[{"role": "user", "content": content_blocks}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text"))

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            json_match = re.search(r'\{[\s\S]*\}', raw)
            data = json.loads(json_match.group()) if json_match else {}

        result = {
            "entities": [e for e in data.get("entities", []) if e.get("name")],
            "relationships": [r for r in data.get("relationships", []) if r.get("source_name") and r.get("target_name")],
            "kpis": [k for k in data.get("kpis", []) if k.get("company_name") and k.get("metric_name")],
        }
        return result

    async def generate():
        task = asyncio.create_task(_do_extraction())
        try:
            while not task.done():
                yield f"data: {json.dumps({'status': 'processing'})}\n\n"
                await asyncio.sleep(5)

            result = await task
            yield f"data: {json.dumps({'status': 'done', 'result': result})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'status': 'error', 'error': str(e)})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Access-Control-Allow-Origin": "*"},
    )


# ---------------------------------------------------------------------------
#  GraphRAG / Contextual Retrieval — V2 retrieval pipeline
# ---------------------------------------------------------------------------

class ContextualChunkRequest(BaseModel):
    """Pre-append a ~100-word contextual header to a chunk before embedding."""
    document_title: str = ""
    document_summary: str = ""
    chunk_text: str
    chunk_index: int = 0
    total_chunks: int = 1


class ContextualChunkResponse(BaseModel):
    enriched_chunk: str
    contextual_header: str


def _fallback_contextual_header(chunk_text: str, document_title: str, chunk_index: int, total_chunks: int) -> str:
    """When Claude is unavailable or returns empty, build a minimal header from chunk and title."""
    if not chunk_text or not chunk_text.strip():
        return ""
    first_line = chunk_text.strip().split("\n")[0].strip()[:120]
    if not first_line:
        first_line = chunk_text.strip()[:120]
    title_part = f"{document_title}: " if document_title and document_title.strip() else ""
    return f"{title_part}Chunk {chunk_index + 1}/{total_chunks}. {first_line}"


@app.post("/contextualize-chunk", response_model=ContextualChunkResponse)
async def contextualize_chunk(request: ContextualChunkRequest):
    """
    Contextual Retrieval: Generate a ~100-word contextual header for a chunk
    that explains what the chunk is about within the larger document.
    This dramatically improves embedding hit rates (per Anthropic's Contextual Retrieval paper).

    Call this endpoint for each chunk *before* embedding it.
    Always returns a non-empty contextual_header when chunk_text is non-empty (Claude or fallback).
    """
    chunk_text = (request.chunk_text or "").strip()
    fallback = _fallback_contextual_header(
        chunk_text, request.document_title or "", request.chunk_index, request.total_chunks or 1
    )

    if not ANTHROPIC_API_KEY:
        enriched = f"{fallback}\n\n{chunk_text}" if fallback else chunk_text
        return ContextualChunkResponse(enriched_chunk=enriched, contextual_header=fallback)

    prompt = (
        f"Document title: {request.document_title}\n"
        f"Document summary: {request.document_summary[:500]}\n"
        f"Chunk {request.chunk_index + 1} of {request.total_chunks}:\n"
        f"---\n{request.chunk_text[:2000]}\n---\n\n"
        "Write a SHORT (50-100 word) contextual header that:\n"
        "1. Identifies what specific section/topic this chunk covers\n"
        "2. Places it in the broader document context\n"
        "3. Mentions key entities (companies, people, metrics)\n"
        "Return ONLY the contextual header text, nothing else."
    )

    try:
        if _anthropic_sdk_available:
            client = _get_anthropic_async_client()
            message = await _call_claude_with_fallback(
                client,
                preferred_model=HAIKU_MODEL,
                max_tokens=200,
                temperature=0.0,
                messages=[{"role": "user", "content": prompt}],
            )
            header = "".join(b.text for b in message.content if hasattr(b, "text")).strip()
        else:
            header = fallback
    except Exception:
        header = fallback

    if not header:
        header = fallback
    enriched = f"{header}\n\n{request.chunk_text}" if header else request.chunk_text
    return ContextualChunkResponse(enriched_chunk=enriched, contextual_header=header)


# ---------------------------------------------------------------------------
#  Agentic Chunking — LLM decides where to split the document by topic
# ---------------------------------------------------------------------------

class AgenticChunkRequest(BaseModel):
    """Send document text to an LLM which returns topic-based sections."""
    document_title: str = ""
    document_text: str = Field(..., description="Full document text (will be truncated to 12k chars)")
    max_sections: int = Field(default=8, description="Maximum number of sections to produce")


class AgenticSection(BaseModel):
    label: str = Field(..., description="Short topic label for this section")
    text: str = Field(..., description="The verbatim text content for this section")


class AgenticChunkResponse(BaseModel):
    sections: List[AgenticSection]
    model_used: str = ""
    fallback: bool = Field(default=False, description="True if LLM call failed and we fell back to paragraphs")


@app.post("/agentic-chunk", response_model=AgenticChunkResponse)
async def agentic_chunk(request: AgenticChunkRequest):
    """
    Agentic Chunking: An LLM reads the document and splits it into
    coherent topic-based sections with labels.

    Each returned section becomes a parent chunk; the frontend splits
    children within each section on sentence boundaries.
    """
    doc_text = request.document_text[:12000]
    title = request.document_title or "Untitled"
    max_sections = min(request.max_sections, 12)

    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        # No Claude — fall back to naive paragraph splitting
        sections = _paragraph_fallback(doc_text, max_sections)
        return AgenticChunkResponse(sections=sections, model_used="", fallback=True)

    prompt = (
        f"You are a document segmentation assistant. Your job is to split the following document into coherent topic-based sections.\n\n"
        f"Document title: {title}\n"
        f"---\n{doc_text}\n---\n\n"
        f"Split this document into at most {max_sections} sections. Each section should cover one coherent topic or theme.\n\n"
        f"Rules:\n"
        f"1. Every word of the original text must appear in exactly one section — do NOT summarize, paraphrase, or drop any text.\n"
        f"2. Sections must be in the same order as the original text.\n"
        f"3. Each section gets a short label (3-8 words) describing the topic.\n"
        f"4. Prefer splitting at natural boundaries: paragraph breaks, heading changes, topic shifts.\n"
        f"5. A section should be at least 200 characters and at most 3000 characters when possible.\n"
        f"6. If the document is very short (< 500 chars), return it as a single section.\n\n"
        f"Respond with ONLY a JSON array, no markdown fences, no explanation. Format:\n"
        f'[{{"label": "Section topic", "text": "verbatim text..."}}, ...]\n'
    )

    try:
        client = _get_anthropic_async_client()
        message = await _call_claude_with_fallback(
            client,
            preferred_model=HAIKU_MODEL,
            max_tokens=8000,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = "".join(b.text for b in message.content if hasattr(b, "text")).strip()
        model_used = getattr(message, "model", HAIKU_MODEL)

        # Parse JSON — handle potential markdown fences
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[1] if "\n" in raw else raw[3:]
            if raw.endswith("```"):
                raw = raw[:-3]
            raw = raw.strip()

        import json as _json
        parsed = _json.loads(raw)

        if not isinstance(parsed, list) or len(parsed) == 0:
            raise ValueError("LLM returned empty or non-list JSON")

        sections = []
        for item in parsed[:max_sections]:
            label = str(item.get("label", "Section")).strip()
            text = str(item.get("text", "")).strip()
            if text:
                sections.append(AgenticSection(label=label, text=text))

        if not sections:
            raise ValueError("No valid sections parsed from LLM output")

        return AgenticChunkResponse(sections=sections, model_used=str(model_used), fallback=False)

    except Exception as e:
        sections = _paragraph_fallback(doc_text, max_sections)
        return AgenticChunkResponse(sections=sections, model_used="", fallback=True)


def _paragraph_fallback(text: str, max_sections: int) -> List[AgenticSection]:
    """Simple paragraph-based fallback when LLM is unavailable."""
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text.strip()] if text.strip() else []

    # Merge tiny paragraphs until we have reasonable sections
    sections: List[AgenticSection] = []
    current = ""
    for p in paragraphs:
        if not current:
            current = p
        elif len(current) + len(p) + 2 < 2000:
            current += "\n\n" + p
        else:
            sections.append(AgenticSection(label=f"Section {len(sections) + 1}", text=current))
            current = p
    if current:
        sections.append(AgenticSection(label=f"Section {len(sections) + 1}", text=current))

    return sections[:max_sections]


class GraphRAGRetrieveRequest(BaseModel):
    """LazyGraphRAG retrieval: vector search → Claude relevance check → expand if needed."""
    query: str
    initial_chunks: List[Dict[str, Any]] = Field(
        ..., description="Initial vector search results: [{id, text, score, metadata}]"
    )
    neighboring_chunks: List[Dict[str, Any]] = Field(
        default_factory=list,
        description="Neighboring/related chunks to expand into if initial results are insufficient"
    )
    min_relevant_chunks: int = Field(default=2, description="Minimum chunks needed before answering")
    user_id: Optional[str] = Field(default=None, description="User ID for ACL filtering")


class RelevanceAssessment(BaseModel):
    chunk_id: str
    is_relevant: bool
    relevance_score: float = 0.0
    reasoning: str = ""


class GraphRAGRetrieveResponse(BaseModel):
    relevant_chunks: List[Dict[str, Any]]
    expanded: bool = False
    assessment_details: List[RelevanceAssessment] = []
    total_assessed: int = 0


@app.post("/graphrag/retrieve", response_model=GraphRAGRetrieveResponse)
async def graphrag_retrieve(request: GraphRAGRetrieveRequest, auth: AuthContext = Depends(get_auth_context)):
    """
    LazyGraphRAG retrieval pipeline:
    1. Receive initial vector search results (pre-filtered by ACL via `acl_metadata_filter(auth)`)
    2. Use Claude to assess relevance of each chunk to the query
    3. If insufficient relevant chunks found, expand to neighboring chunks
    4. Return only genuinely relevant chunks for final synthesis

    The caller must apply `acl_metadata_filter(auth)` during the upstream vector search
    to ensure users never retrieve documents they don't have access to.
    """
    if not ANTHROPIC_API_KEY or not _anthropic_sdk_available:
        # No Claude — return initial chunks as-is (fallback to flat retrieval)
        return GraphRAGRetrieveResponse(
            relevant_chunks=request.initial_chunks,
            expanded=False,
            total_assessed=len(request.initial_chunks),
        )

    client = _get_anthropic_async_client()

    async def _assess_chunk(chunk: Dict[str, Any]) -> RelevanceAssessment:
        """Use Claude to assess whether a chunk is relevant to the query."""
        chunk_text = chunk.get("text", "")[:1500]
        chunk_id = chunk.get("id", "unknown")
        try:
            message = await _call_claude_with_fallback(
                client,
                preferred_model=HAIKU_MODEL,
                max_tokens=150,
                temperature=0.0,
                messages=[{
                    "role": "user",
                    "content": (
                        f"Query: {request.query}\n\n"
                        f"Document chunk:\n{chunk_text}\n\n"
                        "Is this chunk RELEVANT to answering the query? "
                        "Reply with JSON: {\"relevant\": true/false, \"score\": 0.0-1.0, \"reason\": \"...\"}"
                    ),
                }],
            )
            raw = "".join(b.text for b in message.content if hasattr(b, "text"))
            try:
                assessment = json.loads(raw)
            except json.JSONDecodeError:
                assessment = {"relevant": True, "score": 0.5, "reason": "parse error"}
            return RelevanceAssessment(
                chunk_id=chunk_id,
                is_relevant=assessment.get("relevant", True),
                relevance_score=float(assessment.get("score", 0.5)),
                reasoning=assessment.get("reason", ""),
            )
        except Exception as e:
            return RelevanceAssessment(
                chunk_id=chunk_id,
                is_relevant=True,  # Fail open
                relevance_score=0.5,
                reasoning=f"Assessment failed: {e}",
            )

    # Step 1: Assess initial chunks in parallel
    assessments = await asyncio.gather(
        *[_assess_chunk(chunk) for chunk in request.initial_chunks]
    )

    relevant_chunks = [
        chunk for chunk, assessment in zip(request.initial_chunks, assessments)
        if assessment.is_relevant and assessment.relevance_score >= 0.3
    ]
    all_assessments = list(assessments)

    # Step 2: If insufficient, expand to neighboring chunks
    expanded = False
    if len(relevant_chunks) < request.min_relevant_chunks and request.neighboring_chunks:
        expanded = True
        neighbor_assessments = await asyncio.gather(
            *[_assess_chunk(chunk) for chunk in request.neighboring_chunks]
        )
        for chunk, assessment in zip(request.neighboring_chunks, neighbor_assessments):
            if assessment.is_relevant and assessment.relevance_score >= 0.3:
                relevant_chunks.append(chunk)
            all_assessments.append(assessment)

    # Sort by relevance score
    chunk_scores = {a.chunk_id: a.relevance_score for a in all_assessments}
    relevant_chunks.sort(
        key=lambda c: chunk_scores.get(c.get("id", ""), 0.0), reverse=True
    )

    return GraphRAGRetrieveResponse(
        relevant_chunks=relevant_chunks,
        expanded=expanded,
        assessment_details=all_assessments,
        total_assessed=len(all_assessments),
    )


# ---------------------------------------------------------------------------
#  /suggest-connections — AI-powered connection suggestions
# ---------------------------------------------------------------------------

class SuggestConnectionsRequest(BaseModel):
    """Ask Claude to suggest company connections based on documents + existing graph."""
    company_name: str = ""                       # Optional: focus on a specific company
    question: str = ""                            # Optional: the user's question that triggered this
    sources: List[AskSource] = []                # Document sources for context
    existing_connections: List[AskConnection] = []  # Current graph state
    max_suggestions: int = 5

class SuggestedConnection(BaseModel):
    source_company: str
    target_company: str
    connection_type: str              # BD, INV, Knowledge, Partnership, Portfolio
    reasoning: str                    # Why this connection makes sense
    confidence: float = 0.0           # 0.0 – 1.0

class SuggestConnectionsResponse(BaseModel):
    suggestions: List[SuggestedConnection] = []
    context_summary: str = ""         # Brief explanation of what was analyzed


@app.post("/suggest-connections", response_model=SuggestConnectionsResponse)
async def suggest_connections(request: SuggestConnectionsRequest, auth: AuthContext = Depends(get_auth_context)):
    """
    Given document sources and the existing connections graph, use Claude to
    suggest new company connections the user hasn't logged yet.
    """
    # Check Anthropic availability with detailed logging
    if not _anthropic_sdk_available:
        return SuggestConnectionsResponse(
            suggestions=[],
            context_summary="Connection suggestions require Anthropic SDK. Please install: pip install anthropic"
        )
    if not ANTHROPIC_API_KEY:
        return SuggestConnectionsResponse(
            suggestions=[],
            context_summary="Connection suggestions require Anthropic API key. Please set ANTHROPIC_API_KEY in your environment variables."
        )

    # Build context about existing connections
    existing_lines: List[str] = []
    for conn in request.existing_connections:
        existing_lines.append(
            f"- {conn.source_company_name} → {conn.target_company_name} "
            f"({conn.connection_type}, {conn.connection_status})"
        )
    existing_block = "\n".join(existing_lines) if existing_lines else "No existing connections."

    # Build context about sources
    source_lines: List[str] = []
    for idx, src in enumerate(request.sources[:10], start=1):
        title = src.title or src.file_name or f"Source {idx}"
        snippet = (src.snippet or "")[:500]
        source_lines.append(f"[{idx}] {title}\n{snippet}")
    sources_block = "\n\n".join(source_lines) if source_lines else "No sources."

    company_focus = ""
    if request.company_name:
        company_focus = f"\nFocus especially on connections involving: {request.company_name}"

    question_context = ""
    if request.question:
        question_context = f"\nThe user asked: \"{request.question}\""

    prompt = f"""You are Company Assistant, a VC intelligence system. Analyze the following document sources and existing company connections graph.
Suggest up to {request.max_suggestions} NEW company connections that are NOT already in the graph.

{question_context}
{company_focus}

Existing Connections Graph:
{existing_block}

Document Sources:
{sources_block}

For each suggested connection, return a JSON array of objects with:
- "source_company": company name
- "target_company": company name
- "connection_type": one of "BD", "INV", "Knowledge", "Partnership", "Portfolio"
- "reasoning": brief explanation of why this connection makes sense (2-3 sentences)
- "confidence": float between 0.0 and 1.0

Return ONLY the JSON array, no markdown or explanation.
If you cannot suggest any meaningful connections, return an empty array: []"""

    # Use tool_choice for structured output
    tool_def = {
        "name": "suggest_connections",
        "description": "Return an array of suggested company connections.",
        "input_schema": {
            "type": "object",
            "properties": {
                "suggestions": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "source_company": {"type": "string"},
                            "target_company": {"type": "string"},
                            "connection_type": {"type": "string", "enum": ["BD", "INV", "Knowledge", "Partnership", "Portfolio"]},
                            "reasoning": {"type": "string"},
                            "confidence": {"type": "number"},
                        },
                        "required": ["source_company", "target_company", "connection_type", "reasoning"],
                    },
                },
                "context_summary": {"type": "string"},
            },
            "required": ["suggestions", "context_summary"],
        },
    }

    try:
        client = _get_anthropic_async_client()
        message = await _call_claude_with_fallback(
            client,
            max_tokens=2048,
            temperature=0.2,
            tools=[tool_def],
            tool_choice={"type": "tool", "name": "suggest_connections"},
            messages=[{"role": "user", "content": prompt}],
        )

        if message.stop_reason == "tool_use" and message.content:
            tool_use = message.content[0]
            data = tool_use.input
            suggestions = [
                SuggestedConnection(
                    source_company=s.get("source_company", ""),
                    target_company=s.get("target_company", ""),
                    connection_type=s.get("connection_type", "Knowledge"),
                    reasoning=s.get("reasoning", ""),
                    confidence=float(s.get("confidence", 0.5)),
                )
                for s in data.get("suggestions", [])
            ]
            return SuggestConnectionsResponse(
                suggestions=suggestions[:request.max_suggestions],
                context_summary=data.get("context_summary", ""),
            )

        return SuggestConnectionsResponse(suggestions=[], context_summary="Claude did not return structured suggestions.")

    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Connection suggestion failed: {str(e)}")


@app.post("/ingest/google-drive", response_model=GoogleDriveIngestResponse)
async def ingest_google_drive(request: GoogleDriveIngestRequest):
    if not request.access_token:
        raise HTTPException(status_code=400, detail="Access token required for private Google Drive files.")
    
    kind, file_id = parse_google_drive_url(request.url.strip())
    
    if kind not in ["document", "presentation", "spreadsheet"]:
        raise HTTPException(status_code=400, detail="Drive file URLs require a Docs/Slides/Sheets link.")

    # Use Google Drive API v3 to download the file
    # For Google Docs/Sheets/Slides, we need to export them
    if kind == "document":
        # Export as plain text
        api_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=text/plain"
        source_type = "notes"
    elif kind == "presentation":
        # Export as plain text
        api_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=text/plain"
        source_type = "deck"
    elif kind == "spreadsheet":
        # Export as CSV
        api_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=text/csv"
        source_type = "notes"
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")

    headers = {
        "Authorization": f"Bearer {request.access_token}",
        "Accept": "text/plain" if kind != "spreadsheet" else "text/csv"
    }

    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(api_url, headers=headers)
        if res.status_code >= 400:
            error_detail = res.text[:500] if res.text else "No error details"
            raise HTTPException(
                status_code=res.status_code,
                detail=f"Google Drive API failed (status {res.status_code}): {error_detail}. Make sure you have Drive access and the file is accessible."
            )
        content = res.text
        
        # Log if content is empty
        if not content or len(content.strip()) == 0:
            content = f"[Empty content from Google Drive file: {file_id}]"

    title = f"{kind}-{file_id[:8]}"
    return GoogleDriveIngestResponse(title=title, content=content, raw_content=content, sourceType=source_type)


# ──────────────────────────────────────────────────────────────────────────────
# Google Drive Folder-Sync Endpoints
# ──────────────────────────────────────────────────────────────────────────────

GDRIVE_API = "https://www.googleapis.com/drive/v3/files"
GOOGLE_OAUTH_TOKEN_URL = "https://oauth2.googleapis.com/token"
GOOGLE_OAUTH_AUTH_URL = "https://accounts.google.com/o/oauth2/v2/auth"
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET", "")
FRONTEND_ORIGIN = os.getenv("FRONTEND_ORIGIN", "https://general-platform.vercel.app")

# In-memory store: user_id -> (access_token, refresh_token, expires_at). Set GOOGLE_TOKENS_TABLE for persistence.
_user_google_tokens: Dict[str, Tuple[str, str, float]] = {}


async def _get_user_id_from_supabase_token(access_token: str) -> str:
    """Return user id (sub) from Supabase access token. Uses Supabase Auth API."""
    if not SUPABASE_URL or not SUPABASE_SERVICE_KEY:
        raise HTTPException(status_code=503, detail="Supabase not configured.")
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            r = await client.get(
                f"{SUPABASE_URL.rstrip('/')}/auth/v1/user",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "apikey": SUPABASE_SERVICE_KEY,
                },
            )
            if r.status_code != 200:
                raise HTTPException(status_code=401, detail="Invalid or expired Supabase token.")
            data = r.json()
            return data.get("id") or data.get("sub") or ""
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=401, detail=f"Could not verify token: {e}")


class GoogleDriveStartRequest(BaseModel):
    access_token: str  # Supabase user JWT


class GoogleDriveStartResponse(BaseModel):
    redirect_url: str


@app.post("/auth/google-drive/start", response_model=GoogleDriveStartResponse)
async def auth_google_drive_start(request: GoogleDriveStartRequest):
    """Start Google Drive OAuth: verify user, return URL to redirect to. Frontend then redirects to that URL."""
    user_id = await _get_user_id_from_supabase_token(request.access_token)
    if not user_id:
        raise HTTPException(status_code=401, detail="Could not get user id.")
    # State: base64(user_id) so callback knows who to store tokens for. In prod you'd sign it.
    import base64
    state = base64.urlsafe_b64encode(user_id.encode()).decode().rstrip("=")
    redirect_path = f"/auth/google-drive?state={state}"
    base_url = os.getenv("BACKEND_PUBLIC_URL", "https://general-platform.onrender.com").rstrip("/")
    return GoogleDriveStartResponse(redirect_url=f"{base_url}{redirect_path}")


@app.get("/auth/google-drive")
async def auth_google_drive_redirect(state: str):
    """Redirect browser to Google OAuth. state = base64(user_id)."""
    if not GOOGLE_CLIENT_ID:
        raise HTTPException(status_code=503, detail="GOOGLE_CLIENT_ID not set.")
    import base64
    try:
        _b64_decode_user_id(state)  # validate
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid state.")
    redirect_uri = f"{os.getenv('BACKEND_PUBLIC_URL', 'https://general-platform.onrender.com').rstrip('/')}/auth/google-drive/callback"
    params = {
        "client_id": GOOGLE_CLIENT_ID,
        "redirect_uri": redirect_uri,
        "response_type": "code",
        "scope": "https://www.googleapis.com/auth/drive.readonly https://www.googleapis.com/auth/gmail.readonly",
        "access_type": "offline",
        "prompt": "consent",
        "state": state,
    }
    from urllib.parse import urlencode
    url = f"{GOOGLE_OAUTH_AUTH_URL}?{urlencode(params)}"
    return RedirectResponse(url=url)


def _b64_decode_user_id(state: str) -> str:
    import base64
    pad = (4 - len(state) % 4) % 4
    raw = base64.urlsafe_b64decode(state + "=" * pad)
    return raw.decode()


@app.get("/auth/google-drive/callback")
async def auth_google_drive_callback(code: str, state: str):
    """Exchange code for tokens, store by user_id, redirect to frontend."""
    if not GOOGLE_CLIENT_ID or not GOOGLE_CLIENT_SECRET:
        raise HTTPException(status_code=503, detail="Google OAuth not configured.")
    try:
        user_id = _b64_decode_user_id(state)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid state.")
    redirect_uri = f"{os.getenv('BACKEND_PUBLIC_URL', 'https://general-platform.onrender.com').rstrip('/')}/auth/google-drive/callback"
    body = {
        "code": code,
        "client_id": GOOGLE_CLIENT_ID,
        "client_secret": GOOGLE_CLIENT_SECRET,
        "redirect_uri": redirect_uri,
        "grant_type": "authorization_code",
    }
    async with httpx.AsyncClient(timeout=15.0) as client:
        r = await client.post(GOOGLE_OAUTH_TOKEN_URL, data=body)
        if r.status_code >= 400:
            raise HTTPException(status_code=400, detail=f"Google token exchange failed: {r.text[:300]}")
        data = r.json()
    access_token = data.get("access_token")
    refresh_token = data.get("refresh_token")
    if not access_token:
        raise HTTPException(status_code=502, detail="Google did not return access_token.")
    expires_in = data.get("expires_in", 3600)
    _user_google_tokens[user_id] = (access_token, refresh_token or "", time.time() + expires_in)
    return RedirectResponse(url=f"{FRONTEND_ORIGIN.rstrip('/')}/?google_drive=connected")


class GDriveMyTokenResponse(BaseModel):
    access_token: Optional[str] = None  # null when Drive not connected (200 OK, no 404)


@app.get("/gdrive/my-token", response_model=GDriveMyTokenResponse)
async def gdrive_my_token(authorization: Optional[str] = Header(default=None)):
    """Return Google access token for the current user (from Supabase JWT). Refreshes if expired. Returns access_token=null when not connected (200 OK)."""
    if not authorization or "Bearer " not in authorization:
        raise HTTPException(status_code=401, detail="Authorization: Bearer <supabase_access_token> required.")
    token = authorization.replace("Bearer ", "").strip()
    user_id = await _get_user_id_from_supabase_token(token)
    if not user_id:
        raise HTTPException(status_code=401, detail="Could not get user id.")
    stored = _user_google_tokens.get(user_id)
    if not stored:
        return GDriveMyTokenResponse(access_token=None)  # 200 OK, no 404 log spam
    access_token, refresh_token, expires_at = stored
    if refresh_token and time.time() > expires_at - 60:
        # Refresh
        body = {
            "grant_type": "refresh_token",
            "client_id": GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "refresh_token": refresh_token,
        }
        async with httpx.AsyncClient(timeout=15.0) as client:
            r = await client.post(GOOGLE_OAUTH_TOKEN_URL, data=body)
            if r.status_code == 200:
                data = r.json()
                access_token = data.get("access_token", access_token)
                exp = data.get("expires_in", 3600)
                _user_google_tokens[user_id] = (access_token, refresh_token, time.time() + exp)
    return GDriveMyTokenResponse(access_token=access_token)


@app.post("/gdrive/refresh-token", response_model=GDriveRefreshTokenResponse)
async def gdrive_refresh_token(request: GDriveRefreshTokenRequest):
    """Exchange a Google refresh_token for a new access_token. Uses GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET."""
    if not GOOGLE_CLIENT_ID or not GOOGLE_CLIENT_SECRET:
        raise HTTPException(
            status_code=503,
            detail="Google OAuth refresh not configured. Set GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET on the server.",
        )
    body = {
        "grant_type": "refresh_token",
        "client_id": GOOGLE_CLIENT_ID,
        "client_secret": GOOGLE_CLIENT_SECRET,
        "refresh_token": request.refresh_token,
    }
    async with httpx.AsyncClient(timeout=15.0) as client:
        res = await client.post(GOOGLE_OAUTH_TOKEN_URL, data=body)
        if res.status_code >= 400:
            raise HTTPException(
                status_code=401,
                detail=f"Google token refresh failed: {res.text[:300]}",
            )
        data = res.json()
    access_token = data.get("access_token")
    if not access_token:
        raise HTTPException(status_code=502, detail="Google did not return an access_token.")
    return GDriveRefreshTokenResponse(access_token=access_token)


@app.post("/gdrive/list-folders", response_model=GDriveListFoldersResponse)
async def gdrive_list_folders(request: GDriveListFoldersRequest):
    """List sub-folders inside a given Google Drive folder."""
    headers = {"Authorization": f"Bearer {request.access_token}"}
    query = f"'{request.folder_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"
    params = {
        "q": query,
        "fields": "files(id,name,modifiedTime)",
        "pageSize": "1000",
        "supportsAllDrives": "true",
        "includeItemsFromAllDrives": "true",
    }
    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(GDRIVE_API, headers=headers, params=params)
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code, detail=f"Drive API error: {res.text[:500]}")
        data = res.json()
    folders = [
        GDriveFolderEntry(id=f["id"], name=f["name"], modifiedTime=f.get("modifiedTime"))
        for f in data.get("files", [])
    ]
    folders.sort(key=lambda f: (f.name or "").lower())
    return GDriveListFoldersResponse(folders=folders)


@app.post("/gdrive/list-files", response_model=GDriveListFilesResponse)
async def gdrive_list_files(request: GDriveListFilesRequest):
    """List all non-folder files inside a given Google Drive folder."""
    headers = {"Authorization": f"Bearer {request.access_token}"}
    query = f"'{request.folder_id}' in parents and mimeType!='application/vnd.google-apps.folder' and trashed=false"
    params = {
        "q": query,
        "fields": "files(id,name,mimeType,modifiedTime,size)",
        "pageSize": "1000",
        "supportsAllDrives": "true",
        "includeItemsFromAllDrives": "true",
    }
    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(GDRIVE_API, headers=headers, params=params)
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code, detail=f"Drive API error: {res.text[:500]}")
        data = res.json()
    files = [
        GDriveFileEntry(
            id=f["id"],
            name=f["name"],
            mimeType=f.get("mimeType", ""),
            modifiedTime=f.get("modifiedTime"),
            size=f.get("size"),
        )
        for f in data.get("files", [])
    ]
    files.sort(key=lambda f: (f.name or "").lower())
    return GDriveListFilesResponse(files=files)


# Mime-type mappings for Google-native → export
_GOOGLE_EXPORT_MAP = {
    "application/vnd.google-apps.document": ("text/plain", "notes"),
    "application/vnd.google-apps.presentation": ("text/plain", "deck"),
    "application/vnd.google-apps.spreadsheet": ("text/csv", "notes"),
}

# Supported binary/text types we can handle
_SUPPORTED_BINARY = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/msword",
    "application/vnd.ms-excel",
    "application/vnd.ms-powerpoint",
    "text/plain",
    "text/csv",
    "text/markdown",
    "application/json",
}


@app.post("/gdrive/download-file", response_model=GDriveDownloadFileResponse)
async def gdrive_download_file(request: GDriveDownloadFileRequest):
    """Download or export a single Google Drive file by its ID.

    Google-native types (Docs/Sheets/Slides) are exported to text.
    Binary types (PDF, DOCX, PPTX, etc.) are downloaded and have text extracted.
    """
    mime = request.mime_type or ""
    file_id = request.file_id
    file_name = request.file_name or file_id[:12]
    auth_headers = {"Authorization": f"Bearer {request.access_token}"}

    content = ""
    source_type = "notes"
    final_mime = mime

    async with httpx.AsyncClient(timeout=60.0) as client:
        # ── Google-native: export as text ──
        if mime in _GOOGLE_EXPORT_MAP:
            export_mime, source_type = _GOOGLE_EXPORT_MAP[mime]
            url = f"{GDRIVE_API}/{file_id}/export"
            res = await client.get(url, headers=auth_headers, params={"mimeType": export_mime})
            if res.status_code >= 400:
                raise HTTPException(status_code=res.status_code, detail=f"Drive export error: {res.text[:500]}")
            content = res.text
            final_mime = export_mime

        # ── Binary / text files: download raw ──
        elif mime in _SUPPORTED_BINARY or mime.startswith("text/"):
            url = f"{GDRIVE_API}/{file_id}"
            res = await client.get(url, headers=auth_headers, params={"alt": "media"})
            if res.status_code >= 400:
                raise HTTPException(status_code=res.status_code, detail=f"Drive download error: {res.text[:500]}")

            if mime.startswith("text/") or mime == "application/json":
                content = res.text
                source_type = "notes"
            elif mime == "application/pdf":
                # Extract text from PDF bytes
                try:
                    import fitz  # PyMuPDF
                    pdf_bytes = res.content
                    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                    pages = []
                    for page in doc:
                        pages.append(page.get_text())
                    content = "\n\n".join(pages)
                    doc.close()
                    source_type = "deck"
                except ImportError:
                    content = f"[PDF file: {file_name} — PyMuPDF not available for extraction]"
                    source_type = "deck"
                except Exception as e:
                    content = f"[PDF extraction failed for {file_name}: {str(e)[:200]}]"
                    source_type = "deck"
            else:
                # For DOCX/PPTX/XLS — attempt basic text extraction
                try:
                    raw_bytes = res.content
                    # Try python-docx for docx
                    if "wordprocessingml" in mime or mime == "application/msword":
                        try:
                            from docx import Document as DocxDocument
                            doc = DocxDocument(BytesIO(raw_bytes))
                            content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                            source_type = "notes"
                        except Exception:
                            content = f"[Word document: {file_name} — text extraction failed]"
                    elif "presentationml" in mime or mime == "application/vnd.ms-powerpoint":
                        try:
                            from pptx import Presentation
                            prs = Presentation(BytesIO(raw_bytes))
                            slides_text = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if shape.has_text_frame:
                                        slides_text.append(shape.text_frame.text)
                            content = "\n\n".join(slides_text)
                            source_type = "deck"
                        except Exception:
                            content = f"[Presentation: {file_name} — text extraction failed]"
                            source_type = "deck"
                    elif "spreadsheetml" in mime or mime == "application/vnd.ms-excel":
                        try:
                            import openpyxl
                            wb = openpyxl.load_workbook(BytesIO(raw_bytes), read_only=True, data_only=True)
                            rows = []
                            for sheet in wb.sheetnames:
                                ws = wb[sheet]
                                for row in ws.iter_rows(values_only=True):
                                    rows.append(",".join(str(c) if c is not None else "" for c in row))
                            content = "\n".join(rows)
                            source_type = "notes"
                        except Exception:
                            content = f"[Spreadsheet: {file_name} — text extraction failed]"
                    else:
                        content = f"[Unsupported binary type: {mime} for file {file_name}]"
                except Exception as e:
                    content = f"[File extraction failed for {file_name}: {str(e)[:200]}]"
        elif mime.startswith("image/"):
            # Download image and describe with Claude Vision so the model can read pictures
            url = f"{GDRIVE_API}/{file_id}"
            res = await client.get(url, headers=auth_headers, params={"alt": "media"})
            if res.status_code >= 400:
                raise HTTPException(status_code=res.status_code, detail=f"Drive download error: {res.text[:500]}")
            image_bytes = res.content
            if not image_bytes:
                content = f"[Empty image: {file_name}]"
            else:
                content = await _describe_image_with_vision(
                    image_bytes,
                    mime if mime in ("image/png", "image/jpeg", "image/gif", "image/webp") else "image/png",
                    file_name,
                )
            source_type = "notes"
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {mime}")

    if not content or len(content.strip()) == 0:
        content = f"[Empty content from Drive file: {file_name}]"

    return GDriveDownloadFileResponse(
        title=file_name,
        content=content,
        raw_content=content,
        sourceType=source_type,
        mimeType=final_mime,
    )


# ─────────── Gmail API ───────────────────────────────────────────────────────

GMAIL_API = "https://gmail.googleapis.com/gmail/v1/users/me"


def _parse_gmail_headers(headers: List[Dict[str, str]]) -> Dict[str, str]:
    """Extract useful headers from Gmail message payload."""
    out: Dict[str, str] = {}
    for h in headers:
        name = h.get("name", "").lower()
        if name in ("from", "to", "cc", "subject", "date"):
            out[name] = h.get("value", "")
    return out


def _parse_email_addresses(raw: str) -> List[str]:
    """Split a comma-separated header like To/Cc into individual addresses."""
    if not raw:
        return []
    return [addr.strip() for addr in raw.split(",") if addr.strip()]


def _extract_body_from_parts(parts: List[dict], prefer_plain: bool = True) -> Tuple[str, str]:
    """Recursively walk MIME parts and return (plain_text, html_text)."""
    plain, html = "", ""
    for part in parts:
        mime = part.get("mimeType", "")
        body_data = part.get("body", {}).get("data", "")
        nested = part.get("parts", [])

        if nested:
            p, h = _extract_body_from_parts(nested, prefer_plain)
            if p: plain = p
            if h: html = h
        elif mime == "text/plain" and body_data:
            plain = base64.urlsafe_b64decode(body_data + "==").decode("utf-8", errors="replace")
        elif mime == "text/html" and body_data:
            html = base64.urlsafe_b64decode(body_data + "==").decode("utf-8", errors="replace")
    return plain, html


def _strip_html(html: str) -> str:
    """Rough HTML-to-text conversion for email bodies."""
    text = re.sub(r"<style[^>]*>[\s\S]*?</style>", "", html, flags=re.I)
    text = re.sub(r"<script[^>]*>[\s\S]*?</script>", "", text, flags=re.I)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.I)
    text = re.sub(r"</p>", "\n\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r" {2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_attachments_meta(parts: List[dict]) -> List[GmailAttachmentMeta]:
    """Walk MIME parts to find attachment metadata."""
    attachments: List[GmailAttachmentMeta] = []
    for part in parts:
        nested = part.get("parts", [])
        if nested:
            attachments.extend(_extract_attachments_meta(nested))
        filename = part.get("filename", "")
        att_id = part.get("body", {}).get("attachmentId", "")
        if filename and att_id:
            attachments.append(GmailAttachmentMeta(
                id=att_id,
                filename=filename,
                mimeType=part.get("mimeType", "application/octet-stream"),
                size=part.get("body", {}).get("size", 0),
            ))
    return attachments


@app.post("/gmail/list-messages", response_model=GmailListMessagesResponse)
async def gmail_list_messages(request: GmailListMessagesRequest):
    """List Gmail messages matching a query and/or label filter."""
    headers = {"Authorization": f"Bearer {request.access_token}"}
    params: Dict[str, Any] = {"maxResults": min(request.max_results, 500)}
    if request.query:
        params["q"] = request.query
    if request.label_ids:
        params["labelIds"] = request.label_ids
    if request.page_token:
        params["pageToken"] = request.page_token

    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(f"{GMAIL_API}/messages", headers=headers, params=params)
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code,
                                detail=f"Gmail API list error: {res.text[:500]}")
        data = res.json()

    messages_raw = data.get("messages", [])
    snippets: List[GmailMessageSnippet] = []
    for m in messages_raw:
        snippets.append(GmailMessageSnippet(
            id=m["id"],
            threadId=m.get("threadId", m["id"]),
            snippet=None,
        ))

    return GmailListMessagesResponse(
        messages=snippets,
        next_page_token=data.get("nextPageToken"),
        result_size_estimate=data.get("resultSizeEstimate", len(snippets)),
    )


@app.post("/gmail/get-message", response_model=GmailGetMessageResponse)
async def gmail_get_message(request: GmailGetMessageRequest):
    """Fetch a full Gmail message by ID, parse headers, body, and attachments."""
    headers = {"Authorization": f"Bearer {request.access_token}"}

    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(
            f"{GMAIL_API}/messages/{request.message_id}",
            headers=headers,
            params={"format": "full"},
        )
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code,
                                detail=f"Gmail API get-message error: {res.text[:500]}")
        msg = res.json()

    payload = msg.get("payload", {})
    hdrs = _parse_gmail_headers(payload.get("headers", []))
    parts = payload.get("parts", [])

    if parts:
        body_plain, body_html = _extract_body_from_parts(parts)
        attachments = _extract_attachments_meta(parts)
    else:
        body_data = payload.get("body", {}).get("data", "")
        mime = payload.get("mimeType", "")
        if body_data:
            decoded = base64.urlsafe_b64decode(body_data + "==").decode("utf-8", errors="replace")
            body_plain = decoded if mime == "text/plain" else ""
            body_html = decoded if mime == "text/html" else ""
        else:
            body_plain, body_html = "", ""
        attachments = []

    return GmailGetMessageResponse(
        id=msg["id"],
        threadId=msg.get("threadId", msg["id"]),
        subject=hdrs.get("subject", ""),
        sender=hdrs.get("from", ""),
        to=_parse_email_addresses(hdrs.get("to", "")),
        cc=_parse_email_addresses(hdrs.get("cc", "")),
        date=hdrs.get("date"),
        body_text=body_plain,
        body_html=body_html,
        labels=msg.get("labelIds", []),
        attachments=attachments,
    )


@app.post("/gmail/download-attachment", response_model=GmailDownloadAttachmentResponse)
async def gmail_download_attachment(request: GmailDownloadAttachmentRequest):
    """Download a single email attachment by its ID."""
    headers = {"Authorization": f"Bearer {request.access_token}"}

    async with httpx.AsyncClient(timeout=60.0) as client:
        res = await client.get(
            f"{GMAIL_API}/messages/{request.message_id}/attachments/{request.attachment_id}",
            headers=headers,
        )
        if res.status_code >= 400:
            raise HTTPException(status_code=res.status_code,
                                detail=f"Gmail attachment error: {res.text[:500]}")
        data = res.json()

    raw_b64 = data.get("data", "")
    size = data.get("size", 0)

    return GmailDownloadAttachmentResponse(
        data=raw_b64,
        filename="",
        mimeType="application/octet-stream",
        size=size,
    )


@app.post("/ingest/gmail", response_model=GmailIngestResponse)
async def ingest_gmail(request: GmailIngestRequest):
    """Fetch, parse, and return a structured email document ready for storage."""
    msg_req = GmailGetMessageRequest(
        access_token=request.access_token,
        message_id=request.message_id,
    )
    msg = await gmail_get_message(msg_req)

    body = msg.body_text
    if not body.strip() and msg.body_html:
        body = _strip_html(msg.body_html)

    subject = msg.subject or "(no subject)"
    sender = msg.sender or "unknown"
    date_str = msg.date or ""

    header_block = f"From: {sender}\nTo: {', '.join(msg.to)}\n"
    if msg.cc:
        header_block += f"Cc: {', '.join(msg.cc)}\n"
    header_block += f"Date: {date_str}\nSubject: {subject}\n"

    content = f"{header_block}\n{body}"
    title = f"[Email] {subject}"

    return GmailIngestResponse(
        title=title,
        content=content,
        raw_content=content,
        sourceType="gmail",
        email_from=sender,
        email_to=msg.to,
        email_cc=msg.cc,
        email_subject=subject,
        email_date=date_str,
        gmail_thread_id=msg.threadId,
        gmail_labels=msg.labels,
        has_attachments=len(msg.attachments) > 0,
        attachments=msg.attachments,
    )


class ValidationRequest(BaseModel):
    data: str
    dataType: Optional[str] = None

class ValidationResponse(BaseModel):
    isValid: bool
    missingFields: Dict[str, List[str]]  # { "startups": ["geoMarkets"], "investors": ["minTicketSize"] }
    incompleteFields: Dict[str, List[str]]  # Fields that exist but are incomplete
    suggestions: List[str]  # Suggestions for what to add
    extractedData: Dict[str, Any]  # What was successfully extracted

@app.post("/validate-file", response_model=FileValidationResponse)
async def validate_file(file: UploadFile = File(...), dataType: Optional[str] = None):
    """
    Validate an uploaded file (any supported format) and return:
    - row-level errors with explicit row numbers
    - CSV templates with extracted rows prefilled and missing columns preserved
    """
    try:
        file_ext, text_content = await extract_text_content(file)
        conversion_request = ConversionRequest(
            data=text_content,
            dataType=dataType,
            format=file_ext
        )
        conversion_result = await convert_data(conversion_request)

        row_errors = validate_structured_rows(conversion_result.startups, conversion_result.investors)
        errors = (conversion_result.errors or []) + row_errors
        warnings = conversion_result.warnings or []

        startup_csv = build_startup_csv(conversion_result.startups) if conversion_result.startups else build_startup_csv([])
        investor_csv = build_investor_csv(conversion_result.investors) if conversion_result.investors else build_investor_csv([])

        return FileValidationResponse(
            isValid=len(errors) == 0,
            errors=errors,
            warnings=warnings,
            detectedType=conversion_result.detectedType,
            startupCsvTemplate=startup_csv,
            investorCsvTemplate=investor_csv,
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File validation failed: {str(e)}")

@app.post("/validate", response_model=ValidationResponse)
async def validate_data(request: ValidationRequest):
    """
    Validate data and identify what's missing
    This is what the investment team needs - tells them what to add!
    """
    try:
        # First, try to convert the data
        conversion_request = ConversionRequest(
            data=request.data,
            dataType=request.dataType
        )
        conversion_result = await convert_data(conversion_request)
        
        missing_fields = {"startups": [], "investors": []}
        incomplete_fields = {"startups": [], "investors": []}
        suggestions = []
        
        # Check startups
        for startup in conversion_result.startups:
            startup_missing = []
            startup_incomplete = []
            
            if not startup.companyName or (isinstance(startup.companyName, str) and startup.companyName.strip() == ""):
                startup_missing.append("companyName")
            if not startup.geoMarkets or len(startup.geoMarkets) == 0:
                startup_missing.append("geoMarkets")
                suggestions.append(f"Add geographic markets for {startup.companyName}")
            if not startup.industry or (isinstance(startup.industry, str) and startup.industry.strip() == ""):
                startup_missing.append("industry")
            if not startup.fundingTarget or startup.fundingTarget == 0:
                startup_missing.append("fundingTarget")
                suggestions.append(f"Add funding target amount for {startup.companyName}")
            if not startup.fundingStage or (isinstance(startup.fundingStage, str) and startup.fundingStage.strip() == ""):
                startup_missing.append("fundingStage")
                suggestions.append(f"Add funding stage (Pre-seed, Seed, Series A, etc.) for {startup.companyName}")
            
            if startup_missing:
                missing_fields["startups"].extend(startup_missing)
        
        # Check investors
        for investor in conversion_result.investors:
            investor_missing = []
            investor_incomplete = []
            
            if not investor.firmName or (isinstance(investor.firmName, str) and investor.firmName.strip() == ""):
                investor_missing.append("firmName")
            if not investor.memberName or (isinstance(investor.memberName, str) and investor.memberName.strip() == ""):
                investor_missing.append("memberName")
                suggestions.append(f"Add investor member name (person) for {investor.firmName or 'this investor'}")
            if not investor.geoFocus or len(investor.geoFocus) == 0:
                investor_missing.append("geoFocus")
                suggestions.append(f"Add geographic focus for {investor.firmName}")
            if not investor.industryPreferences or len(investor.industryPreferences) == 0:
                investor_missing.append("industryPreferences")
                suggestions.append(f"Add industry preferences for {investor.firmName}")
            if not investor.stagePreferences or len(investor.stagePreferences) == 0:
                investor_missing.append("stagePreferences")
                suggestions.append(f"Add stage preferences (Seed, Series A, etc.) for {investor.firmName}")
            if not investor.minTicketSize or investor.minTicketSize == 0:
                investor_missing.append("minTicketSize")
                suggestions.append(f"Add minimum ticket size for {investor.firmName}")
            if not investor.maxTicketSize or investor.maxTicketSize == 0:
                investor_missing.append("maxTicketSize")
                suggestions.append(f"Add maximum ticket size for {investor.firmName}")
            if not investor.totalSlots or investor.totalSlots == 0:
                investor_missing.append("totalSlots")
                suggestions.append(f"Add number of meeting slots for {investor.firmName}")
            
            if investor_missing:
                missing_fields["investors"].extend(investor_missing)
        
        # Remove duplicates
        missing_fields["startups"] = list(set(missing_fields["startups"]))
        missing_fields["investors"] = list(set(missing_fields["investors"]))
        
        is_valid = (
            len(missing_fields["startups"]) == 0 and
            len(missing_fields["investors"]) == 0 and
            len(conversion_result.errors) == 0
        )
        
        return ValidationResponse(
            isValid=is_valid,
            missingFields=missing_fields,
            incompleteFields=incomplete_fields,
            suggestions=suggestions,
            extractedData={
                "startups": [s.dict() for s in conversion_result.startups],
                "investors": [i.dict() for i in conversion_result.investors],
                "detectedType": conversion_result.detectedType,
                "confidence": conversion_result.confidence
            }
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Validation failed: {str(e)}")

@app.on_event("startup")
async def startup_event():
    """Log configuration on startup for debugging"""
    if ANTHROPIC_API_KEY:
        # Show first/last 4 chars for verification (don't expose full key)
        key_preview = f"{ANTHROPIC_API_KEY[:4]}...{ANTHROPIC_API_KEY[-4:]}" if len(ANTHROPIC_API_KEY) > 8 else "***"


# ══════════════════════════════════════════════════════════════
#  COMPANY PLATFORM — Company-Agnostic Extensions
#  These endpoints let any company paste a description and get
#  a custom AI persona. Team members share the same context.
# ══════════════════════════════════════════════════════════════

from prompt_builder import generate_system_prompt as _generate_system_prompt, FALLBACK_PROMPT as _FALLBACK_PROMPT


class CompanySetupRequest(BaseModel):
    organization_id: str
    company_description: str
    company_name: Optional[str] = None


class CompanySetupResponse(BaseModel):
    system_prompt: str
    organization_id: str
    message: str


@app.options("/company/setup")
async def company_setup_options():
    """CORS preflight for /company/setup."""
    return {}

@app.post("/company/setup", response_model=CompanySetupResponse)
async def setup_company(request: CompanySetupRequest):
    """MD pastes company description → generate system prompt → store on org."""
    if not request.company_description.strip():
        raise HTTPException(status_code=400, detail="Company description is required")

    try:
        system_prompt = await _generate_system_prompt(request.company_description)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate AI context: {str(e)[:200]}")

    try:
        sb = get_supabase()
        if sb:
            sb.table("organizations").update({
                "company_description": request.company_description.strip(),
                "system_prompt": system_prompt,
            }).eq("id", request.organization_id).execute()
    except Exception as e:
        pass  # non-fatal: prompt was generated, frontend can show it

    return CompanySetupResponse(
        system_prompt=system_prompt,
        organization_id=request.organization_id,
        message="Company context saved. AI is now customized for your team.",
    )


@app.post("/company/regenerate-prompt", response_model=CompanySetupResponse)
async def regenerate_prompt(request: CompanySetupRequest):
    """Re-generate the system prompt if MD updates the company description."""
    return await setup_company(request)


@app.get("/company/{org_id}/context")
async def get_company_context(org_id: str):
    """Team members call this to get the shared AI context."""
    try:
        sb = get_supabase()
        if sb:
            result = sb.table("organizations").select(
                "id, name, company_description, system_prompt"
            ).eq("id", org_id).single().execute()
            if result.data:
                raw_prompt = result.data.get("system_prompt", "") or _FALLBACK_PROMPT
                raw_prompt = re.sub(r"\bOrbit\s+AI\b", "Company Assistant", raw_prompt, flags=re.IGNORECASE)
                return {
                    "organization_id": result.data["id"],
                    "company_name": result.data.get("name", ""),
                    "company_description": result.data.get("company_description", ""),
                    "system_prompt": raw_prompt,
                }
    except Exception as e:
        pass

    return {
        "organization_id": org_id,
        "company_description": "",
        "system_prompt": _FALLBACK_PROMPT,
    }


# ---------------------------------------------------------------------------
#  Agentic RAG — Tools, Handlers, System Prompt, /ask/agent/stream (same as remix)
# ---------------------------------------------------------------------------

AGENT_TOOLS = [
    {
        "name": "search_portfolio",
        "description": (
            "Search portfolio companies by structured metadata filters. "
            "Use this for questions about company counts, locations, sectors, funding stages, or listing companies. "
            "Returns matching company names and properties."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "country": {"type": "string", "description": "Filter by country or location."},
                "sector": {"type": "string", "description": "Filter by industry/sector."},
                "stage": {"type": "string", "description": "Filter by funding stage."},
                "name": {"type": "string", "description": "Search by company name (partial match)."},
                "business_model": {"type": "string", "description": "Filter by business model."},
                "list_all": {"type": "boolean", "description": "If true, return ALL portfolio companies. Default false."},
            },
            "required": [],
        },
    },
    {
        "name": "search_documents",
        "description": (
            "Semantic search across all uploaded documents (pitch decks, memos, meeting notes, reports). "
            "Returns relevant text chunks with document titles."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "The search query."},
                "company_name": {"type": "string", "description": "Optional: limit search to documents about this specific company."},
                "top_k": {"type": "integer", "description": "Number of results to return (default 10, max 20)."},
            },
            "required": ["query"],
        },
    },
    {
        "name": "get_company_details",
        "description": (
            "Get detailed information about a specific portfolio company including all properties, "
            "KPI summaries, document count, and relationships."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "company_name": {"type": "string", "description": "The exact or approximate company name to look up."},
            },
            "required": ["company_name"],
        },
    },
    {
        "name": "search_knowledge_graph",
        "description": (
            "Search the knowledge graph for entities and their relationships. "
            "Use this for questions about connections, investors, founders, competitors, partnerships."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "entity_name": {"type": "string", "description": "Name of the entity to search for."},
                "entity_type": {"type": "string", "description": "Entity type filter: 'company', 'person', 'fund', 'round', 'sector', 'location'."},
                "relationship_type": {"type": "string", "description": "Filter by relationship type."},
                "max_depth": {"type": "integer", "description": "Max traversal depth (1-3, default 1)."},
            },
            "required": ["entity_name"],
        },
    },
]


async def _agent_search_portfolio(tool_input: dict, event_id: str) -> str:
    sb = get_supabase()
    country = (tool_input.get("country") or "").strip().lower()
    sector = (tool_input.get("sector") or "").strip().lower()
    stage = (tool_input.get("stage") or "").strip().lower()
    name = (tool_input.get("name") or "").strip().lower()
    biz_model = (tool_input.get("business_model") or "").strip().lower()
    list_all = tool_input.get("list_all", False)

    query = sb.table("kg_entities").select("name, properties").eq("event_id", event_id).eq("entity_type", "company")
    result = query.execute()
    rows = result.data or []

    if not rows:
        return "No portfolio companies found in the database."

    def matches(row) -> bool:
        if list_all:
            return True
        props = row.get("properties") or {}
        text_blob = json.dumps(props).lower()
        company_name_lower = (row.get("name") or "").lower()
        if name and name not in company_name_lower and name not in text_blob:
            return False
        if country:
            geo_fields = " ".join(str(props.get(f, "")) for f in [
                "country", "headquarters", "location", "hq",
                "geo_focus", "geo_markets", "geography", "region", "regions",
            ]).lower()
            if country not in geo_fields and country not in text_blob:
                return False
        if sector:
            industry_fields = " ".join(str(props.get(f, "")) for f in ["industry", "sector", "vertical"]).lower()
            bio = str(props.get("bio", "")).lower()
            if sector not in industry_fields and sector not in bio:
                return False
        if stage:
            stage_fields = " ".join(str(props.get(f, "")) for f in ["funding_stage", "stage", "round", "funding_round"]).lower()
            if stage not in stage_fields:
                return False
        if biz_model:
            model_fields = " ".join(str(props.get(f, "")) for f in ["business_model", "model"]).lower()
            bio = str(props.get("bio", "")).lower()
            if biz_model not in model_fields and biz_model not in bio:
                return False
        return True

    matched = [r for r in rows if matches(r)]

    if not matched:
        all_names = ", ".join(r.get("name", "?") for r in rows[:30])
        return f"No companies matched the filters. Total portfolio: {len(rows)} companies. Names: {all_names}"

    lines = [f"Found {len(matched)} matching companies (out of {len(rows)} total in portfolio):\n"]
    for r in matched:
        props = r.get("properties") or {}
        parts = [f"- **{r.get('name', '?')}**"]
        for key in ["industry", "funding_stage", "headquarters", "country", "location", "business_model", "bio"]:
            val = props.get(key)
            if val:
                label = key.replace("_", " ").title()
                display = str(val)[:200] if key == "bio" else str(val)
                parts.append(f"  {label}: {display}")
        geo = props.get("geo_focus") or props.get("geo_markets") or props.get("geography")
        if geo:
            if isinstance(geo, list):
                parts.append(f"  Geo Focus: {', '.join(str(g) for g in geo)}")
            else:
                parts.append(f"  Geo Focus: {geo}")
        lines.append("\n".join(parts))

    return "\n\n".join(lines)


async def _agent_search_documents(tool_input: dict, event_id: str) -> str:
    sb = get_supabase()
    query_text = (tool_input.get("query") or "").strip()
    company_name = (tool_input.get("company_name") or "").strip()
    top_k = min(tool_input.get("top_k", 10), 20)

    if not query_text:
        return "Error: 'query' is required for document search."

    try:
        if EMBEDDINGS_PROVIDER == "voyage":
            embedding = await generate_embedding_voyage(query_text, "query")
        elif EMBEDDINGS_PROVIDER == "openai":
            embedding = await generate_embedding_openai(query_text)
        else:
            embedding = await generate_embedding_ollama(query_text)
    except Exception as e:
        return f"Error generating embedding: {str(e)}"

    try:
        result = sb.rpc("match_document_chunks", {
            "query_embedding": embedding,
            "match_count": top_k,
            "filter_event_id": event_id,
        }).execute()
    except Exception as e:
        return f"Error in semantic search: {str(e)}"

    chunks = result.data or []
    if not chunks:
        return f"No documents found matching: '{query_text}'"

    doc_ids = list(set(c.get("document_id", "") for c in chunks if c.get("document_id")))
    doc_map = {}
    if doc_ids:
        try:
            docs_result = sb.table("documents").select("id, title, file_name").in_("id", doc_ids).execute()
            doc_map = {d["id"]: d for d in (docs_result.data or [])}
        except Exception:
            pass

    if company_name:
        company_lower = company_name.lower()
        filtered = [c for c in chunks if company_lower in (doc_map.get(c.get("document_id", ""), {}).get("title") or doc_map.get(c.get("document_id", ""), {}).get("file_name") or "").lower()]
        if filtered:
            chunks = filtered

    lines = [f"Found {len(chunks)} relevant document chunks:\n"]
    for i, chunk in enumerate(chunks, 1):
        doc = doc_map.get(chunk.get("document_id", ""), {})
        title = doc.get("title") or doc.get("file_name") or "Unknown document"
        similarity = chunk.get("similarity", 0)
        text = chunk.get("parent_text") or chunk.get("chunk_text") or ""
        text = text[:600]
        lines.append(f"[{i}] **{title}** (relevance: {similarity:.2f})\n{text}")

    return "\n\n".join(lines)


async def _agent_get_company_details(tool_input: dict, event_id: str) -> str:
    sb = get_supabase()
    company_name = (tool_input.get("company_name") or "").strip()
    if not company_name:
        return "Error: 'company_name' is required."

    try:
        find_result = sb.rpc("kg_find_entity", {"search_name": company_name, "filter_event_id": event_id}).execute()
    except Exception as e:
        return f"Error finding entity: {str(e)}"

    entities = find_result.data or []
    if not entities:
        try:
            fallback = sb.table("kg_entities").select("id, name, properties").eq("event_id", event_id).eq("entity_type", "company").ilike("normalized_name", f"%{company_name.lower()}%").limit(5).execute()
            entities = fallback.data or []
        except Exception:
            pass

    if not entities:
        return f"No company found matching '{company_name}'. Try searching the portfolio with search_portfolio tool."

    entity = entities[0]
    entity_id = entity.get("id")

    try:
        card_result = sb.rpc("get_company_card", {"p_company_entity_id": entity_id, "p_filter_event_id": event_id}).execute()
    except Exception:
        props = entity.get("properties") or {}
        lines = [f"**{entity.get('name', company_name)}**\n"]
        for k, v in props.items():
            if v and k not in ("auto_created", "source", "property_sources"):
                lines.append(f"- {k.replace('_', ' ').title()}: {str(v)[:300]}")
        return "\n".join(lines)

    cards = card_result.data or []
    if not cards:
        props = entity.get("properties") or {}
        lines = [f"**{entity.get('name', company_name)}**\n"]
        for k, v in props.items():
            if v and k not in ("auto_created", "source", "property_sources"):
                lines.append(f"- {k.replace('_', ' ').title()}: {str(v)[:300]}")
        return "\n".join(lines)

    card = cards[0]
    props = card.get("company_properties") or {}
    lines = [f"**{card.get('company_name', company_name)}**\n"]
    for key in ["industry", "funding_stage", "business_model", "headquarters", "country",
                 "location", "bio", "website", "linkedin_url", "email", "phone",
                 "founded_year", "team_size", "amount_seeking", "valuation",
                 "arr", "mrr", "burn_rate", "runway_months",
                 "problem", "solution", "tam", "competitive_edge"]:
        val = props.get(key)
        if val:
            label = key.replace("_", " ").title()
            display = str(val)[:500] if key in ("bio", "problem", "solution", "tam", "competitive_edge") else str(val)
            lines.append(f"- **{label}**: {display}")
    geo = props.get("geo_focus") or props.get("geo_markets")
    if geo:
        lines.append(f"- **Geo Focus**: {', '.join(geo) if isinstance(geo, list) else geo}")
    founders = props.get("founders")
    if founders and isinstance(founders, list):
        founder_strs = []
        for f in founders:
            if isinstance(f, dict):
                founder_strs.append(f"{f.get('name', '?')} ({f.get('role', 'founder')})")
            else:
                founder_strs.append(str(f))
        lines.append(f"- **Founders**: {', '.join(founder_strs)}")
    doc_count = card.get("document_count", 0)
    rel_count = card.get("relationship_count", 0)
    kpi_count = card.get("kpi_count", 0)
    related = card.get("related_companies") or []
    lines.append(f"\n**Stats**: {doc_count} documents, {kpi_count} KPIs, {rel_count} relationships")
    if related:
        lines.append(f"**Related Companies**: {', '.join(related[:10])}")
    kpi_summary = card.get("kpi_summary")
    if kpi_summary and isinstance(kpi_summary, (dict, list)):
        lines.append(f"**KPI Summary**: {json.dumps(kpi_summary, default=str)[:500]}")
    return "\n".join(lines)


async def _agent_search_knowledge_graph(tool_input: dict, event_id: str) -> str:
    sb = get_supabase()
    entity_name = (tool_input.get("entity_name") or "").strip()
    entity_type = (tool_input.get("entity_type") or "").strip()
    relationship_type = (tool_input.get("relationship_type") or "").strip()
    max_depth = min(tool_input.get("max_depth", 1), 3)

    if not entity_name:
        return "Error: 'entity_name' is required."

    try:
        find_result = sb.rpc("kg_find_entity", {"search_name": entity_name, "filter_event_id": event_id}).execute()
    except Exception as e:
        return f"Error finding entity: {str(e)}"

    entities = find_result.data or []
    if entity_type:
        entities = [e for e in entities if (e.get("entity_type") or "").lower() == entity_type.lower()]
    if not entities:
        return f"No entity found matching '{entity_name}'" + (f" (type: {entity_type})" if entity_type else "")

    lines = [f"Found {len(entities)} entity/entities matching '{entity_name}':\n"]
    for entity in entities[:5]:
        eid = entity.get("id")
        ename = entity.get("name", "?")
        etype = entity.get("entity_type", "?")
        props = entity.get("properties") or {}
        lines.append(f"### {ename} ({etype})")
        for k, v in list(props.items())[:10]:
            if v and k not in ("auto_created", "source", "property_sources"):
                lines.append(f"  - {k}: {str(v)[:200]}")
        try:
            neighbors_result = sb.rpc("kg_neighbors", {"entity_id": eid, "max_depth": max_depth}).execute()
            neighbors = neighbors_result.data or []
        except Exception:
            neighbors = []
        if relationship_type:
            neighbors = [n for n in neighbors if (n.get("relation_type") or "").lower() == relationship_type.lower()]
        if neighbors:
            lines.append(f"  **Relationships** ({len(neighbors)}):")
            for n in neighbors[:15]:
                direction = n.get("direction", "?")
                rel_type = n.get("relation_type", "?")
                neighbor_name = n.get("entity_name", "?")
                neighbor_type = n.get("entity_type", "?")
                arrow = "\u2192" if direction == "outgoing" else "\u2190"
                lines.append(f"    {arrow} {rel_type} \u2014 {neighbor_name} ({neighbor_type})")
        else:
            lines.append("  No relationships found.")
    return "\n".join(lines)


async def _execute_agent_tool(tool_name: str, tool_input: dict, event_id: str) -> str:
    try:
        if tool_name == "search_portfolio":
            return await _agent_search_portfolio(tool_input, event_id)
        elif tool_name == "search_documents":
            return await _agent_search_documents(tool_input, event_id)
        elif tool_name == "get_company_details":
            return await _agent_get_company_details(tool_input, event_id)
        elif tool_name == "search_knowledge_graph":
            return await _agent_search_knowledge_graph(tool_input, event_id)
        else:
            return f"Unknown tool: {tool_name}"
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"Tool error ({tool_name}): {str(e)[:300]}"


AGENT_SYSTEM_PROMPT = """You are Company Assistant, a VC portfolio intelligence assistant for a venture capital firm.

You have access to tools that let you search the firm's portfolio database, documents, and knowledge graph.
ALWAYS use your tools to find information before answering. Never guess or say "I don't have access" without searching first.

## Tool Selection Rules

1. **Portfolio metadata questions** (counts, lists, filtering by country/sector/stage):
   Use `search_portfolio` with appropriate filters.

2. **Document content questions** (pitch details, meeting notes, financials):
   Use `search_documents` with a clear query.

3. **Specific company deep-dives** (detailed info about one company):
   Use `get_company_details` first, then `search_documents` if you need more.

4. **Relationship questions** (investors, founders, competitors, connections):
   Use `search_knowledge_graph`.

5. **Complex multi-step questions**:
   Use multiple tools in sequence.

## Important Rules

- ALWAYS search before claiming data is unavailable.
- If a search returns no results, try broader terms or a different tool.
- Cite information sources clearly in your response.
- Be precise with numbers.
- Keep responses well-structured with clear formatting.
- When listing companies, include key details (sector, stage, country) for each.
"""

MAX_AGENT_ITERATIONS = 4


class AgentAskImage(BaseModel):
    """Base64 image for vision: PNG/JPEG supported."""
    media_type: str = Field(alias="mediaType")  # "image/png" or "image/jpeg"
    data: str  # base64-encoded

    model_config = {"populate_by_name": True}


class AgentAskRequest(BaseModel):
    question: str = ""
    event_id: str = Field(default="", alias="eventId")
    previous_messages: List[ChatMessage] = Field(default_factory=list, alias="previousMessages")
    web_search_enabled: bool = Field(default=False, alias="webSearchEnabled")
    folder_ids: List[str] = Field(default_factory=list, alias="folderIds")
    images: List[AgentAskImage] = Field(default_factory=list)

    model_config = {"populate_by_name": True}


async def _load_org_system_prompt(event_id: str) -> str:
    """Resolve event_id → organization → system_prompt. Falls back to AGENT_SYSTEM_PROMPT."""
    try:
        sb = get_supabase()
        if sb:
            ev = sb.table("events").select("organization_id").eq("id", event_id).single().execute()
            org_id = (ev.data or {}).get("organization_id")
            if org_id:
                org = sb.table("organizations").select("system_prompt").eq("id", org_id).single().execute()
                prompt = (org.data or {}).get("system_prompt", "").strip()
                # Normalize legacy branding: stored prompts may still say "Orbit AI"
                if prompt:
                    prompt = re.sub(r"\bOrbit\s+AI\b", "Company Assistant", prompt, flags=re.IGNORECASE)
                if prompt and len(prompt) > 50:
                    return prompt
    except Exception:
        pass
    return AGENT_SYSTEM_PROMPT


@app.post("/ask/agent/stream")
async def ask_agent_stream(request: AgentAskRequest, auth: AuthContext = Depends(get_auth_context)):
    """Agentic RAG endpoint — Claude decides which tools to call, executes them,
    and generates a final answer. Streams SSE events for real-time UI updates."""
    try:
        if get_supabase() is None:
            raise HTTPException(status_code=503, detail="Supabase not configured. Set SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY.")
        question = (request.question or "").strip()
        images = request.images or []
        if not question and not images:
            raise HTTPException(status_code=400, detail="question or at least one image is required.")
        if question:
            resolved_question = await rewrite_query_with_llm(question, request.previous_messages or [])
        else:
            resolved_question = "What do you see in this image? Please describe and analyze it."

        event_id = (request.event_id or "").strip()
        if not event_id:
            raise HTTPException(status_code=400, detail="event_id is required.")

        system_prompt = await _load_org_system_prompt(event_id)

        messages: List[dict] = []
        for msg in (request.previous_messages or [])[-10:]:
            messages.append({"role": msg.role, "content": msg.content})

        # Build last user message: text + optional image blocks for Claude vision
        if not images:
            messages.append({"role": "user", "content": resolved_question})
        else:
            content_blocks: List[dict] = [{"type": "text", "text": resolved_question}]
            for img in images:
                media = (img.media_type or "image/png").lower()
                if media not in ("image/png", "image/jpeg", "image/jpg"):
                    media = "image/png"
                if media == "image/jpg":
                    media = "image/jpeg"
                content_blocks.append({
                    "type": "image",
                    "source": {"type": "base64", "media_type": media, "data": img.data},
                })
            messages.append({"role": "user", "content": content_blocks})

        tools = list(AGENT_TOOLS)
        if request.web_search_enabled:
            tools.append(ANTHROPIC_WEB_SEARCH_TOOL)

        is_comp = is_comprehensive_question(question)
        max_tokens = 8000 if is_comp else ASK_MAX_TOKENS

        async def generate():
            try:
                yield f"data: {json.dumps({'ping': True})}\n\n"
                yield f"data: {json.dumps({'status': 'Analyzing your question...'})}\n\n"

                if not _anthropic_sdk_available:
                    yield f"data: {json.dumps({'error': 'Anthropic SDK not available'})}\n\n"
                    yield "data: [DONE]\n\n"
                    return

                client = _get_anthropic_async_client()
                current_messages = list(messages)
                model_name = ANTHROPIC_MODEL_FALLBACKS[0] if ANTHROPIC_MODEL_FALLBACKS else "claude-sonnet-4-20250514"

                for iteration in range(MAX_AGENT_ITERATIONS):
                    response = None
                    last_create_error = None
                    max_retries = 3
                    for attempt in range(max_retries):
                        try:
                            response = await client.messages.create(
                                model=model_name,
                                max_tokens=max_tokens,
                                temperature=0.1,
                                system=system_prompt,
                                messages=current_messages,
                                tools=tools,
                            )
                            last_create_error = None
                            break
                        except anthropic.NotFoundError:
                            last_create_error = None
                            for fallback in ANTHROPIC_MODEL_FALLBACKS[1:]:
                                try:
                                    response = await client.messages.create(
                                        model=fallback,
                                        max_tokens=max_tokens,
                                        temperature=0.1,
                                        system=system_prompt,
                                        messages=current_messages,
                                        tools=tools,
                                    )
                                    model_name = fallback
                                    break
                                except Exception:
                                    continue
                            else:
                                yield f"data: {json.dumps({'error': 'All models failed'})}\n\n"
                                yield "data: [DONE]\n\n"
                                return
                            break  # fallback succeeded, exit retry loop
                        except Exception as e:
                            last_create_error = e
                            if _is_anthropic_overloaded_or_rate_limited(e) and attempt < max_retries - 1:
                                delay = (2 ** attempt) * 2 + random.uniform(0, 1)
                                await asyncio.sleep(delay)
                                continue
                            yield f"data: {json.dumps({'error': _friendly_anthropic_error_message(e)})}\n\n"
                            yield "data: [DONE]\n\n"
                            return
                    if response is None and last_create_error is not None:
                        yield f"data: {json.dumps({'error': _friendly_anthropic_error_message(last_create_error)})}\n\n"
                        yield "data: [DONE]\n\n"
                        return

                    tool_calls = []
                    text_parts = []
                    web_search_citations: Dict[str, str] = {}

                    for block in response.content:
                        if block.type == "text":
                            text_parts.append(block.text)
                        elif block.type == "tool_use":
                            tool_calls.append(block)
                        elif block.type == "server_tool_use":
                            yield f"data: {json.dumps({'status': 'Searching the web...'})}\n\n"
                        elif block.type == "web_search_tool_result":
                            result_content = getattr(block, "content", [])
                            if isinstance(result_content, list):
                                for item in result_content:
                                    if hasattr(item, "type") and getattr(item, "type", "") == "web_search_result":
                                        url = getattr(item, "url", "")
                                        title = getattr(item, "title", "")
                                        if url:
                                            web_search_citations[url] = title

                    if not tool_calls:
                        full_text = "\n".join(text_parts)
                        if full_text:
                            chunk_size = 80
                            for i in range(0, len(full_text), chunk_size):
                                yield f"data: {json.dumps({'text': full_text[i:i+chunk_size]})}\n\n"
                                await asyncio.sleep(0)
                        if web_search_citations:
                            sources_text = "\n\n**Web Sources:**"
                            for i, (url, title) in enumerate(web_search_citations.items(), 1):
                                sources_text += f"\n[{i}] [{title}]({url})"
                            yield f"data: {json.dumps({'text': sources_text})}\n\n"
                        yield "data: [DONE]\n\n"
                        return

                    tool_names = [tc.name for tc in tool_calls]
                    tool_label = ", ".join(tool_names)
                    yield f"data: {json.dumps({'status': f'Searching: {tool_label}...'})}\n\n"

                    assistant_content = []
                    for block in response.content:
                        if block.type == "text":
                            assistant_content.append({"type": "text", "text": block.text})
                        elif block.type == "tool_use":
                            assistant_content.append({
                                "type": "tool_use",
                                "id": block.id,
                                "name": block.name,
                                "input": block.input,
                            })

                    current_messages.append({"role": "assistant", "content": assistant_content})

                    tool_results = []
                    for tc in tool_calls:
                        result_text = await _execute_agent_tool(tc.name, tc.input, event_id)
                        tool_results.append({
                            "type": "tool_result",
                            "tool_use_id": tc.id,
                            "content": result_text,
                        })

                    current_messages.append({"role": "user", "content": tool_results})

                yield f"data: {json.dumps({'status': 'Generating final answer...'})}\n\n"
                try:
                    async with client.messages.stream(
                        model=model_name,
                        max_tokens=max_tokens,
                        temperature=0.1,
                        system=system_prompt,
                        messages=current_messages,
                    ) as final_stream:
                        async for event in final_stream:
                            if event.type == "content_block_delta" and hasattr(event.delta, "type"):
                                if event.delta.type == "text_delta" and hasattr(event.delta, "text"):
                                    yield f"data: {json.dumps({'text': event.delta.text})}\n\n"
                except Exception as e:
                    yield f"data: {json.dumps({'error': _friendly_anthropic_error_message(e)})}\n\n"

                yield "data: [DONE]\n\n"

            except Exception as e:
                import traceback
                traceback.print_exc()
                yield f"data: {json.dumps({'error': _friendly_anthropic_error_message(e)})}\n\n"
                yield "data: [DONE]\n\n"

        return StreamingResponse(
            generate(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Methods": "POST, OPTIONS",
                "Access-Control-Allow-Headers": "*",
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Agent error: {str(e)[:200]}")


def get_supabase():
    """Get Supabase client using existing env vars."""
    url = os.getenv("SUPABASE_URL", "")
    key = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "") or os.getenv("SUPABASE_SERVICE_KEY", "")
    if not url or not key:
        return None
    try:
        from supabase import create_client
        return create_client(url, key)
    except Exception:
        return None


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", "10000"))
    uvicorn.run(app, host="0.0.0.0", port=port)
